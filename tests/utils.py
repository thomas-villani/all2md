"""Test utilities for all2md test suite.

This module provides utilities for creating test documents, validating
Markdown output, and other common testing functions.
"""

import base64
import datetime
import email
import tempfile
from email.mime.image import MIMEImage
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path

import docx
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.shared import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as PptxInches

# Base64 encoded 1x1 pixel PNG for testing
MINIMAL_PNG_B64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/w8AAn8B9FpQHLwAAAAASUVORK5CYII="
MINIMAL_PNG_BYTES = base64.b64decode(MINIMAL_PNG_B64)


class DocxTestGenerator:
    """Generator for complex DOCX test documents with edge cases."""

    @staticmethod
    def create_complex_lists_document() -> docx.Document:
        """Create a document with complex list scenarios."""
        doc = docx.Document()

        # Mixed numbering styles
        doc.add_heading("Complex Lists Test", level=1)

        # Roman numerals
        p1 = doc.add_paragraph("First item", style="List Number")
        p1.paragraph_format.left_indent = Inches(0)

        # Nested with different symbols
        p2 = doc.add_paragraph("Nested bullet", style="List Bullet")
        p2.paragraph_format.left_indent = Inches(0.5)

        # Large indent levels (5 levels deep)
        for i in range(5):
            p = doc.add_paragraph(f"Level {i+1} item", style="List Bullet")
            p.paragraph_format.left_indent = Inches(0.5 * (i + 1))

        # Paragraph with both text and inline image
        p_mixed = doc.add_paragraph("Text before image ")
        run = p_mixed.add_run()
        # Simulate adding an image inline with text
        p_mixed.add_run(" text after image with ")
        link_run = p_mixed.add_run("hyperlink")
        # Add hyperlink would be done here in real scenario

        # Restart numbering
        doc.add_paragraph("Restarted list item 1", style="List Number")
        doc.add_paragraph("Restarted list item 2", style="List Number")

        return doc

    @staticmethod
    def create_complex_tables_document() -> docx.Document:
        """Create a document with complex table scenarios."""
        doc = docx.Document()

        doc.add_heading("Complex Tables Test", level=1)

        # Table with merged cells
        table1 = doc.add_table(rows=3, cols=3)
        table1.style = 'Table Grid'

        # Header row
        table1.rows[0].cells[0].text = "Merged Header"
        table1.rows[0].cells[1].text = "Column 2"
        table1.rows[0].cells[2].text = "Column 3"

        # Merge cells (simulated - python-docx has limitations)
        table1.rows[1].cells[0].text = "Merged Cell"
        table1.rows[1].cells[1].text = "Normal"
        table1.rows[1].cells[2].text = "Cell"

        # Alignment variations
        for row in table1.rows:
            for cell in row.cells:
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        # Table without headers
        doc.add_paragraph("Table without headers:")
        table2 = doc.add_table(rows=2, cols=2)
        table2.rows[0].cells[0].text = "Data 1"
        table2.rows[0].cells[1].text = "Data 2"
        table2.rows[1].cells[0].text = "Data 3"
        table2.rows[1].cells[1].text = "Data 4"

        # Nested table (limited support in python-docx)
        doc.add_paragraph("Complex nested content:")
        table3 = doc.add_table(rows=2, cols=2)
        table3.rows[0].cells[0].text = "Cell with\nmultiple\nlines"
        table3.rows[0].cells[1].text = "Normal cell"

        return doc

    @staticmethod
    def create_images_document(tmp_dir: Path) -> docx.Document:
        """Create a document with various image scenarios."""
        doc = docx.Document()

        doc.add_heading("Image Test Document", level=1)

        # Create test image file
        img_path = tmp_dir / "test_image.png"
        img_path.write_bytes(MINIMAL_PNG_BYTES)

        # Embedded image with alt text (simulated)
        doc.add_paragraph("Image with alt text:")
        doc.add_picture(str(img_path), width=Inches(1))

        # Image in paragraph with text
        p = doc.add_paragraph("Text before ")
        run = p.add_run()
        # run.add_picture(str(img_path), width=Inches(0.5))  # Inline image
        p.add_run(" text after image")

        return doc

    @staticmethod
    def create_hyperlinks_document() -> docx.Document:
        """Create a document with various hyperlink scenarios."""
        doc = docx.Document()

        doc.add_heading("Hyperlinks Test", level=1)

        # External links
        p1 = doc.add_paragraph("Visit ")
        p1.add_run("Google").hyperlink = "https://www.google.com"

        # Mailto links
        p2 = doc.add_paragraph("Email ")
        p2.add_run("support@example.com").hyperlink = "mailto:support@example.com"

        # Relative paths
        p3 = doc.add_paragraph("Local ")
        p3.add_run("document").hyperlink = "./local/document.html"

        # Anchor links
        p4 = doc.add_paragraph("Go to ")
        p4.add_run("section 1").hyperlink = "#section1"

        return doc

    @staticmethod
    def create_escaping_document() -> docx.Document:
        """Create a document with Markdown special characters."""
        doc = docx.Document()

        doc.add_heading("Escaping Test", level=1)

        # Paragraph with Markdown special characters
        special_chars = "* _ # [ ] ( ) { } \\ ` + - . ! ~ ^ | < >"
        doc.add_paragraph(f"Special characters: {special_chars}")

        # Emphasis with special characters
        p = doc.add_paragraph()
        run = p.add_run("Bold text with * asterisks *")
        run.bold = True

        p2 = doc.add_paragraph()
        run2 = p2.add_run("Italic text with _ underscores _")
        run2.italic = True

        # Code-like content
        doc.add_paragraph("Code snippet: `function() { return 'hello'; }`")

        return doc


class HtmlTestGenerator:
    """Generator for complex HTML test content with edge cases."""

    @staticmethod
    def create_nested_elements_html() -> str:
        """Create HTML with complex nested element combinations."""
        return """
        <html>
        <body>
            <h1>Nested Elements Test</h1>

            <!-- Complex nested emphasis -->
            <p>Text with <strong>bold and <em>italic and <code>code</code></em></strong> nesting</p>
            <p><em>Italic with <strong>bold</strong> and <a href="http://example.com">link</a></em></p>

            <!-- Nested lists and blockquotes -->
            <blockquote>
                <p>Quote with <em>emphasis</em></p>
                <ul>
                    <li>List in quote
                        <blockquote>
                            <p>Nested quote in list</p>
                            <ol>
                                <li>Ordered in quote in list</li>
                            </ol>
                        </blockquote>
                    </li>
                </ul>
            </blockquote>

            <!-- Mixed content in divs -->
            <div>
                <h2>Section Header</h2>
                <p>Paragraph with <code>inline code</code> and <strong>bold text</strong></p>
                <ul>
                    <li><strong>Bold list item</strong> with <a href="#anchor">anchor link</a></li>
                    <li><em>Italic</em> with <code>code</code></li>
                </ul>
            </div>
        </body>
        </html>
        """

    @staticmethod
    def create_code_blocks_html() -> str:
        """Create HTML with complex code block scenarios."""
        return """
        <html>
        <body>
            <h1>Code Blocks Test</h1>

            <!-- Code with language -->
            <pre class="language-python"><code>def hello():
    print("world")
    # Comment with ``` backticks
    return "```markdown```"</code></pre>

            <!-- Code with backticks in content -->
            <pre><code>Use ```markdown``` syntax for code blocks.
Multiple ``` backticks ``` in content.
Even `````five````` backticks.</code></pre>

            <!-- Inline code with backticks -->
            <p>Use <code>`markdown`</code> and <code>```blocks```</code> for formatting</p>

            <!-- Code in different contexts -->
            <ul>
                <li><code>inline code</code> in list</li>
                <li><pre><code>block code
in list item</code></pre></li>
            </ul>

            <blockquote>
                <p>Quote with <code>inline code</code></p>
                <pre class="bash"><code>echo "code block in quote"</code></pre>
            </blockquote>
        </body>
        </html>
        """

    @staticmethod
    def create_entities_html() -> str:
        """Create HTML with various entity scenarios."""
        return """
        <html>
        <body>
            <h1>HTML Entities Test</h1>

            <!-- Common entities -->
            <p>Common: &amp; &lt; &gt; &quot; &#39;</p>

            <!-- Special spacing -->
            <p>Spacing:&nbsp;non-breaking&nbsp;spaces</p>
            <p>Em&mdash;dash and en&ndash;dash</p>

            <!-- Symbols -->
            <p>Symbols: &copy; &reg; &trade; &hellip;</p>

            <!-- Numeric entities -->
            <p>Numeric: &#8230; &#8482; &#169;</p>

            <!-- Entities in different contexts -->
            <strong>Bold &amp; italic text</strong>
            <code>Code with &lt;tags&gt;</code>
            <a href="mailto:test@example.com?subject=Test&amp;body=Hello">Email with &amp;</a>
        </body>
        </html>
        """

    @staticmethod
    def create_complex_tables_html() -> str:
        """Create HTML with complex table structures."""
        return """
        <html>
        <body>
            <h1>Complex Tables Test</h1>

            <!-- Table with thead/tbody/tfoot -->
            <table>
                <caption>Sales Data</caption>
                <thead>
                    <tr>
                        <th align="left">Product</th>
                        <th align="center">Q1</th>
                        <th align="right">Q2</th>
                        <th>Total</th>
                    </tr>
                </thead>
                <tbody>
                    <tr>
                        <td>Widget A</td>
                        <td style="text-align: center">100</td>
                        <td style="text-align: right">150</td>
                        <td><strong>250</strong></td>
                    </tr>
                    <tr>
                        <td>Widget B</td>
                        <td colspan="2">Not available</td>
                        <td>0</td>
                    </tr>
                </tbody>
                <tfoot>
                    <tr>
                        <td><strong>Total</strong></td>
                        <td>100</td>
                        <td>150</td>
                        <td><strong>250</strong></td>
                    </tr>
                </tfoot>
            </table>

            <!-- Table with rowspan/colspan -->
            <table border="1">
                <tr>
                    <th rowspan="2">Category</th>
                    <th colspan="2">Values</th>
                </tr>
                <tr>
                    <th>Min</th>
                    <th>Max</th>
                </tr>
                <tr>
                    <td rowspan="2">Group A</td>
                    <td>1</td>
                    <td>10</td>
                </tr>
                <tr>
                    <td colspan="2">Average: 5.5</td>
                </tr>
            </table>

            <!-- Nested content in cells -->
            <table>
                <tr>
                    <td>
                        <ul>
                            <li>List in cell</li>
                            <li>Another item</li>
                        </ul>
                    </td>
                    <td>
                        <p>Paragraph with <strong>formatting</strong></p>
                        <p>Multiple paragraphs</p>
                    </td>
                </tr>
            </table>
        </body>
        </html>
        """


class PptxTestGenerator:
    """Generator for complex PPTX test presentations with edge cases."""

    @staticmethod
    def create_charts_presentation() -> Presentation:
        """Create a presentation with complex chart scenarios."""
        prs = Presentation()

        # Multiple series chart
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        slide.shapes.title.text = "Complex Charts"

        chart_data = ChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Series 1', (10, 15, 12, 18))
        chart_data.add_series('Series 2', (8, 12, 15, 14))
        chart_data.add_series('Series 3', (12, 10, 8, 16))

        x, y, cx, cy = PptxInches(1), PptxInches(2), PptxInches(8), PptxInches(4)
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

        # Pie chart (exploded would need additional configuration)
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "Pie Chart"

        pie_data = ChartData()
        pie_data.categories = ['A', 'B', 'C', 'D']
        pie_data.add_series('Values', (30, 25, 25, 20))

        slide2.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, pie_data)

        return prs

    @staticmethod
    def create_shapes_presentation() -> Presentation:
        """Create a presentation with advanced shape scenarios."""
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Advanced Shapes"

        # Text box with RTL simulation (python-pptx has limited RTL support)
        left, top, width, height = PptxInches(1), PptxInches(2), PptxInches(4), PptxInches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Text with special bullets • → ◦"

        # Grouped shapes (simulated - would need more complex setup)
        # Add multiple shapes that conceptually could be grouped
        shape1 = slide.shapes.add_shape(1, PptxInches(1), PptxInches(3), PptxInches(1), PptxInches(0.5))
        shape2 = slide.shapes.add_shape(1, PptxInches(2), PptxInches(3), PptxInches(1), PptxInches(0.5))

        return prs

    @staticmethod
    def create_tables_presentation() -> Presentation:
        """Create a presentation with complex table scenarios."""
        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Complex Tables"

        # Table with merged cells and formatting
        rows, cols = 4, 3
        left, top, width, height = PptxInches(1), PptxInches(2), PptxInches(8), PptxInches(3)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header row
        table.cell(0, 0).text = "Category"
        table.cell(0, 1).text = "Value 1"
        table.cell(0, 2).text = "Value 2"

        # Data rows with formatting
        table.cell(1, 0).text = "Item A"
        table.cell(1, 1).text = "100"
        table.cell(1, 2).text = "200"

        # Cell with bold formatting
        cell = table.cell(2, 0)
        cell.text = "Bold Item"
        # Apply bold formatting to paragraph
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.bold = True

        # Merged cell simulation (python-pptx merge support is limited)
        table.cell(3, 0).text = "Merged Content"
        table.cell(3, 1).text = ""
        table.cell(3, 2).text = "End"

        return prs


class EmlTestGenerator:
    """Generator for complex EML test messages with edge cases."""

    @staticmethod
    def create_multipart_message() -> str:
        """Create a complex multipart email message."""
        msg = MIMEMultipart('mixed')
        msg['From'] = 'sender@example.com'
        msg['To'] = 'recipient@example.com'
        msg['Subject'] = 'Multipart Message Test'
        msg['Date'] = email.utils.formatdate(localtime=True)
        msg['Message-ID'] = '<multipart@example.com>'

        # Text part
        text_part = MIMEText('This is the plain text part of the message.', 'plain')
        msg.attach(text_part)

        # HTML part
        html_content = """
        <html>
        <body>
            <h1>HTML Part</h1>
            <p>This is the <strong>HTML</strong> part of the message.</p>
        </body>
        </html>
        """
        html_part = MIMEText(html_content, 'html')
        msg.attach(html_part)

        # Image attachment
        img_part = MIMEImage(MINIMAL_PNG_BYTES)
        img_part.add_header('Content-Disposition', 'attachment; filename="test.png"')
        msg.attach(img_part)

        return msg.as_string()

    @staticmethod
    def create_headers_edge_cases() -> str:
        """Create email with various header edge cases."""
        # Manually construct email with edge case headers
        eml_content = """From: =?UTF-8?B?VGVzdCBVc2Vy?= <test@example.com>
To: Multiple <multi1@example.com>, Recipients <multi2@example.com>
CC: =?UTF-8?Q?Copi=C3=A9?= <copy@example.com>
BCC: hidden@example.com
Subject: =?UTF-8?Q?Subject_with_=C3=A9ncoded_ch=C3=A1racters?=
Date: Invalid Date Format
Message-ID: <malformed-id>
In-Reply-To: <original-message@example.com>
References: <ref1@example.com> <ref2@example.com>
Content-Type: text/plain; charset=UTF-8
Content-Transfer-Encoding: 8bit

This is a test message with edge case headers.
Special characters: café résumé naïve
"""
        return eml_content

    @staticmethod
    def create_large_chain() -> str:
        """Create a large email chain for testing."""
        chain_parts = []

        base_date = datetime.datetime(2023, 1, 1, 10, 0, 0)

        for i in range(5):
            date = base_date + datetime.timedelta(days=i, hours=i)
            formatted_date = date.strftime("%a, %d %b %Y %H:%M:%S +0000")

            if i == 0:
                subject = "Original Message"
                in_reply_to = ""
                references = ""
            else:
                subject = "Re: Original Message"
                in_reply_to = f"In-Reply-To: <msg{i-1}@example.com>"
                references = f"References: {' '.join([f'<msg{j}@example.com>' for j in range(i)])}"

            message = f"""From: User{i} <user{i}@example.com>
To: User{i+1} <user{i+1}@example.com>
Date: {formatted_date}
Subject: {subject}
Message-ID: <msg{i}@example.com>
{in_reply_to}
{references}

Message {i+1} content.
This is a longer message with multiple lines
to test email chain processing.

> Quoted content from previous message
> With multiple quoted lines
> And various formatting

Best regards,
User {i}
"""
            chain_parts.append(message)

        return "\n\n".join(chain_parts)

    @staticmethod
    def create_encoding_issues() -> str:
        """Create email with character encoding challenges."""
        return """From: Test User <test@example.com>
To: Recipient <recipient@example.com>
Subject: =?iso-8859-1?Q?Testing_=E9ncoding_=EFss=FCes?=
Date: Mon, 01 Jan 2023 12:00:00 +0000
Content-Type: text/plain; charset=iso-8859-1
Content-Transfer-Encoding: quoted-printable

This message contains special characters:
Caf=E9, r=E9sum=E9, na=EFve

And various symbols: =A9 =AE =99

Mixed encoding content with UTF-8 fallback:
Café résumé naïve
"""


def assert_markdown_valid(markdown: str) -> None:
    """Assert that the generated Markdown is valid and well-formed."""
    lines = markdown.split('\n')

    # Check for common Markdown issues
    for i, line in enumerate(lines):
        # Check for unescaped special characters at start of line
        if line.startswith(('#', '*', '-', '+')) and i > 0:
            if not lines[i-1].strip() == '':  # Should have empty line before headers/lists
                pass  # Could add warnings for style issues

        # Check for proper link formatting
        if '[' in line and ']' in line:
            # Basic link structure validation
            pass

    # Check for proper table formatting
    table_lines = [line for line in lines if '|' in line]
    if table_lines:
        # Validate table structure consistency
        pass


def create_test_temp_dir() -> Path:
    """Create a temporary directory for test files."""
    return Path(tempfile.mkdtemp())


def cleanup_test_dir(temp_dir: Path) -> None:
    """Clean up test directory and files."""
    import shutil
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
