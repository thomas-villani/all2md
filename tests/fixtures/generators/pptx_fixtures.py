"""PPTX test fixture generators using python-pptx.

This module provides functions to programmatically create PPTX documents
for testing various aspects of PPTX-to-Markdown conversion.
"""

import tempfile
from io import BytesIO
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_pptx_with_basic_slides() -> Presentation:
    """Create a PPTX with basic slides for testing slide conversion.

    Returns
    -------
    Presentation
        PowerPoint presentation with title, content, and section slides.

    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Test Presentation"
    subtitle.text = "Generated for testing PPTX to Markdown conversion"

    # Content slide with bullet points
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Main Content Slide"
    tf = body.text_frame
    tf.text = "First bullet point"

    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Sub-bullet point"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Another sub-bullet"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Third main bullet point"
    p.level = 0

    # Section header slide
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Section 2: Additional Content"

    # Another content slide
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Features and Benefits"
    tf = body.text_frame
    tf.text = "Easy to use interface"

    p = tf.add_paragraph()
    p.text = "Comprehensive feature set"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Document conversion"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Multiple format support"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Reliable performance"
    p.level = 0

    return prs


def create_pptx_with_tables() -> Presentation:
    """Create a PPTX with tables for testing table conversion.

    Returns
    -------
    Presentation
        PowerPoint presentation with table slides.

    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table Test Presentation"

    # Table slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title shape
    left = top = Inches(1)
    width = height = Inches(1)
    title_shape = slide.shapes.add_textbox(left, top, Inches(8), height)
    title_frame = title_shape.text_frame
    title_frame.text = "Sample Data Table"

    # Add table
    rows, cols = 4, 3
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header row
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(0, 2).text = "City"

    # Data rows
    data = [
        ("Alice Johnson", "25", "New York"),
        ("Bob Smith", "30", "San Francisco"),
        ("Carol Davis", "28", "Los Angeles")
    ]

    for i, (name, age, city) in enumerate(data, 1):
        table.cell(i, 0).text = name
        table.cell(i, 1).text = age
        table.cell(i, 2).text = city

    # Second slide with complex table
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Quarterly Sales Data"

    # Add complex table
    table2 = slide.shapes.add_table(4, 4, Inches(1.5), Inches(2.5), Inches(7), Inches(3)).table

    # Headers
    headers = ["Product", "Q1", "Q2", "Total"]
    for i, header in enumerate(headers):
        table2.cell(0, i).text = header

    # Data with totals
    sales_data = [
        ("Product A", "100", "150", "250"),
        ("Product B", "200", "180", "380"),
        ("Total", "300", "330", "630")
    ]

    for i, row_data in enumerate(sales_data, 1):
        for j, cell_data in enumerate(row_data):
            table2.cell(i, j).text = cell_data

    return prs


def create_pptx_with_shapes() -> Presentation:
    """Create a PPTX with shapes for testing shape conversion.

    Returns
    -------
    Presentation
        PowerPoint presentation with various shapes and text boxes.

    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Shapes Test Presentation"

    # Shapes slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Various Shapes and Text Boxes"

    # Add text box with formatted text
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(2))
    tf = text_box.text_frame
    tf.text = "This is a text box with multiple paragraphs."

    p = tf.add_paragraph()
    p.text = "Second paragraph with different formatting."
    p.font.size = Pt(14)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Third paragraph in italic."
    p.font.italic = True

    # Add rectangle with text
    rectangle = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(6), Inches(2), Inches(2.5), Inches(1.5)
    )
    rect_frame = rectangle.text_frame
    rect_frame.text = "Rectangle with text content"

    # Add another text box positioned differently
    text_box2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(7), Inches(1.5))
    tf2 = text_box2.text_frame
    tf2.text = "Text box at bottom of slide with conclusion text."

    # Slide with callout boxes (text boxes positioned as callouts)
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Callout and Annotation Examples"

    # Main content area
    main_text = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(3))
    main_frame = main_text.text_frame
    main_frame.text = "Main content area with primary information."

    # Callout boxes
    callout1 = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(2), Inches(1))
    callout1_frame = callout1.text_frame
    callout1_frame.text = "Important note: This is a callout."

    callout2 = slide.shapes.add_textbox(Inches(6.5), Inches(4), Inches(2), Inches(1))
    callout2_frame = callout2.text_frame
    callout2_frame.text = "Additional info: More details here."

    return prs


def create_pptx_with_charts() -> Presentation:
    """Create a PPTX with charts for testing chart conversion.

    Returns
    -------
    Presentation
        PowerPoint presentation with various chart types.

    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Chart Test Presentation"

    # Chart slide - Column chart
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Sales Data - Column Chart"

    # Chart data
    chart_data = ChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('Product A', (100, 150, 120, 180))
    chart_data.add_series('Product B', (200, 180, 220, 190))

    # Add chart
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    # Line chart slide
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Trends Over Time - Line Chart"

    # Line chart data
    line_data = ChartData()
    line_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    line_data.add_series('Revenue', (50, 65, 55, 70, 80, 85))
    line_data.add_series('Expenses', (30, 35, 40, 38, 42, 45))

    # Add line chart
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, line_data
    )

    # Pie chart slide
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Market Share - Pie Chart"

    # Pie chart data
    pie_data = ChartData()
    pie_data.categories = ['Company A', 'Company B', 'Company C', 'Others']
    pie_data.add_series('Market Share', (40, 25, 20, 15))

    # Add pie chart
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, pie_data
    )

    return prs


def create_pptx_with_formatting() -> Presentation:
    """Create a PPTX with text formatting for testing emphasis detection.

    Returns
    -------
    Presentation
        PowerPoint presentation with various text formatting.

    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Text Formatting Test"

    # Content slide with formatted text
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Text Formatting Examples"
    tf = body.text_frame
    tf.text = "Regular text"

    p = tf.add_paragraph()
    p.text = "Bold text"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Italic text"
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = "Bold and italic text"
    p.font.bold = True
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = "Underlined text"
    p.font.underline = True

    p = tf.add_paragraph()
    p.text = "Large text"
    p.font.size = Pt(24)

    p = tf.add_paragraph()
    p.text = "Small text"
    p.font.size = Pt(8)

    # Slide with text box formatting
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Advanced Text Formatting"
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True

    # Mixed formatting text box
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = text_box.text_frame

    # First paragraph with mixed formatting
    p1 = tf.paragraphs[0]
    p1.text = "This paragraph contains "

    run = p1.runs[0]
    run.text += "bold text "
    # Note: In python-pptx, we need to work with runs for inline formatting
    # This is a simplified version for testing

    p2 = tf.add_paragraph()
    p2.text = "Centered text"
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Right-aligned text"
    p3.alignment = PP_ALIGN.RIGHT

    return prs


def create_minimal_pptx(title: str = "Test Presentation", content: str = "Test content") -> Presentation:
    """Create a minimal PPTX presentation for basic testing.

    Parameters
    ----------
    title : str, optional
        Presentation title, by default "Test Presentation"
    content : str, optional
        Slide content, by default "Test content"

    Returns
    -------
    Presentation
        Simple presentation with title and content slide.

    """
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Content"
    body = slide.placeholders[1]
    body.text = content

    return prs


def save_pptx_to_bytes(prs: Presentation) -> bytes:
    """Save a PPTX presentation to bytes for testing.

    Parameters
    ----------
    prs : Presentation
        Presentation to convert to bytes

    Returns
    -------
    bytes
        Presentation as bytes

    """
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.read()


def create_pptx_file(prs: Presentation, file_path: Optional[Path] = None) -> Path:
    """Save a PPTX presentation to a temporary file.

    Parameters
    ----------
    prs : Presentation
        Presentation to save
    file_path : Path, optional
        File path to save to, by default creates temporary file

    Returns
    -------
    Path
        Path to the saved file

    """
    if file_path is None:
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        file_path = Path(temp_file.name)
        temp_file.close()

    prs.save(str(file_path))
    return file_path
