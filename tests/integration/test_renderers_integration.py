#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/integration/test_renderers_integration.py
"""Integration tests for renderers.

Tests cover:
- End-to-end rendering workflows
- Round-trip conversions (e.g., Markdown -> AST -> HTML -> verify)
- Integration with parsers
- Real-world document scenarios
- Cross-renderer compatibility

"""

import pytest
from pathlib import Path

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from reportlab.platypus import SimpleDocTemplate
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from ebooklib import epub
    EBOOKLIB_AVAILABLE = True
except ImportError:
    EBOOKLIB_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from all2md.ast import (
    BlockQuote,
    CodeBlock,
    Document,
    Heading,
    Image,
    Link,
    List,
    ListItem,
    Paragraph,
    Strong,
    Table,
    TableCell,
    TableRow,
    Text,
    ThematicBreak,
)
from all2md.options import (
    DocxRendererOptions,
    EpubRendererOptions,
    HtmlRendererOptions,
    MarkdownOptions,
    PdfRendererOptions,
    PptxRendererOptions,
)
from all2md.renderers import MarkdownRenderer, HtmlRenderer

if DOCX_AVAILABLE:
    from all2md.renderers.docx import DocxRenderer

if REPORTLAB_AVAILABLE:
    from all2md.renderers.pdf import PdfRenderer

if EBOOKLIB_AVAILABLE:
    from all2md.renderers.epub import EpubRenderer

if PPTX_AVAILABLE:
    from all2md.renderers.pptx import PptxRenderer


def create_sample_document():
    """Create a sample AST document for testing."""
    return Document(
        metadata={"title": "Sample Document", "author": "Test Author"},
        children=[
            Heading(level=1, content=[Text(content="Document Title")]),
            Paragraph(content=[
                Text(content="This is a paragraph with "),
                Strong(content=[Text(content="bold text")]),
                Text(content=" and a "),
                Link(url="https://example.com", content=[Text(content="link")]),
                Text(content=".")
            ]),
            Heading(level=2, content=[Text(content="Lists")]),
            List(ordered=False, items=[
                ListItem(children=[Paragraph(content=[Text(content="First item")])]),
                ListItem(children=[Paragraph(content=[Text(content="Second item")])]),
                ListItem(children=[Paragraph(content=[Text(content="Third item")])])
            ]),
            Heading(level=2, content=[Text(content="Code Example")]),
            CodeBlock(content='def hello():\n    print("Hello, world!")', language="python"),
            Heading(level=2, content=[Text(content="Table")]),
            Table(
                header=TableRow(cells=[
                    TableCell(content=[Text(content="Name")]),
                    TableCell(content=[Text(content="Value")])
                ]),
                rows=[
                    TableRow(cells=[
                        TableCell(content=[Text(content="Alpha")]),
                        TableCell(content=[Text(content="1")])
                    ]),
                    TableRow(cells=[
                        TableCell(content=[Text(content="Beta")]),
                        TableCell(content=[Text(content="2")])
                    ])
                ]
            ),
            Heading(level=2, content=[Text(content="Quote")]),
            BlockQuote(children=[
                Paragraph(content=[Text(content="This is a blockquote.")])
            ])
        ]
    )


@pytest.mark.integration
class TestMarkdownRendering:
    """Integration tests for Markdown rendering."""

    def test_full_document_to_markdown(self):
        """Test rendering complete document to Markdown."""
        doc = create_sample_document()
        renderer = MarkdownRenderer(MarkdownOptions())
        result = renderer.render_to_string(doc)

        # Verify all major elements are present
        assert "# Document Title" in result
        assert "**bold text**" in result
        assert "[link](https://example.com)" in result
        assert "* First item" in result
        assert "```python" in result
        assert "| Name" in result
        assert "> This is a blockquote" in result

    def test_markdown_different_flavors(self):
        """Test rendering with different Markdown flavors."""
        doc = create_sample_document()

        # GFM
        gfm_renderer = MarkdownRenderer(MarkdownOptions(flavor="gfm"))
        gfm_result = gfm_renderer.render_to_string(doc)
        assert "| Name" in gfm_result  # Tables supported in GFM

        # CommonMark
        cm_renderer = MarkdownRenderer(MarkdownOptions(
            flavor="commonmark",
            unsupported_table_mode="html"
        ))
        cm_result = cm_renderer.render_to_string(doc)
        # Tables as HTML in CommonMark
        assert "<table>" in cm_result


@pytest.mark.integration
class TestHtmlRendering:
    """Integration tests for HTML rendering."""

    def test_full_document_to_html_standalone(self):
        """Test rendering complete document to standalone HTML."""
        doc = create_sample_document()
        renderer = HtmlRenderer(HtmlRendererOptions(standalone=True))
        result = renderer.render_to_string(doc)

        # Verify HTML structure
        assert "<!DOCTYPE html>" in result
        assert "<html" in result
        assert "<head>" in result
        assert "<body>" in result
        assert "</html>" in result

        # Verify content
        assert "<h1" in result
        assert "Document Title" in result
        assert "<strong>bold text</strong>" in result
        assert '<a href="https://example.com">link</a>' in result
        assert "<ul>" in result
        assert "<table>" in result
        assert "<blockquote>" in result

    def test_full_document_to_html_fragment(self):
        """Test rendering complete document to HTML fragment."""
        doc = create_sample_document()
        renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        result = renderer.render_to_string(doc)

        # Should not have document structure
        assert "<!DOCTYPE html>" not in result
        assert "<html" not in result

        # But should have content
        assert "<h1" in result
        assert "<p>" in result

    def test_html_with_toc(self):
        """Test HTML rendering with table of contents."""
        doc = create_sample_document()
        renderer = HtmlRenderer(HtmlRendererOptions(standalone=True, include_toc=True))
        result = renderer.render_to_string(doc)

        assert '<nav id="table-of-contents">' in result
        assert "Document Title" in result
        assert "Lists" in result
        assert "Table" in result

    def test_html_css_options(self):
        """Test different CSS styling options."""
        doc = Document(children=[Paragraph(content=[Text(content="Test")])])

        # Embedded CSS
        embedded_renderer = HtmlRenderer(HtmlRendererOptions(
            standalone=True,
            css_style="embedded"
        ))
        embedded_result = embedded_renderer.render_to_string(doc)
        assert "<style>" in embedded_result

        # No CSS
        no_css_renderer = HtmlRenderer(HtmlRendererOptions(
            standalone=True,
            css_style="none"
        ))
        no_css_result = no_css_renderer.render_to_string(doc)
        assert "<style>" not in no_css_result


@pytest.mark.integration
@pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
@pytest.mark.docx
class TestDocxRendering:
    """Integration tests for DOCX rendering."""

    def test_full_document_to_docx(self, tmp_path):
        """Test rendering complete document to DOCX."""
        doc = create_sample_document()
        renderer = DocxRenderer()
        output_file = tmp_path / "full_document.docx"
        renderer.render(doc, output_file)

        assert output_file.exists()

        # Verify DOCX content
        docx_doc = DocxDocument(str(output_file))

        # Check for tables
        assert len(docx_doc.tables) >= 1

        # Check for content in paragraphs
        all_text = " ".join(p.text for p in docx_doc.paragraphs)
        assert "Document Title" in all_text
        assert "bold text" in all_text
        assert "First item" in all_text

    def test_docx_with_custom_styles(self, tmp_path):
        """Test DOCX rendering with custom styles."""
        doc = create_sample_document()
        options = DocxRendererOptions(
            default_font="Arial",
            default_font_size=12,
            code_font="Consolas"
        )
        renderer = DocxRenderer(options)
        output_file = tmp_path / "custom_styles.docx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_docx_metadata(self, tmp_path):
        """Test DOCX metadata handling."""
        doc = create_sample_document()
        renderer = DocxRenderer()
        output_file = tmp_path / "with_metadata.docx"
        renderer.render(doc, output_file)

        docx_doc = DocxDocument(str(output_file))
        assert docx_doc.core_properties.title == "Sample Document"
        assert docx_doc.core_properties.author == "Test Author"


@pytest.mark.integration
@pytest.mark.skipif(not REPORTLAB_AVAILABLE, reason="reportlab not installed")
@pytest.mark.pdf
class TestPdfRendering:
    """Integration tests for PDF rendering."""

    def test_full_document_to_pdf(self, tmp_path):
        """Test rendering complete document to PDF."""
        doc = create_sample_document()
        renderer = PdfRenderer()
        output_file = tmp_path / "full_document.pdf"
        renderer.render(doc, output_file)

        assert output_file.exists()
        assert output_file.stat().st_size > 0

    def test_pdf_page_sizes(self, tmp_path):
        """Test PDF rendering with different page sizes."""
        doc = Document(children=[Paragraph(content=[Text(content="Test")])])

        # Letter
        letter_renderer = PdfRenderer(PdfRendererOptions(page_size="letter"))
        letter_file = tmp_path / "letter.pdf"
        letter_renderer.render(doc, letter_file)
        assert letter_file.exists()

        # A4
        a4_renderer = PdfRenderer(PdfRendererOptions(page_size="a4"))
        a4_file = tmp_path / "a4.pdf"
        a4_renderer.render(doc, a4_file)
        assert a4_file.exists()

        # Legal
        legal_renderer = PdfRenderer(PdfRendererOptions(page_size="legal"))
        legal_file = tmp_path / "legal.pdf"
        legal_renderer.render(doc, legal_file)
        assert legal_file.exists()

    def test_pdf_with_custom_fonts(self, tmp_path):
        """Test PDF with custom font settings."""
        doc = create_sample_document()
        options = PdfRendererOptions(
            font_name="Times-Roman",
            font_size=12,
            code_font="Courier"
        )
        renderer = PdfRenderer(options)
        output_file = tmp_path / "custom_fonts.pdf"
        renderer.render(doc, output_file)

        assert output_file.exists()


@pytest.mark.integration
class TestCrossRendererConsistency:
    """Tests for consistency across different renderers."""

    def test_same_document_all_renderers(self, tmp_path):
        """Test that same AST renders successfully to all formats."""
        doc = create_sample_document()

        # Markdown
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert len(md_result) > 0

        # HTML
        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert len(html_result) > 0

        # DOCX (if available)
        if DOCX_AVAILABLE:
            docx_renderer = DocxRenderer()
            docx_file = tmp_path / "consistency.docx"
            docx_renderer.render(doc, docx_file)
            assert docx_file.exists()

        # PDF (if available)
        if REPORTLAB_AVAILABLE:
            pdf_renderer = PdfRenderer()
            pdf_file = tmp_path / "consistency.pdf"
            pdf_renderer.render(doc, pdf_file)
            assert pdf_file.exists()

    def test_content_preservation(self):
        """Test that essential content is preserved across renderers."""
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Test Title")]),
            Paragraph(content=[Text(content="Test content")])
        ])

        # Markdown
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "Test Title" in md_result
        assert "Test content" in md_result

        # HTML
        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert "Test Title" in html_result
        assert "Test content" in html_result


@pytest.mark.integration
class TestComplexScenarios:
    """Tests for complex real-world scenarios."""

    def test_deeply_nested_structures(self):
        """Test handling of deeply nested structures."""
        doc = Document(children=[
            List(ordered=False, items=[
                ListItem(children=[
                    Paragraph(content=[Text(content="Item 1")]),
                    List(ordered=False, items=[
                        ListItem(children=[
                            Paragraph(content=[Text(content="Nested 1")])
                        ])
                    ])
                ])
            ])
        ])

        # Should render without errors
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "Item 1" in md_result
        assert "Nested 1" in md_result

        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert "Item 1" in html_result
        assert "Nested 1" in html_result

    def test_mixed_formatting(self):
        """Test paragraph with multiple formatting types."""
        doc = Document(children=[
            Paragraph(content=[
                Text(content="This is "),
                Strong(content=[Text(content="bold")]),
                Text(content=" and this is "),
                Link(url="http://example.com", content=[Text(content="a link")]),
                Text(content=".")
            ])
        ])

        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "**bold**" in md_result
        assert "[a link](http://example.com)" in md_result

        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert "<strong>bold</strong>" in html_result
        assert '<a href="http://example.com">a link</a>' in html_result

    def test_large_table(self, tmp_path):
        """Test handling of large tables."""
        # Create table with many rows
        rows = [
            TableRow(cells=[
                TableCell(content=[Text(content=f"Cell {i}-{j}")])
                for j in range(5)
            ])
            for i in range(20)
        ]

        doc = Document(children=[
            Table(
                header=TableRow(cells=[
                    TableCell(content=[Text(content=f"Col {i}")])
                    for i in range(5)
                ]),
                rows=rows
            )
        ])

        # Should handle large tables
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "Cell 0-0" in md_result
        assert "Cell 19-4" in md_result

        if DOCX_AVAILABLE:
            docx_renderer = DocxRenderer()
            docx_file = tmp_path / "large_table.docx"
            docx_renderer.render(doc, docx_file)
            assert docx_file.exists()


@pytest.mark.integration
class TestErrorHandling:
    """Tests for error handling and edge cases."""

    def test_empty_document_all_renderers(self, tmp_path):
        """Test that empty documents are handled gracefully."""
        doc = Document()

        # Markdown
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert md_result == ""

        # HTML
        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert html_result == ""

        # DOCX (if available)
        if DOCX_AVAILABLE:
            docx_renderer = DocxRenderer()
            docx_file = tmp_path / "empty.docx"
            docx_renderer.render(doc, docx_file)
            assert docx_file.exists()

        # PDF (if available)
        if REPORTLAB_AVAILABLE:
            pdf_renderer = PdfRenderer()
            pdf_file = tmp_path / "empty.pdf"
            pdf_renderer.render(doc, pdf_file)
            assert pdf_file.exists()

    def test_missing_optional_attributes(self):
        """Test handling of missing optional attributes."""
        # Image without title
        doc = Document(children=[
            Paragraph(content=[
                Image(url="image.png", alt_text="Image")
            ])
        ])

        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "![Image](image.png)" in md_result

        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert '<img src="image.png" alt="Image">' in html_result


@pytest.mark.integration
class TestFileOutput:
    """Tests for file output methods."""

    def test_markdown_to_file(self, tmp_path):
        """Test Markdown rendering to file."""
        doc = create_sample_document()
        renderer = MarkdownRenderer()
        output_file = tmp_path / "output.md"
        renderer.render(doc, output_file)

        assert output_file.exists()
        content = output_file.read_text()
        assert "Document Title" in content

    def test_html_to_file(self, tmp_path):
        """Test HTML rendering to file."""
        doc = create_sample_document()
        renderer = HtmlRenderer(HtmlRendererOptions(standalone=True))
        output_file = tmp_path / "output.html"
        renderer.render(doc, output_file)

        assert output_file.exists()
        content = output_file.read_text()
        assert "<!DOCTYPE html>" in content


@pytest.mark.integration
class TestRealWorldDocuments:
    """Tests with realistic document structures."""

    def test_technical_documentation(self, tmp_path):
        """Test rendering technical documentation structure."""
        doc = Document(children=[
            Heading(level=1, content=[Text(content="API Documentation")]),
            Heading(level=2, content=[Text(content="Overview")]),
            Paragraph(content=[Text(content="This API provides access to our services.")]),
            Heading(level=2, content=[Text(content="Endpoints")]),
            Heading(level=3, content=[Text(content="GET /users")]),
            Paragraph(content=[Text(content="Retrieves all users.")]),
            CodeBlock(content='GET /api/users\nAuthorization: Bearer token', language="http"),
            Heading(level=3, content=[Text(content="Response")]),
            CodeBlock(content='{\n  "users": []\n}', language="json"),
        ])

        # Should render successfully to all formats
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "# API Documentation" in md_result

        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=True, include_toc=True))
        html_result = html_renderer.render_to_string(doc)
        assert "API Documentation" in html_result

        if DOCX_AVAILABLE:
            docx_renderer = DocxRenderer()
            docx_file = tmp_path / "api_docs.docx"
            docx_renderer.render(doc, docx_file)
            assert docx_file.exists()

    def test_blog_post_structure(self):
        """Test rendering blog post structure."""
        doc = Document(
            metadata={"title": "My Blog Post", "author": "John Doe"},
            children=[
                Heading(level=1, content=[Text(content="My Blog Post")]),
                Paragraph(content=[Text(content="Posted on January 1, 2025")]),
                Paragraph(content=[
                    Text(content="This is the introduction to my blog post.")
                ]),
                Heading(level=2, content=[Text(content="Main Point")]),
                Paragraph(content=[Text(content="Here's the main content.")]),
                BlockQuote(children=[
                    Paragraph(content=[Text(content="A relevant quote.")])
                ]),
                Heading(level=2, content=[Text(content="Conclusion")]),
                Paragraph(content=[Text(content="Final thoughts.")]),
            ]
        )

        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert "My Blog Post" in md_result

        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=True))
        html_result = html_renderer.render_to_string(doc)
        assert "My Blog Post" in html_result
        assert "<blockquote>" in html_result


@pytest.mark.skipif(not EBOOKLIB_AVAILABLE, reason="ebooklib not installed")
@pytest.mark.integration
class TestEpubRendering:
    """Integration tests for EPUB rendering."""

    def test_full_document_to_epub(self, tmp_path):
        """Test rendering complete document to EPUB."""
        doc = create_sample_document()
        renderer = EpubRenderer(EpubRendererOptions(
            title="Sample Book",
            author="Test Author"
        ))
        output_file = tmp_path / "sample.epub"
        renderer.render(doc, output_file)

        # Verify file created and is valid EPUB
        assert output_file.exists()
        book = epub.read_epub(str(output_file))
        assert book.get_metadata('DC', 'title')[0][0] == "Sample Book"
        assert book.get_metadata('DC', 'creator')[0][0] == "Test Author"

    def test_epub_chapter_splitting_strategies(self, tmp_path):
        """Test different chapter splitting strategies."""
        # Document with both separators and headings
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Chapter 1")]),
            Paragraph(content=[Text(content="Content 1")]),
            ThematicBreak(),
            Heading(level=1, content=[Text(content="Chapter 2")]),
            Paragraph(content=[Text(content="Content 2")])
        ])

        # Test separator mode
        sep_renderer = EpubRenderer(EpubRendererOptions(chapter_split_mode="separator"))
        sep_file = tmp_path / "separator.epub"
        sep_renderer.render(doc, sep_file)
        assert sep_file.exists()

        # Test heading mode
        heading_renderer = EpubRenderer(EpubRendererOptions(chapter_split_mode="heading"))
        heading_file = tmp_path / "heading.epub"
        heading_renderer.render(doc, heading_file)
        assert heading_file.exists()

        # Test auto mode
        auto_renderer = EpubRenderer(EpubRendererOptions(chapter_split_mode="auto"))
        auto_file = tmp_path / "auto.epub"
        auto_renderer.render(doc, auto_file)
        assert auto_file.exists()

    def test_epub_with_toc(self, tmp_path):
        """Test EPUB with table of contents."""
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Introduction")]),
            Paragraph(content=[Text(content="Intro content")]),
            Heading(level=1, content=[Text(content="Main Content")]),
            Paragraph(content=[Text(content="Main content")]),
            Heading(level=1, content=[Text(content="Conclusion")]),
            Paragraph(content=[Text(content="Conclusion")])
        ])

        renderer = EpubRenderer(EpubRendererOptions(
            chapter_split_mode="heading",
            generate_toc=True
        ))
        output_file = tmp_path / "with_toc.epub"
        renderer.render(doc, output_file)

        # Verify TOC was created
        book = epub.read_epub(str(output_file))
        assert len(book.toc) > 0

    def test_epub_metadata_from_document(self, tmp_path):
        """Test EPUB metadata extraction from document."""
        doc = Document(
            metadata={
                "title": "My Document",
                "author": "Jane Smith",
                "subject": "Testing",
                "date": "2025-01-01"
            },
            children=[
                Paragraph(content=[Text(content="Content")])
            ]
        )

        renderer = EpubRenderer()
        output_file = tmp_path / "metadata.epub"
        renderer.render(doc, output_file)

        # Verify metadata
        book = epub.read_epub(str(output_file))
        assert book.get_metadata('DC', 'title')[0][0] == "My Document"
        assert book.get_metadata('DC', 'creator')[0][0] == "Jane Smith"

    def test_epub_complex_content(self, tmp_path):
        """Test EPUB with complex content structures."""
        doc = Document(
            metadata={"title": "Complex EPUB"},
            children=[
                Heading(level=1, content=[Text(content="Chapter 1: Introduction")]),
                Paragraph(content=[
                    Text(content="This is "),
                    Strong(content=[Text(content="important")]),
                    Text(content=".")
                ]),
                List(ordered=True, items=[
                    ListItem(children=[Paragraph(content=[Text(content="First point")])]),
                    ListItem(children=[Paragraph(content=[Text(content="Second point")])])
                ]),
                ThematicBreak(),
                Heading(level=1, content=[Text(content="Chapter 2: Code")]),
                CodeBlock(content='def example():\n    return 42', language="python"),
                ThematicBreak(),
                Heading(level=1, content=[Text(content="Chapter 3: Data")]),
                Table(
                    header=TableRow(cells=[
                        TableCell(content=[Text(content="Name")]),
                        TableCell(content=[Text(content="Value")])
                    ]),
                    rows=[
                        TableRow(cells=[
                            TableCell(content=[Text(content="X")]),
                            TableCell(content=[Text(content="10")])
                        ])
                    ]
                )
            ]
        )

        renderer = EpubRenderer()
        output_file = tmp_path / "complex.epub"
        renderer.render(doc, output_file)

        # Verify EPUB is valid
        book = epub.read_epub(str(output_file))
        assert book is not None


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
@pytest.mark.integration
@pytest.mark.pptx
class TestPptxRendering:
    """Integration tests for PPTX rendering."""

    def test_full_document_to_pptx(self, tmp_path):
        """Test rendering complete document to PPTX."""
        doc = create_sample_document()
        renderer = PptxRenderer()
        output_file = tmp_path / "sample.pptx"
        renderer.render(doc, output_file)

        # Verify file created and is valid PPTX
        assert output_file.exists()
        prs = Presentation(str(output_file))
        assert len(prs.slides) > 0

    def test_pptx_slide_splitting_strategies(self, tmp_path):
        """Test different slide splitting strategies."""
        # Document with both separators and headings
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Slide 1")]),
            Paragraph(content=[Text(content="Content 1")]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Slide 2")]),
            Paragraph(content=[Text(content="Content 2")]),
            ThematicBreak(),
            Paragraph(content=[Text(content="Slide 3")])
        ])

        # Test separator mode (should create 3 slides)
        sep_renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="separator"))
        sep_file = tmp_path / "separator.pptx"
        sep_renderer.render(doc, sep_file)
        sep_prs = Presentation(str(sep_file))
        assert len(sep_prs.slides) == 3

        # Test heading mode (should create 3 slides: H2, H2, content)
        heading_renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="heading"))
        heading_file = tmp_path / "heading.pptx"
        heading_renderer.render(doc, heading_file)
        heading_prs = Presentation(str(heading_file))
        assert len(heading_prs.slides) == 3

        # Test auto mode (prefers separator)
        auto_renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="auto"))
        auto_file = tmp_path / "auto.pptx"
        auto_renderer.render(doc, auto_file)
        auto_prs = Presentation(str(auto_file))
        assert len(auto_prs.slides) == 3

    def test_pptx_slide_titles(self, tmp_path):
        """Test slide title handling."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="First Slide")]),
            Paragraph(content=[Text(content="Content 1")]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Second Slide")]),
            Paragraph(content=[Text(content="Content 2")])
        ])

        renderer = PptxRenderer(PptxRendererOptions(
            slide_split_mode="separator",
            use_heading_as_slide_title=True
        ))
        output_file = tmp_path / "titles.pptx"
        renderer.render(doc, output_file)

        # Verify slides have titles
        prs = Presentation(str(output_file))
        assert prs.slides[0].shapes.title.text == "First Slide"

    def test_pptx_complex_content(self, tmp_path):
        """Test PPTX with complex content structures."""
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Presentation Title")]),
            Paragraph(content=[Text(content="Introduction")]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Key Points")]),
            List(ordered=False, items=[
                ListItem(children=[Paragraph(content=[Text(content="Point 1")])]),
                ListItem(children=[Paragraph(content=[Text(content="Point 2")])]),
                ListItem(children=[Paragraph(content=[Text(content="Point 3")])])
            ]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Code Example")]),
            CodeBlock(content='def hello():\n    print("Hello")', language="python"),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Data")]),
            Table(
                header=TableRow(cells=[
                    TableCell(content=[Text(content="Metric")]),
                    TableCell(content=[Text(content="Value")])
                ]),
                rows=[
                    TableRow(cells=[
                        TableCell(content=[Text(content="Sales")]),
                        TableCell(content=[Text(content="100")])
                    ]),
                    TableRow(cells=[
                        TableCell(content=[Text(content="Revenue")]),
                        TableCell(content=[Text(content="$1000")])
                    ])
                ]
            )
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "complex.pptx"
        renderer.render(doc, output_file)

        # Verify presentation structure
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 4  # 4 slides separated by ThematicBreak

    def test_pptx_formatting_preservation(self, tmp_path):
        """Test that text formatting is preserved in PPTX."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Formatting Test")]),
            Paragraph(content=[
                Text(content="Normal "),
                Strong(content=[Text(content="bold")]),
                Text(content=" and "),
                Emphasis(content=[Text(content="italic")]),
                Text(content=" and "),
                Code(content="code"),
                Text(content=".")
            ])
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "formatting.pptx"
        renderer.render(doc, output_file)

        # Verify file created successfully
        assert output_file.exists()
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.skipif(not (EBOOKLIB_AVAILABLE and PPTX_AVAILABLE), reason="ebooklib or python-pptx not installed")
@pytest.mark.integration
class TestNewRendererConsistency:
    """Test consistency between new renderers and existing ones."""

    def test_same_ast_all_renderers(self, tmp_path):
        """Test that same AST renders successfully to all formats including new ones."""
        doc = create_sample_document()

        # Markdown
        md_renderer = MarkdownRenderer()
        md_result = md_renderer.render_to_string(doc)
        assert len(md_result) > 0

        # HTML
        html_renderer = HtmlRenderer(HtmlRendererOptions(standalone=False))
        html_result = html_renderer.render_to_string(doc)
        assert len(html_result) > 0

        # DOCX (if available)
        if DOCX_AVAILABLE:
            docx_renderer = DocxRenderer()
            docx_file = tmp_path / "all_formats.docx"
            docx_renderer.render(doc, docx_file)
            assert docx_file.exists()

        # PDF (if available)
        if REPORTLAB_AVAILABLE:
            pdf_renderer = PdfRenderer()
            pdf_file = tmp_path / "all_formats.pdf"
            pdf_renderer.render(doc, pdf_file)
            assert pdf_file.exists()

        # EPUB
        epub_renderer = EpubRenderer()
        epub_file = tmp_path / "all_formats.epub"
        epub_renderer.render(doc, epub_file)
        assert epub_file.exists()

        # PPTX
        pptx_renderer = PptxRenderer()
        pptx_file = tmp_path / "all_formats.pptx"
        pptx_renderer.render(doc, pptx_file)
        assert pptx_file.exists()

    def test_markdown_to_epub_to_markdown_roundtrip(self, tmp_path):
        """Test parsing EPUB and rendering back to Markdown."""
        from all2md.parsers.epub import EpubToAstConverter

        # Create EPUB from AST
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Chapter 1")]),
            Paragraph(content=[Text(content="Some content here.")])
        ])

        epub_renderer = EpubRenderer()
        epub_file = tmp_path / "test.epub"
        epub_renderer.render(doc, epub_file)

        # Parse EPUB back to AST
        epub_parser = EpubToAstConverter()
        parsed_doc = epub_parser.parse(epub_file)

        # Render to Markdown
        md_renderer = MarkdownRenderer()
        result = md_renderer.render_to_string(parsed_doc)

        # Verify content preserved
        assert "Chapter 1" in result
        assert "Some content" in result

    def test_markdown_to_pptx_conversion(self, tmp_path):
        """Test converting Markdown-style document to PPTX."""
        # Create document with slide-like structure
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Presentation")]),
            Paragraph(content=[Text(content="Welcome")]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Agenda")]),
            List(ordered=True, items=[
                ListItem(children=[Paragraph(content=[Text(content="Introduction")])]),
                ListItem(children=[Paragraph(content=[Text(content="Main Points")])]),
                ListItem(children=[Paragraph(content=[Text(content="Conclusion")])])
            ])
        ])

        # Render to PPTX
        renderer = PptxRenderer()
        output_file = tmp_path / "converted.pptx"
        renderer.render(doc, output_file)

        # Verify structure
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 2  # Two slides separated by ThematicBreak
