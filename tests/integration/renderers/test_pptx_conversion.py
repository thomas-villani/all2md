#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/integration/renderers/test_pptx_conversion.py
"""Integration tests for PPTX renderer.

Tests cover:
- End-to-end PPTX rendering workflows
- Slide splitting strategies
- Slide title handling
- Complex content structures
- Formatting preservation

"""

import pytest

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from all2md.ast import (
    BlockQuote,
    Code,
    CodeBlock,
    Document,
    Emphasis,
    Heading,
    Link,
    List,
    ListItem,
    Paragraph,
    Strong,
    Table,
    TableCell,
    TableRow,
    Text,
    ThematicBreak,
)
from all2md.options import PptxRendererOptions

if PPTX_AVAILABLE:
    from all2md.renderers.pptx import PptxRenderer


def create_sample_document():
    """Create a sample AST document for testing.

    Returns
    -------
    Document
        A sample document with various elements for testing.

    """
    return Document(
        metadata={"title": "Sample Document", "author": "Test Author"},
        children=[
            Heading(level=1, content=[Text(content="Document Title")]),
            Paragraph(
                content=[
                    Text(content="This is a paragraph with "),
                    Strong(content=[Text(content="bold text")]),
                    Text(content=" and a "),
                    Link(url="https://example.com", content=[Text(content="link")]),
                    Text(content="."),
                ]
            ),
            Heading(level=2, content=[Text(content="Lists")]),
            List(
                ordered=False,
                items=[
                    ListItem(children=[Paragraph(content=[Text(content="First item")])]),
                    ListItem(children=[Paragraph(content=[Text(content="Second item")])]),
                    ListItem(children=[Paragraph(content=[Text(content="Third item")])]),
                ],
            ),
            Heading(level=2, content=[Text(content="Code Example")]),
            CodeBlock(content='def hello():\n    print("Hello, world!")', language="python"),
            Heading(level=2, content=[Text(content="Table")]),
            Table(
                header=TableRow(
                    cells=[TableCell(content=[Text(content="Name")]), TableCell(content=[Text(content="Value")])]
                ),
                rows=[
                    TableRow(
                        cells=[TableCell(content=[Text(content="Alpha")]), TableCell(content=[Text(content="1")])]
                    ),
                    TableRow(cells=[TableCell(content=[Text(content="Beta")]), TableCell(content=[Text(content="2")])]),
                ],
            ),
            Heading(level=2, content=[Text(content="Quote")]),
            BlockQuote(children=[Paragraph(content=[Text(content="This is a blockquote.")])]),
        ],
    )


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
@pytest.mark.integration
@pytest.mark.pptx
class TestPptxRendering:
    """Integration tests for PPTX rendering."""

    def test_full_document_to_pptx(self, tmp_path):
        """Test rendering complete document to PPTX."""
        doc = create_sample_document()
        renderer = PptxRenderer()
        output_file = tmp_path / "sample.pptx"
        renderer.render(doc, output_file)

        # Verify file created and is valid PPTX
        assert output_file.exists()
        prs = Presentation(str(output_file))
        assert len(prs.slides) > 0

    def test_pptx_slide_splitting_strategies(self, tmp_path):
        """Test different slide splitting strategies."""
        # Document with both separators and headings
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide 1")]),
                Paragraph(content=[Text(content="Content 1")]),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Slide 2")]),
                Paragraph(content=[Text(content="Content 2")]),
                ThematicBreak(),
                Paragraph(content=[Text(content="Slide 3")]),
            ]
        )

        # Test separator mode (should create 3 slides)
        sep_renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="separator"))
        sep_file = tmp_path / "separator.pptx"
        sep_renderer.render(doc, sep_file)
        sep_prs = Presentation(str(sep_file))
        assert len(sep_prs.slides) == 3

        # Test heading mode (should create 2 slides: H2 "Slide 1", H2 "Slide 2" with remaining content)
        # Note: In heading mode, content after the last heading is grouped with that heading
        heading_renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="heading"))
        heading_file = tmp_path / "heading.pptx"
        heading_renderer.render(doc, heading_file)
        heading_prs = Presentation(str(heading_file))
        assert len(heading_prs.slides) == 2

        # Test auto mode (prefers separator)
        auto_renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="auto"))
        auto_file = tmp_path / "auto.pptx"
        auto_renderer.render(doc, auto_file)
        auto_prs = Presentation(str(auto_file))
        assert len(auto_prs.slides) == 3

    def test_pptx_slide_titles(self, tmp_path):
        """Test slide title handling."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="First Slide")]),
                Paragraph(content=[Text(content="Content 1")]),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Second Slide")]),
                Paragraph(content=[Text(content="Content 2")]),
            ]
        )

        renderer = PptxRenderer(PptxRendererOptions(slide_split_mode="separator", use_heading_as_slide_title=True))
        output_file = tmp_path / "titles.pptx"
        renderer.render(doc, output_file)

        # Verify slides have titles
        prs = Presentation(str(output_file))
        assert prs.slides[0].shapes.title.text == "First Slide"

    def test_pptx_complex_content(self, tmp_path):
        """Test PPTX with complex content structures."""
        doc = Document(
            children=[
                Heading(level=1, content=[Text(content="Presentation Title")]),
                Paragraph(content=[Text(content="Introduction")]),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Key Points")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Point 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Point 2")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Point 3")])]),
                    ],
                ),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Code Example")]),
                CodeBlock(content='def hello():\n    print("Hello")', language="python"),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Data")]),
                Table(
                    header=TableRow(
                        cells=[TableCell(content=[Text(content="Metric")]), TableCell(content=[Text(content="Value")])]
                    ),
                    rows=[
                        TableRow(
                            cells=[TableCell(content=[Text(content="Sales")]), TableCell(content=[Text(content="100")])]
                        ),
                        TableRow(
                            cells=[
                                TableCell(content=[Text(content="Revenue")]),
                                TableCell(content=[Text(content="$1000")]),
                            ]
                        ),
                    ],
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "complex.pptx"
        renderer.render(doc, output_file)

        # Verify presentation structure
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 4  # 4 slides separated by ThematicBreak

    def test_pptx_formatting_preservation(self, tmp_path):
        """Test that text formatting is preserved in PPTX."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Formatting Test")]),
                Paragraph(
                    content=[
                        Text(content="Normal "),
                        Strong(content=[Text(content="bold")]),
                        Text(content=" and "),
                        Emphasis(content=[Text(content="italic")]),
                        Text(content=" and "),
                        Code(content="code"),
                        Text(content="."),
                    ]
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "formatting.pptx"
        renderer.render(doc, output_file)

        # Verify file created successfully
        assert output_file.exists()
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1
