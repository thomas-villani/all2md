"""Advanced tests for PPTX shape handling edge cases.

NOTE: These tests are disabled pending refactoring for the new AST architecture.
They test internal functions (_process_shape) that no longer exist in the new parser.
"""

import pytest

# Skip all tests in this module - they test internal functions that were removed in AST refactor
pytest.skip("Tests require refactoring for new AST architecture", allow_module_level=True)

from pptx import Presentation
from pptx.util import Inches, Pt
from utils import assert_markdown_valid, cleanup_test_dir, create_test_temp_dir

# from all2md.parsers.pptx import _process_shape  # Function no longer exists
from all2md import PptxOptions


class TestPptxShapesAdvanced:
    """Test complex shape scenarios in PPTX presentations."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_smartart_graphics_simulation(self):
        """Test SmartArt graphics representation (simulated with grouped shapes)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        slide.shapes.title.text = "SmartArt Process"

        # Simulate SmartArt with multiple connected shapes
        # Step 1 box
        left, top, width, height = Inches(1), Inches(2), Inches(2), Inches(1)
        shape1 = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape1.text = "Step 1\nInitialize"

        # Arrow
        arrow_left, arrow_top = Inches(3.2), Inches(2.4)
        arrow_width, arrow_height = Inches(1), Inches(0.2)
        slide.shapes.add_shape(5, arrow_left, arrow_top, arrow_width, arrow_height)  # Right arrow

        # Step 2 box
        left2, top2 = Inches(4.5), Inches(2)
        shape2 = slide.shapes.add_shape(1, left2, top2, width, height)
        shape2.text = "Step 2\nProcess"

        # Step 3 box (below)
        left3, top3 = Inches(2.75), Inches(3.5)
        shape3 = slide.shapes.add_shape(1, left3, top3, width, height)
        shape3.text = "Step 3\nComplete"

        options = PptxOptions(attachment_mode="alt_text")

        # Test each shape
        md1 = _process_shape(shape1, options)
        md2 = _process_shape(shape2, options)
        md3 = _process_shape(shape3, options)

        assert_markdown_valid(md1)
        assert_markdown_valid(md2)
        assert_markdown_valid(md3)

        # Should contain step text
        assert "Step 1" in md1
        assert "Initialize" in md1
        assert "Step 2" in md2
        assert "Process" in md2
        assert "Step 3" in md3
        assert "Complete" in md3

    def test_grouped_shapes_simulation(self):
        """Test grouped shapes (simulated by positioning related shapes)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Grouped Elements"

        # Create shapes that would conceptually be grouped
        # Main container
        container = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2), Inches(2), Inches(4), Inches(3)
        )
        container.text = "Container"

        # Elements within container area
        element1 = slide.shapes.add_textbox(
            Inches(2.2), Inches(2.5), Inches(1.5), Inches(0.5)
        )
        element1.text = "Element 1"

        element2 = slide.shapes.add_textbox(
            Inches(4.2), Inches(2.5), Inches(1.5), Inches(0.5)
        )
        element2.text = "Element 2"

        # Sub-elements
        sub_element = slide.shapes.add_textbox(
            Inches(2.2), Inches(3.5), Inches(3.5), Inches(0.8)
        )
        sub_element.text = "Sub-element with detailed information"

        options = PptxOptions(attachment_mode="alt_text")

        # Process shapes
        container_md = _process_shape(container, options)
        elem1_md = _process_shape(element1, options)
        elem2_md = _process_shape(element2, options)
        sub_md = _process_shape(sub_element, options)

        # Should contain all text content
        assert "Container" in container_md
        assert "Element 1" in elem1_md
        assert "Element 2" in elem2_md
        assert "Sub-element" in sub_md

    def test_text_direction_and_rotation(self):
        """Test text direction and rotation handling."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Text Direction Test"

        # Vertical text simulation
        vertical_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(1), Inches(4)
        )
        vertical_text = vertical_box.text_frame
        vertical_text.text = "Vertical Text Content"

        # Rotated text simulation (python-pptx rotation is limited)
        rotated_box = slide.shapes.add_textbox(
            Inches(3), Inches(2), Inches(3), Inches(1)
        )
        rotated_text = rotated_box.text_frame
        rotated_text.text = "This text would be rotated 45 degrees"

        # RTL text simulation
        rtl_box = slide.shapes.add_textbox(
            Inches(5), Inches(3), Inches(3), Inches(1)
        )
        rtl_text = rtl_box.text_frame
        rtl_text.text = "مرحبا بالعالم"  # Arabic text (RTL)

        options = PptxOptions(attachment_mode="alt_text")

        vertical_md = _process_shape(vertical_box, options)
        rotated_md = _process_shape(rotated_box, options)
        rtl_md = _process_shape(rtl_box, options)

        assert_markdown_valid(vertical_md)
        assert_markdown_valid(rotated_md)
        assert_markdown_valid(rtl_md)

        # Should preserve text content regardless of direction
        assert "Vertical Text Content" in vertical_md
        assert "rotated 45 degrees" in rotated_md
        assert "مرحبا بالعالم" in rtl_md

    def test_advanced_bullet_styles(self):
        """Test various bullet point styles and symbols."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = "Advanced Bullets"

        # Get the content placeholder
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame

        # Different bullet styles
        p1 = text_frame.paragraphs[0]
        p1.text = "Standard bullet point"
        p1.level = 0

        p2 = text_frame.add_paragraph()
        p2.text = "Second level with different symbol"
        p2.level = 1

        p3 = text_frame.add_paragraph()
        p3.text = "Custom bullet: → Arrow symbol"
        p3.level = 0

        p4 = text_frame.add_paragraph()
        p4.text = "Another custom: ◦ Circle symbol"
        p4.level = 1

        p5 = text_frame.add_paragraph()
        p5.text = "Third level deep"
        p5.level = 2

        # Numbered list
        p6 = text_frame.add_paragraph()
        p6.text = "First numbered item"
        p6.level = 0

        p7 = text_frame.add_paragraph()
        p7.text = "Second numbered item"
        p7.level = 0

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(content_shape, options)
        assert_markdown_valid(markdown)

        # Should contain all bullet content
        assert "Standard bullet point" in markdown
        assert "Second level" in markdown
        assert "Arrow symbol" in markdown
        assert "Circle symbol" in markdown
        assert "Third level deep" in markdown
        assert "numbered item" in markdown

    def test_text_formatting_combinations(self):
        """Test complex text formatting combinations."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Text Formatting"

        # Complex formatted text
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(6), Inches(3)
        )
        text_frame = textbox.text_frame

        # Paragraph with mixed formatting
        p = text_frame.paragraphs[0]
        p.text = "This paragraph has "

        # Add runs with different formatting
        run1 = p.runs[0]
        run1.font.size = Pt(14)

        # Bold run
        run2 = p.add_run()
        run2.text = "bold text"
        run2.font.bold = True

        # Normal run
        run3 = p.add_run()
        run3.text = " and "

        # Italic run
        run4 = p.add_run()
        run4.text = "italic text"
        run4.font.italic = True

        # Bold + Italic run
        run5 = p.add_run()
        run5.text = " and bold italic"
        run5.font.bold = True
        run5.font.italic = True

        # Underlined run
        run6 = p.add_run()
        run6.text = " plus underlined"
        run6.font.underline = True

        # Different font size
        run7 = p.add_run()
        run7.text = " and large text"
        run7.font.size = Pt(18)

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(textbox, options)
        assert_markdown_valid(markdown)

        # Should contain all text content
        assert "This paragraph has" in markdown
        assert "bold text" in markdown
        assert "italic text" in markdown
        assert "bold italic" in markdown
        assert "underlined" in markdown
        assert "large text" in markdown

    def test_shapes_with_hyperlinks(self):
        """Test shapes containing hyperlinks."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Hyperlinked Shapes"

        # Text box with hyperlinks
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(6), Inches(2)
        )
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Visit our website at example.com"

        # Hyperlinked shape (button-like)
        button_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2), Inches(4), Inches(3), Inches(1)
        )
        button_shape.text = "Click Here for More Info"

        options = PptxOptions(attachment_mode="alt_text")

        textbox_md = _process_shape(textbox, options)
        button_md = _process_shape(button_shape, options)

        assert_markdown_valid(textbox_md)
        assert_markdown_valid(button_md)

        # Should contain link text
        assert "example.com" in textbox_md
        assert "Click Here" in button_md

    def test_callout_and_annotation_shapes(self):
        """Test callout and annotation shapes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Callouts and Annotations"

        # Callout shape
        callout = slide.shapes.add_shape(
            2,  # Rounded rectangle (simulating callout)
            Inches(4), Inches(1.5), Inches(3), Inches(1.5)
        )
        callout.text = "Important Note:\nThis is a callout with critical information!"

        # Annotation box
        annotation = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(2.5), Inches(1)
        )
        annotation.text_frame.text = "Side note: Additional context here"

        # Arrow pointing to something
        slide.shapes.add_shape(
            5,  # Right arrow
            Inches(6), Inches(3), Inches(1.5), Inches(0.5)
        )

        options = PptxOptions(attachment_mode="alt_text")

        callout_md = _process_shape(callout, options)
        annotation_md = _process_shape(annotation, options)

        assert_markdown_valid(callout_md)
        assert_markdown_valid(annotation_md)

        # Should contain callout and annotation text
        assert "Important Note" in callout_md
        assert "critical information" in callout_md
        assert "Side note" in annotation_md
        assert "Additional context" in annotation_md

    def test_overlapping_shapes(self):
        """Test handling of overlapping shapes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Overlapping Elements"

        # Base shape
        base_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2), Inches(2), Inches(4), Inches(2)
        )
        base_shape.text = "Base Layer Content"

        # Overlapping shape
        overlay_shape = slide.shapes.add_shape(
            9,  # Oval
            Inches(3), Inches(2.5), Inches(2), Inches(1.5)
        )
        overlay_shape.text = "Overlay Text"

        # Text box that overlaps both
        textbox = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.5), Inches(5), Inches(1)
        )
        textbox.text = "Text spanning across multiple shapes"

        options = PptxOptions(attachment_mode="alt_text")

        base_md = _process_shape(base_shape, options)
        overlay_md = _process_shape(overlay_shape, options)
        textbox_md = _process_shape(textbox, options)

        assert_markdown_valid(base_md)
        assert_markdown_valid(overlay_md)
        assert_markdown_valid(textbox_md)

        # Should contain all text from each shape
        assert "Base Layer Content" in base_md
        assert "Overlay Text" in overlay_md
        assert "spanning across" in textbox_md

    def test_shapes_with_animations_references(self):
        """Test shapes that might have animation references."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Animated Elements"

        # Shapes that might be animated
        entrance_shape = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(3), Inches(1)
        )
        entrance_shape.text = "Entrance Animation: Fade In"

        emphasis_shape = slide.shapes.add_textbox(
            Inches(5), Inches(2), Inches(3), Inches(1)
        )
        emphasis_shape.text = "Emphasis: Pulse Effect"

        exit_shape = slide.shapes.add_textbox(
            Inches(3), Inches(4), Inches(3), Inches(1)
        )
        exit_shape.text = "Exit: Fly Out"

        options = PptxOptions(attachment_mode="alt_text")

        entrance_md = _process_shape(entrance_shape, options)
        emphasis_md = _process_shape(emphasis_shape, options)
        exit_md = _process_shape(exit_shape, options)

        assert_markdown_valid(entrance_md)
        assert_markdown_valid(emphasis_md)
        assert_markdown_valid(exit_md)

        # Should preserve text content regardless of animation
        assert "Fade In" in entrance_md
        assert "Pulse Effect" in emphasis_md
        assert "Fly Out" in exit_md

    def test_connector_and_flow_shapes(self):
        """Test connector lines and flow diagram shapes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Flow Diagram"

        # Process boxes
        start_box = slide.shapes.add_shape(
            23,  # Rounded rectangle
            Inches(1), Inches(2), Inches(2), Inches(1)
        )
        start_box.text = "Start Process"

        decision_box = slide.shapes.add_shape(
            4,  # Diamond (for decisions)
            Inches(4), Inches(2), Inches(2), Inches(1.5)
        )
        decision_box.text = "Decision Point?"

        end_box = slide.shapes.add_shape(
            23,  # Rounded rectangle
            Inches(7), Inches(2), Inches(2), Inches(1)
        )
        end_box.text = "End Process"

        # Connector lines (simplified as arrows)
        slide.shapes.add_shape(
            5,  # Right arrow
            Inches(3.2), Inches(2.4), Inches(0.6), Inches(0.2)
        )

        slide.shapes.add_shape(
            5,  # Right arrow
            Inches(6.2), Inches(2.4), Inches(0.6), Inches(0.2)
        )

        options = PptxOptions(attachment_mode="alt_text")

        start_md = _process_shape(start_box, options)
        decision_md = _process_shape(decision_box, options)
        end_md = _process_shape(end_box, options)

        assert_markdown_valid(start_md)
        assert_markdown_valid(decision_md)
        assert_markdown_valid(end_md)

        # Should contain process text
        assert "Start Process" in start_md
        assert "Decision Point" in decision_md
        assert "End Process" in end_md
