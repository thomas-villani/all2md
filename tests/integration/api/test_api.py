#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/integration/test_api.py
"""Integration tests for all2md API functions and full pipeline testing.

This module provides comprehensive integration tests for the main API functions:
- from_ast: Render AST to various formats
- from_markdown: Convert Markdown to other formats
- convert: General format-to-format conversion
- to_markdown: Full conversion pipeline with format detection

Tests cover various input/output types, format combinations, options handling,
transform pipeline integration, and full end-to-end conversion workflows.
"""

import tempfile
from importlib.util import find_spec
from io import BytesIO
from pathlib import Path

import pytest

try:
    import docx
    from docx import Document as DocxDocument
    from pptx import Presentation

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    docx = None
    Presentation = None

try:
    from reportlab.platypus import SimpleDocTemplate  # noqa: F401

    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

EBOOKLIB_AVAILABLE = find_spec("ebooklib") is not None
PPTX_AVAILABLE = find_spec("pptx") is not None

from utils import (
    DocxTestGenerator,
    EmlTestGenerator,
    assert_markdown_valid,
    cleanup_test_dir,
    create_test_temp_dir,
)

from all2md import (
    DocxOptions,
    HtmlOptions,
    MarkdownParserOptions,
    convert,
    from_ast,
    from_markdown,
    to_ast,
    to_markdown,
)
from all2md.ast import (
    BlockQuote,
    CodeBlock,
    Document,
    Emphasis,
    Heading,
    List,
    ListItem,
    Paragraph,
    Strong,
    Table,
    TableCell,
    TableRow,
    Text,
)
from all2md.exceptions import All2MdError
from all2md.options import (
    HtmlRendererOptions,
    MarkdownOptions,
    PdfOptions,
)


def create_sample_ast_document():
    """Create a sample AST document for testing.

    Returns
    -------
    Document
        Sample document with various node types for comprehensive testing.

    """
    return Document(
        metadata={"title": "Test Document", "author": "Test Suite"},
        children=[
            Heading(level=1, content=[Text(content="Main Title")]),
            Paragraph(
                content=[
                    Text(content="This is a paragraph with "),
                    Strong(content=[Text(content="bold")]),
                    Text(content=" and "),
                    Emphasis(content=[Text(content="italic")]),
                    Text(content=" text."),
                ]
            ),
            Heading(level=2, content=[Text(content="Lists")]),
            List(
                ordered=False,
                items=[
                    ListItem(children=[Paragraph(content=[Text(content="First item")])]),
                    ListItem(children=[Paragraph(content=[Text(content="Second item")])]),
                    ListItem(children=[Paragraph(content=[Text(content="Third item")])]),
                ],
            ),
            Heading(level=2, content=[Text(content="Code")]),
            CodeBlock(content='print("Hello, World!")', language="python"),
            Heading(level=2, content=[Text(content="Table")]),
            Table(
                header=TableRow(
                    cells=[TableCell(content=[Text(content="Name")]), TableCell(content=[Text(content="Value")])]
                ),
                rows=[
                    TableRow(
                        cells=[TableCell(content=[Text(content="Item A")]), TableCell(content=[Text(content="100")])]
                    ),
                    TableRow(
                        cells=[TableCell(content=[Text(content="Item B")]), TableCell(content=[Text(content="200")])]
                    ),
                ],
            ),
            BlockQuote(children=[Paragraph(content=[Text(content="A quoted passage.")])]),
        ],
    )


@pytest.mark.integration
class TestFromAst:
    """Integration tests for from_ast function."""

    def test_from_ast_to_markdown_return_string(self):
        """Test rendering AST to Markdown str."""
        doc = create_sample_ast_document()
        result = from_ast(doc, "markdown")

        assert isinstance(result, str)
        assert "# Main Title" in result
        assert "**bold**" in result
        assert "*italic*" in result
        assert "* First item" in result
        assert "```python" in result
        assert "| Name" in result
        assert "> A quoted passage" in result

    def test_from_ast_to_markdown_with_file_output(self, tmp_path):
        """Test rendering AST to Markdown file."""
        doc = create_sample_ast_document()
        output_file = tmp_path / "output.md"

        result = from_ast(doc, "markdown", output=output_file)

        assert result is None
        assert output_file.exists()
        content = output_file.read_text()
        assert "# Main Title" in content
        assert "**bold**" in content

    def test_from_ast_to_html_standalone(self):
        """Test rendering AST to standalone HTML."""
        doc = create_sample_ast_document()
        result = from_ast(doc, "html", renderer_options=HtmlRendererOptions(standalone=True))

        assert isinstance(result, str)
        assert "<!DOCTYPE html>" in result
        assert "<html" in result
        assert "<h1" in result
        assert "Main Title" in result
        assert "<strong>bold</strong>" in result
        assert "<em>italic</em>" in result

    def test_from_ast_to_html_fragment(self):
        """Test rendering AST to HTML fragment."""
        doc = create_sample_ast_document()
        result = from_ast(doc, "html", renderer_options=HtmlRendererOptions(standalone=False))

        assert isinstance(result, str)
        assert "<!DOCTYPE html>" not in result
        assert "<h1" in result
        assert "Main Title" in result

    def test_from_ast_with_renderer_options(self):
        """Test from_ast with various renderer options."""
        doc = create_sample_ast_document()

        # Test with emphasis symbol preference
        result = from_ast(doc, "markdown", renderer_options=MarkdownOptions(emphasis_symbol="_"))

        assert "_italic_" in result

    def test_from_ast_with_transforms(self):
        """Test from_ast with transform pipeline."""
        doc = create_sample_ast_document()

        # Apply remove-images transform
        result = from_ast(doc, "markdown", transforms=["remove-images"])

        assert isinstance(result, str)
        assert "Main Title" in result

    def test_from_ast_to_path_output(self, tmp_path):
        """Test from_ast with Path output."""
        doc = create_sample_ast_document()
        output_file = tmp_path / "test.md"

        result = from_ast(doc, "markdown", output=output_file)

        assert result is None
        assert output_file.exists()

    def test_from_ast_to_io_output(self, tmp_path):
        """Test from_ast with IO output."""
        doc = create_sample_ast_document()
        output_file = tmp_path / "test.md"

        with open(output_file, "w") as f:
            result = from_ast(doc, "markdown", output=f)

        assert result is None
        assert output_file.exists()
        content = output_file.read_text()
        assert "Main Title" in content

    @pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
    def test_from_ast_to_docx(self, tmp_path):
        """Test rendering AST to DOCX format."""
        doc = create_sample_ast_document()
        output_file = tmp_path / "test.docx"

        result = from_ast(doc, "docx", output=output_file)

        assert result is None
        assert output_file.exists()

        # Verify DOCX content
        docx_doc = DocxDocument(str(output_file))
        all_text = " ".join(p.text for p in docx_doc.paragraphs)
        assert "Main Title" in all_text
        assert "bold" in all_text

    @pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
    def test_from_ast_to_docx_bytes(self):
        """Test rendering AST to DOCX returns bytes."""
        doc = create_sample_ast_document()

        result = from_ast(doc, "docx")

        # Should return bytes
        assert isinstance(result, bytes)

        # Verify it's a valid DOCX (ZIP signature)
        assert len(result) > 0
        assert result.startswith(b"PK")  # ZIP signature for DOCX files

    @pytest.mark.skipif(not REPORTLAB_AVAILABLE, reason="reportlab not installed")
    def test_from_ast_to_pdf(self, tmp_path):
        """Test rendering AST to PDF format."""
        doc = create_sample_ast_document()
        output_file = tmp_path / "test.pdf"

        result = from_ast(doc, "pdf", output=output_file)

        assert result is None
        assert output_file.exists()
        assert output_file.stat().st_size > 0

    def test_from_ast_with_kwargs_override(self):
        """Test from_ast with kwargs overriding renderer options."""
        doc = create_sample_ast_document()
        base_options = MarkdownOptions(emphasis_symbol="*")

        result = from_ast(doc, "markdown", renderer_options=base_options, emphasis_symbol="_")  # Override

        assert "_italic_" in result


@pytest.mark.integration
class TestFromMarkdown:
    """Integration tests for from_markdown function."""

    def test_from_markdown_string_to_html(self, tmp_path):
        """Test converting Markdown string to HTML."""
        # Write markdown to file first since from_markdown expects a file path
        md_file = tmp_path / "test.md"
        md_file.write_text("# Hello\n\nThis is **bold** text.")

        result = from_markdown(str(md_file), "html")

        assert isinstance(result, str)
        assert "<h1" in result
        assert "Hello" in result
        assert "<strong>bold</strong>" in result

    def test_from_markdown_file_to_html(self, tmp_path):
        """Test converting Markdown file to HTML."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Test\n\nSome content.")

        result = from_markdown(str(md_file), "html")

        assert isinstance(result, str)
        assert "<h1" in result
        assert "Test" in result

    def test_from_markdown_to_html_with_output_file(self, tmp_path):
        """Test converting Markdown to HTML file."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Document\n\nParagraph text.")
        output_file = tmp_path / "output.html"

        result = from_markdown(str(md_file), "html", output=output_file)

        assert result is None
        assert output_file.exists()
        content = output_file.read_text()
        assert "Document" in content

    @pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
    def test_from_markdown_to_docx(self, tmp_path):
        """Test converting Markdown to DOCX."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Title\n\n**Bold** and *italic*.")
        output_file = tmp_path / "output.docx"

        result = from_markdown(str(md_file), "docx", output=output_file)

        assert result is None
        assert output_file.exists()

        docx_doc = DocxDocument(str(output_file))
        all_text = " ".join(p.text for p in docx_doc.paragraphs)
        assert "Title" in all_text

    def test_from_markdown_with_parser_options(self, tmp_path):
        """Test from_markdown with parser options."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nContent")

        result = from_markdown(str(md_file), "html", parser_options=MarkdownParserOptions(flavor="commonmark"))

        assert isinstance(result, str)
        assert "Title" in result

    def test_from_markdown_with_renderer_options(self, tmp_path):
        """Test from_markdown with renderer options."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nContent")

        result = from_markdown(str(md_file), "html", renderer_options=HtmlRendererOptions(standalone=True))

        assert isinstance(result, str)
        assert "<!DOCTYPE html>" in result

    def test_from_markdown_with_transforms(self, tmp_path):
        """Test from_markdown with transforms applied."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\n![Image](test.png)\n\nContent")

        result = from_markdown(str(md_file), "html", transforms=["remove-images"])

        assert isinstance(result, str)
        assert "Title" in result
        assert "Content" in result

    def test_from_markdown_bytes_input(self, tmp_path):
        """Test from_markdown with bytes input via file."""
        md_file = tmp_path / "test.md"
        md_file.write_bytes(b"# Test\n\nContent")

        result = from_markdown(str(md_file), "html")

        assert isinstance(result, str)
        assert "Test" in result

    def test_from_markdown_path_input(self, tmp_path):
        """Test from_markdown with Path input."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Path Test\n\nContent")

        result = from_markdown(md_file, "html")

        assert isinstance(result, str)
        assert "Path Test" in result

    @pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
    def test_from_markdown_to_docx_bytes(self, tmp_path):
        """Test from_markdown to DOCX returns bytes."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nContent")

        result = from_markdown(str(md_file), "docx")

        # Should return bytes
        assert isinstance(result, bytes)

        # Verify it's a valid DOCX
        assert len(result) > 0
        assert result.startswith(b"PK")  # ZIP signature


@pytest.mark.integration
class TestConvert:
    """Integration tests for convert function."""

    def test_convert_html_to_markdown(self, tmp_path):
        """Test converting HTML to Markdown."""
        html_file = tmp_path / "input.html"
        html_file.write_text("<h1>Title</h1><p>Content with <strong>bold</strong>.</p>")

        result = convert(str(html_file), target_format="markdown")

        assert isinstance(result, str)
        assert "# Title" in result or "Title" in result
        assert "**bold**" in result

    def test_convert_markdown_to_html(self, tmp_path):
        """Test converting Markdown to HTML."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\n**Bold** text.")

        result = convert(str(md_file), target_format="html", source_format="markdown")

        assert isinstance(result, str)
        assert "Title" in result
        assert "Bold" in result or "bold" in result

    def test_convert_with_auto_source_detection(self, tmp_path):
        """Test convert with auto source format detection."""
        html_file = tmp_path / "test.html"
        html_file.write_text("<html><body><h1>Test</h1></body></html>")

        result = convert(str(html_file), target_format="markdown")

        assert isinstance(result, str)
        assert "Test" in result

    def test_convert_with_auto_target_detection(self, tmp_path):
        """Test convert with auto target format detection from output."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Test\n\nContent")
        output_file = tmp_path / "output.html"

        result = convert(str(md_file), output=output_file, source_format="markdown")

        assert result is None
        assert output_file.exists()
        content = output_file.read_text()
        assert "Test" in content

    def test_convert_with_parser_and_renderer_options(self, tmp_path):
        """Test convert with both parser and renderer options."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nContent")

        result = convert(
            str(md_file),
            source_format="markdown",
            target_format="html",
            parser_options=MarkdownParserOptions(),
            renderer_options=HtmlRendererOptions(standalone=False),
        )

        assert isinstance(result, str)
        assert "<!DOCTYPE html>" not in result
        assert "Title" in result

    def test_convert_with_transforms(self, tmp_path):
        """Test convert with transform pipeline."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\n![Image](test.png)\n\nContent")

        result = convert(str(md_file), source_format="markdown", target_format="html", transforms=["remove-images"])

        assert isinstance(result, str)
        assert "Title" in result
        assert "Content" in result

    def test_convert_bytes_input(self):
        """Test convert with bytes input."""
        html_bytes = b"<h1>Title</h1><p>Content</p>"

        result = convert(html_bytes, source_format="html", target_format="markdown")

        assert isinstance(result, str)
        assert "Title" in result

    def test_convert_io_input(self):
        """Test convert with IO input."""
        html_content = b"<h1>Test</h1><p>Content</p>"
        html_io = BytesIO(html_content)

        result = convert(html_io, source_format="html", target_format="markdown")

        assert isinstance(result, str)
        assert "Test" in result

    def test_convert_to_file_output(self, tmp_path):
        """Test convert with file output."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Title\n\nContent")
        output_file = tmp_path / "output.html"

        result = convert(str(md_file), output=output_file, source_format="markdown", target_format="html")

        assert result is None
        assert output_file.exists()

    def test_convert_to_io_output(self, tmp_path):
        """Test convert with IO output."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Title\n\nContent")
        output_file = tmp_path / "output.html"

        with open(output_file, "w") as f:
            result = convert(str(md_file), output=f, source_format="markdown", target_format="html")

        assert result is None
        assert output_file.exists()

    @pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
    def test_convert_markdown_to_docx(self, tmp_path):
        """Test converting Markdown to DOCX."""
        md_file = tmp_path / "input.md"
        md_file.write_text("# Document\n\n**Important** content.")
        output_file = tmp_path / "output.docx"

        result = convert(str(md_file), output=output_file, source_format="markdown", target_format="docx")

        assert result is None
        assert output_file.exists()

        docx_doc = DocxDocument(str(output_file))
        all_text = " ".join(p.text for p in docx_doc.paragraphs)
        assert "Document" in all_text

    def test_convert_with_flavor_shorthand(self, tmp_path):
        """Test convert with flavor shorthand."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nContent")

        result = convert(str(md_file), source_format="markdown", target_format="markdown", flavor="commonmark")

        assert isinstance(result, str)
        assert "Title" in result

    def test_convert_with_kwargs_split(self):
        """Test convert with kwargs split between parser and renderer."""
        html_content = b"<h1>Test</h1><p>Content</p>"

        result = convert(
            html_content,
            source_format="html",
            target_format="markdown",
            convert_nbsp=True,  # HtmlOptions (parser)
            emphasis_symbol="_",  # MarkdownOptions (renderer)
        )

        assert isinstance(result, str)
        assert "Test" in result

    def test_convert_roundtrip_markdown_html(self, tmp_path):
        """Test roundtrip conversion Markdown -> HTML -> Markdown."""
        md_file = tmp_path / "original.md"
        md_file.write_text("# Title\n\n**Bold** and *italic* text.")

        # Markdown -> HTML
        html_file = tmp_path / "temp.html"
        convert(str(md_file), output=html_file, source_format="markdown", target_format="html")
        html_content = html_file.read_text()
        assert "Title" in html_content

        # HTML -> Markdown
        recovered_markdown = convert(str(html_file), source_format="html", target_format="markdown")
        assert isinstance(recovered_markdown, str)
        assert "Title" in recovered_markdown
        assert "**Bold**" in recovered_markdown or "bold" in recovered_markdown.lower()

    @pytest.mark.skipif(not DOCX_AVAILABLE, reason="python-docx not installed")
    def test_convert_docx_to_markdown(self, tmp_path):
        """Test converting DOCX to Markdown."""
        # Create a simple DOCX file
        from docx import Document as CreateDoc

        doc = CreateDoc()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test paragraph.")

        docx_file = tmp_path / "test.docx"
        doc.save(str(docx_file))

        # Convert to Markdown
        result = convert(str(docx_file), target_format="markdown")

        assert isinstance(result, str)
        assert "Test Document" in result
        assert "test paragraph" in result


@pytest.mark.integration
class TestAPIEdgeCases:
    """Test edge cases and error handling for API functions."""

    def test_from_ast_empty_document(self):
        """Test from_ast with empty document."""
        doc = Document()
        result = from_ast(doc, "markdown")

        assert isinstance(result, str)
        assert result == ""

    def test_from_markdown_empty_content(self, tmp_path):
        """Test from_markdown with empty content."""
        md_file = tmp_path / "empty.md"
        md_file.write_text("")

        result = from_markdown(str(md_file), "html")

        assert isinstance(result, str)

    def test_convert_empty_file(self, tmp_path):
        """Test convert with empty file."""
        empty_file = tmp_path / "empty.md"
        empty_file.write_text("")

        result = convert(str(empty_file), source_format="markdown", target_format="markdown")

        assert isinstance(result, str)
        assert result == ""

    def test_from_ast_with_complex_transforms(self):
        """Test from_ast with multiple transforms."""
        doc = create_sample_ast_document()

        result = from_ast(doc, "markdown", transforms=["remove-images", "heading-offset"])

        assert isinstance(result, str)
        # After heading offset, h1 becomes h2
        assert "## Main Title" in result

    def test_convert_with_explicit_formats_override_detection(self, tmp_path):
        """Test that explicit formats override auto-detection."""
        # Create HTML file but process as markdown to avoid txt format issues
        html_file = tmp_path / "test.html"
        html_file.write_text("<h1>HTML</h1>")

        # Convert HTML to markdown (should process HTML properly)
        result = convert(str(html_file), source_format="html", target_format="markdown")

        # Should convert HTML to markdown
        assert isinstance(result, str)
        assert "HTML" in result


@pytest.mark.integration
class TestAPIConsistency:
    """Test consistency across different API functions."""

    def test_to_markdown_vs_convert_consistency(self, tmp_path):
        """Test that to_markdown and convert produce same results."""
        from all2md import to_markdown

        html_file = tmp_path / "test.html"
        html_file.write_text("<html><body><h1>Title</h1><p>Content</p></body></html>")

        result1 = to_markdown(str(html_file))
        result2 = convert(str(html_file), target_format="markdown")

        # Results should be similar (exact match may vary due to rendering details)
        # Both return str
        assert isinstance(result2, str)
        assert "Title" in result1 and "Title" in result2
        assert "Content" in result1 and "Content" in result2

    def test_to_ast_plus_from_ast_vs_to_markdown(self, tmp_path):
        """Test that to_ast + from_ast equals to_markdown."""
        html_file = tmp_path / "test.html"
        html_file.write_text("<html><body><h1>Title</h1><p>Content</p></body></html>")

        # Direct conversion
        from all2md import to_markdown

        direct = to_markdown(str(html_file))

        # Via AST
        ast_doc = to_ast(str(html_file))
        via_ast = from_ast(ast_doc, "markdown")

        # Should be identical
        assert isinstance(via_ast, str)
        assert direct == via_ast

    def test_from_markdown_vs_convert_consistency(self, tmp_path):
        """Test that from_markdown and convert produce same results."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\n**Bold** text.")

        result1 = from_markdown(str(md_file), "html")
        result2 = convert(str(md_file), source_format="markdown", target_format="html")

        # Should be identical
        assert isinstance(result1, str)
        assert isinstance(result2, str)
        assert result1 == result2


@pytest.mark.integration
class TestProgressCallbackIntegration:
    """Test progress callback integration through all API functions."""

    def test_to_markdown_with_progress_callback(self, tmp_path):
        """Test that to_markdown forwards progress_callback through the pipeline."""
        from all2md import to_markdown

        html_file = tmp_path / "test.html"
        html_file.write_text("<html><body><h1>Title</h1><p>Content</p></body></html>")

        events = []

        def progress_handler(event):
            events.append(event)

        # Convert with progress callback
        markdown = to_markdown(str(html_file), progress_callback=progress_handler)

        # Should produce markdown
        assert isinstance(markdown, str)
        assert "Title" in markdown

        # Should have received progress events
        assert len(events) > 0

        # Should have started and finished events
        event_types = [e.event_type for e in events]
        assert "started" in event_types
        assert "finished" in event_types

    def test_to_markdown_with_progress_and_transforms(self, tmp_path):
        """Test progress callback receives transform events in to_markdown."""
        from all2md import to_markdown

        html_file = tmp_path / "test.html"
        html_file.write_text("<html><body><h1>Title</h1><p>Content</p></body></html>")

        events = []

        def progress_handler(event):
            events.append(event)

        # Convert with transforms and progress callback
        markdown = to_markdown(str(html_file), transforms=["remove-images"], progress_callback=progress_handler)

        assert "Title" in markdown

        # Should have received progress events including transform events
        item_done_events = [e for e in events if e.event_type == "item_done"]
        assert len(item_done_events) > 0

    def test_from_ast_with_progress_callback(self):
        """Test that from_ast forwards progress_callback."""
        doc = create_sample_ast_document()
        events = []

        def progress_handler(event):
            events.append(event)

        result = from_ast(doc, "markdown", progress_callback=progress_handler)

        assert isinstance(result, str)
        assert "Main Title" in result

        # Should have received progress events
        assert len(events) > 0
        event_types = [e.event_type for e in events]
        assert "started" in event_types
        assert "finished" in event_types

    def test_convert_with_progress_callback(self, tmp_path):
        """Test that convert forwards progress_callback."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nContent")

        events = []

        def progress_handler(event):
            events.append(event)

        result = convert(
            str(md_file), source_format="markdown", target_format="html", progress_callback=progress_handler
        )

        assert isinstance(result, str)
        assert "Title" in result

        # Should have received progress events
        assert len(events) > 0
        event_types = [e.event_type for e in events]
        assert "started" in event_types
        assert "finished" in event_types

    def test_progress_callback_error_handling(self, tmp_path):
        """Test that progress callback errors don't break conversion."""
        from all2md import to_markdown

        html_file = tmp_path / "test.html"
        html_file.write_text("<h1>Title</h1>")

        def failing_callback(event):
            raise RuntimeError("Callback failed!")

        # Should not raise - callback errors are logged
        markdown = to_markdown(str(html_file), progress_callback=failing_callback)

        # Conversion should complete successfully
        assert "Title" in markdown


@pytest.mark.integration
class TestIntegration:
    """Integration tests for full pipeline testing."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_parse_file_auto_detection(self):
        """Test automatic file format detection."""
        # Create test files of different formats

        # DOCX file
        doc = docx.Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test paragraph.")
        docx_file = self.temp_dir / "test.docx"
        doc.save(str(docx_file))

        # HTML file
        html_content = "<html><body><h1>Test HTML</h1><p>HTML content</p></body></html>"
        html_file = self.temp_dir / "test.html"
        html_file.write_text(html_content)

        # PPTX file
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        pptx_file = self.temp_dir / "test.pptx"
        prs.save(str(pptx_file))

        # EML file
        eml_content = EmlTestGenerator.create_headers_edge_cases()
        eml_file = self.temp_dir / "test.eml"
        eml_file.write_text(eml_content)

        # Test parsing each format
        docx_result = to_markdown(str(docx_file))
        html_result = to_markdown(str(html_file))
        pptx_result = to_markdown(str(pptx_file))
        eml_result = to_markdown(str(eml_file))

        # All should produce valid markdown
        assert_markdown_valid(docx_result)
        assert_markdown_valid(html_result)
        assert_markdown_valid(pptx_result)
        assert_markdown_valid(eml_result)

        # Should contain expected content
        assert "Test Document" in docx_result
        assert "Test HTML" in html_result
        assert "Test Presentation" in pptx_result
        # EML content depends on parsing implementation

    def test_parse_file_with_options(self):
        """Test parse_file with various options."""
        # Create test DOCX
        doc = DocxTestGenerator.create_complex_lists_document()
        docx_file = self.temp_dir / "complex.docx"
        doc.save(str(docx_file))

        # Test with different options (note: parse_file doesn't currently support format-specific options)
        result1 = to_markdown(str(docx_file))
        assert_markdown_valid(result1)

        result2 = to_markdown(str(docx_file))
        assert_markdown_valid(result2)

        # Both should contain list content
        assert "First item" in result1
        assert "First item" in result2

    def test_markdown_options_consistency(self):
        """Test that parse_file works with different content types."""
        # HTML with special characters
        html_content = "<p>Text with * asterisks and _ underscores</p>"
        html_file = self.temp_dir / "special.html"
        html_file.write_text(html_content)

        # Test basic parsing (options not yet supported in parse_file)
        html_result = to_markdown(str(html_file))

        # Should contain the text content
        assert "Text with" in html_result
        assert "asterisks" in html_result
        assert "underscores" in html_result

    def test_attachment_handling_consistency(self):
        """Test that files with images can be parsed consistently."""
        # Create documents with images/attachments

        # DOCX with image simulation
        doc = DocxTestGenerator.create_images_document(self.temp_dir)
        docx_file = self.temp_dir / "images.docx"
        doc.save(str(docx_file))

        # HTML with images
        html_content = """<html><body>
            <h1>HTML with Images</h1>
            <img src="image1.png" alt="Image 1">
            <p>Text between images</p>
            <img src="image2.jpg" alt="Image 2">
        </body></html>"""
        html_file = self.temp_dir / "images.html"
        html_file.write_text(html_content)

        # Test basic parsing (format-specific options not yet supported in parse_file)
        docx_result = to_markdown(str(docx_file))
        html_result = to_markdown(str(html_file))

        assert_markdown_valid(docx_result)
        assert_markdown_valid(html_result)

        # Should contain text content
        assert "HTML with Images" in html_result
        assert "Text between images" in html_result

        # HTML should contain image references (default behavior)
        assert "Image 1" in html_result
        assert "Image 2" in html_result

    def test_complex_document_integration(self):
        """Test complex documents with multiple features."""
        # Create comprehensive test document
        doc = docx.Document()

        # Header
        doc.add_heading("Comprehensive Test Document", level=1)

        # Paragraph with formatting
        p1 = doc.add_paragraph("This document tests ")
        run1 = p1.add_run("bold")
        run1.bold = True
        p1.add_run(" and ")
        run2 = p1.add_run("italic")
        run2.italic = True
        p1.add_run(" formatting.")

        # Lists
        doc.add_paragraph("First list item", style="List Bullet")
        doc.add_paragraph("Second list item", style="List Bullet")
        doc.add_paragraph("Nested item", style="List Bullet")

        # Table
        table = doc.add_table(rows=2, cols=3)
        table.rows[0].cells[0].text = "Column A"
        table.rows[0].cells[1].text = "Column B"
        table.rows[0].cells[2].text = "Column C"
        table.rows[1].cells[0].text = "Data 1"
        table.rows[1].cells[1].text = "Data 2"
        table.rows[1].cells[2].text = "Data 3"

        # Special characters
        doc.add_paragraph("Special chars: * _ # [ ] ( ) { }")

        docx_file = self.temp_dir / "comprehensive.docx"
        doc.save(str(docx_file))

        # Test with different option combinations
        parser_options1 = DocxOptions(preserve_tables=True)
        renderer_options1 = MarkdownOptions(escape_special=True)

        result1 = to_markdown(str(docx_file), parser_options=parser_options1, renderer_options=renderer_options1)
        assert_markdown_valid(result1)

        result2 = to_markdown(str(docx_file))
        assert_markdown_valid(result2)

        # Both should contain all elements
        for result in [result1, result2]:
            assert "# Comprehensive Test Document" in result
            assert "**bold**" in result
            assert "*italic*" in result
            assert "First list item" in result
            assert "| Column A | Column B | Column C |" in result
            assert "Special chars" in result

    def test_error_handling_integration(self):
        """Test error handling in integration scenarios."""
        # Test with non-existent file
        with pytest.raises(FileNotFoundError):
            with open("/non/existent/file.docx", "rb") as f:
                to_markdown(f)

        # Test with unsupported file type - should fall back to text parsing
        unsupported_file = self.temp_dir / "test.xyz"
        unsupported_file.write_text("Unsupported content")

        result = to_markdown(str(unsupported_file))
        # Should fall back to treating as plain text
        assert result == "Unsupported content"

        # Test with corrupted file
        corrupted_file = self.temp_dir / "corrupted.docx"
        corrupted_file.write_bytes(b"Not a real DOCX file")

        # Should handle gracefully or raise appropriate error
        try:
            result = to_markdown(str(corrupted_file))
            # If it doesn't raise an error, result should be a string
            assert isinstance(result, str)
        except Exception as e:
            # Should be a meaningful error (including custom All2MdError)
            from all2md.exceptions import MalformedFileError

            assert isinstance(e, (ValueError, IOError, OSError, MalformedFileError, All2MdError))

    def test_performance_with_large_documents(self):
        """Test performance and handling of large documents."""
        # Create larger document
        doc = docx.Document()
        doc.add_heading("Large Document Test", level=1)

        # Add many paragraphs
        for i in range(100):
            doc.add_paragraph(f"This is paragraph {i + 1} with some content to make it longer.")

        # Add large table
        table = doc.add_table(rows=20, cols=5)
        for row_idx in range(20):
            for col_idx in range(5):
                table.rows[row_idx].cells[col_idx].text = f"Cell {row_idx},{col_idx}"

        # Add many list items
        for i in range(50):
            doc.add_paragraph(f"List item {i + 1}", style="List Bullet")

        large_file = self.temp_dir / "large.docx"
        doc.save(str(large_file))

        # Should handle large document without issues
        result = to_markdown(str(large_file))
        assert_markdown_valid(result)

        # Should contain first and last content
        assert "Large Document Test" in result
        assert "paragraph 1 " in result
        assert "paragraph 100" in result
        assert "List item 1" in result
        assert "List item 50" in result

    def test_file_like_object_support(self):
        """Test support for file-like objects."""
        # Test with StringIO for HTML
        html_content = b"<html><body><h1>StringIO Test</h1><p>Content from StringIO</p></body></html>"
        html_io = BytesIO(html_content)

        # parse_file requires both file object and filename
        # StringIO doesn't have a name, so we simulate it
        try:
            result = to_markdown(html_io)
            assert_markdown_valid(result)
            assert "StringIO Test" in result
        except (TypeError, AttributeError):
            # If not supported, that's also valid
            pass

    def test_batch_processing_simulation(self):
        """Test processing multiple files in sequence."""
        files_and_results = []

        # Create multiple test files
        for i in range(5):
            # HTML files
            html_content = f"<html><body><h1>Document {i + 1}</h1><p>Content {i + 1}</p></body></html>"
            html_file = self.temp_dir / f"doc{i + 1}.html"
            html_file.write_text(html_content)

            result = to_markdown(str(html_file))
            files_and_results.append((str(html_file), result))

        # All should be processed successfully
        assert len(files_and_results) == 5

        for file_path, result in files_and_results:
            assert_markdown_valid(result)
            # Extract number from filename to verify content
            file_num = Path(file_path).stem[-1]
            assert f"Document {file_num}" in result
            assert f"Content {file_num}" in result

    def test_option_inheritance_and_precedence(self):
        """Test basic HTML parsing through parse_file."""
        # Create simple test HTML (avoiding features that trigger bugs)
        html_content = """<html><body>
            <h1>Option Test</h1>
            <p>Text with bold and italic</p>
        </body></html>"""
        html_file = self.temp_dir / "options.html"
        html_file.write_text(html_content)

        # Test basic HTML parsing (format-specific options not yet supported in parse_file)
        result = to_markdown(str(html_file))
        assert_markdown_valid(result)

        # Should contain heading and text
        assert "Option Test" in result
        assert "Text with" in result


@pytest.mark.integration
class TestNewAPI:
    """Test the new to_markdown API with enhanced format detection and options handling."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_explicit_format_parameter(self):
        """Test explicit format specification bypasses detection."""
        # Create HTML content but force it to be processed as text
        html_content = b"<html><body><h1>Test</h1><p>Content</p></body></html>"
        html_io = BytesIO(html_content)

        # Process as HTML (auto-detection)
        result_html = to_markdown(html_io, source_format="html")
        assert "# Test" in result_html or "Test" in result_html
        assert "Content" in result_html
        assert "<html>" not in result_html  # Should be converted

        # Reset and process as plain text (force format)
        html_io.seek(0)
        result_text = to_markdown(html_io, source_format="plaintext")
        assert "<html>" in result_text  # Should preserve HTML tags
        assert "<body>" in result_text

    def test_options_parameter_with_kwargs_override(self):
        """Test that kwargs override options parameter values."""
        html_content = b"<html><body><p>Test content</p></body></html>"
        html_io = BytesIO(html_content)

        # Create options with one setting
        md_opts = MarkdownOptions(use_hash_headings=False)
        base_parser_options = HtmlOptions(convert_nbsp=True)

        # Override with kwargs
        result = to_markdown(
            html_io,
            parser_options=base_parser_options,
            renderer_options=md_opts,
            use_hash_headings=True,  # Override the False setting
            source_format="html",
        )

        assert_markdown_valid(result)

    def test_kwargs_only_options_creation(self):
        """Test creating options from kwargs alone."""
        html_content = b"<html><body><h1>Header</h1><p>Content</p></body></html>"
        html_io = BytesIO(html_content)

        result = to_markdown(
            html_io,
            source_format="html",
            use_hash_headings=True,
            convert_nbsp=False,
            emphasis_symbol="_",  # MarkdownOptions field
        )

        assert_markdown_valid(result)
        assert "Header" in result

    def test_content_based_format_detection(self):
        """Test format detection from file content when filename unavailable."""
        # PDF magic bytes
        pdf_content = b"%PDF-1.4\ntest content"
        pdf_io = BytesIO(pdf_content)

        # Should detect as PDF even without filename
        # Note: This will fail conversion due to invalid PDF, but format detection should work
        try:
            result = to_markdown(pdf_io)
        except Exception:
            pass  # Expected - invalid PDF content

        # Test HTML detection
        html_content = b"<!DOCTYPE html><html><body><h1>Test</h1></body></html>"
        html_io = BytesIO(html_content)

        result = to_markdown(html_io)
        assert_markdown_valid(result)
        assert "Test" in result

    def test_filename_based_detection_with_file_objects(self):
        """Test format detection from filename attribute on file objects."""
        html_content = b"<html><body><h1>Filename Test</h1></body></html>"

        # Create a BytesIO with a name attribute
        html_io = BytesIO(html_content)
        html_io.name = "test.html"

        result = to_markdown(html_io)
        assert_markdown_valid(result)
        assert "Filename Test" in result

    def test_format_specific_options_mapping(self):
        """Test that format-specific options are properly handled."""
        # Test with PDF options (even if we can't fully test PDF conversion here)
        pdf_options = PdfOptions(pages=[0], attachment_mode="alt_text")

        # This should work with the options system even if PDF processing fails
        simple_text = b"Simple text content"
        text_io = BytesIO(simple_text)

        result = to_markdown(
            text_io,
            parser_options=pdf_options,
            source_format="plaintext",  # Force as text since we don't have real PDF
        )

        assert result.strip() == "Simple text content"

    def test_backward_compatibility(self):
        """Test that old API usage still works."""
        # Old style: just file path
        html_content = "<html><body><h1>Compatibility</h1></body></html>"
        html_file = self.temp_dir / "compat.html"
        html_file.write_text(html_content)

        result = to_markdown(str(html_file))
        assert_markdown_valid(result)
        assert "Compatibility" in result

        # Old style: file object
        html_io = BytesIO(html_content.encode())
        result = to_markdown(html_io)
        assert_markdown_valid(result)

    def test_multiple_format_detections(self):
        """Test detection for various file formats."""
        # Test CSV detection - CSV has its own converter
        csv_content = b"name,age,city\nJohn,25,NYC\nJane,30,LA"
        csv_io = BytesIO(csv_content)
        csv_io.name = "test.csv"  # Give it a name for format detection

        result = to_markdown(csv_io, source_format="csv")
        assert "|" in result  # Should be markdown table
        assert "John" in result
        assert "Jane" in result

        # Test TSV detection - TSV is handled by CSV converter
        tsv_content = b"name\tage\tcity\nBob\t35\tSF"
        tsv_io = BytesIO(tsv_content)
        tsv_io.name = "test.tsv"  # Give it a name for format detection

        result = to_markdown(tsv_io, source_format="csv")
        assert "|" in result  # Should be markdown table
        assert "Bob" in result

    def test_markdown_options_inheritance(self):
        """Test that MarkdownOptions are properly handled across formats."""
        html_content = b"<html><body><em>italic</em> <strong>bold</strong></body></html>"
        html_io = BytesIO(html_content)

        result = to_markdown(
            html_io,
            source_format="html",
            emphasis_symbol="_",
            bullet_symbols="*-+",  # Should affect italic rendering
        )

        assert_markdown_valid(result)
        # Check that emphasis symbol preference is handled

    def test_comprehensive_detection_strategy(self):
        """Test the comprehensive detection with multiple fallback layers."""
        # Test 1: Extension + MIME type detection for uncommon extension
        with tempfile.NamedTemporaryFile(suffix=".unknown", delete=False) as tmp:
            html_content = b"<html><body><h1>Unknown Extension</h1></body></html>"
            tmp.write(html_content)
            tmp.flush()

            # Should fall back to content detection
            result = to_markdown(tmp.name)
            assert_markdown_valid(result)
            assert "Unknown Extension" in result

        # Test 2: MIME type detection for files with known MIME but unknown extension
        csv_like_content = b"name,age\nBob,25\nAlice,30"
        csv_io = BytesIO(csv_like_content)

        # Force detection through content analysis
        result = to_markdown(csv_io)
        assert "|" in result  # Should detect as CSV and create table
        assert "Bob" in result

    def test_detection_logging_and_priority(self):
        """Test that detection methods are tried in correct priority order."""
        import logging

        # Enable debug logging to capture detection flow
        logger = logging.getLogger("all2md")
        logger.setLevel(logging.DEBUG)

        # Create in-memory log capture
        log_capture = []
        handler = logging.Handler()
        handler.emit = lambda record: log_capture.append(record.getMessage())
        logger.addHandler(handler)

        try:
            # Test filename detection (should be first and succeed)
            html_file = self.temp_dir / "test.html"
            html_file.write_text("<html><body><h1>Priority Test</h1></body></html>")

            to_markdown(str(html_file))

            # Check that logs show filename detection succeeded
            # The actual log message is "Format detected from filename: html"
            log_messages = [msg for msg in log_capture if "Format detected from" in msg]
            assert any("html" in msg for msg in log_messages), f"Log messages: {log_messages}"

        finally:
            logger.removeHandler(handler)
            logger.setLevel(logging.WARNING)  # Reset to default

    def test_edge_case_format_detection(self):
        """Test detection for edge cases and unusual file types."""
        # Test with empty file
        empty_io = BytesIO(b"")
        result = to_markdown(empty_io, source_format="plaintext")
        assert result == ""

        # Test with binary data that doesn't match any signature
        binary_data = b"\x00\x01\x02\x03\x04random binary data"
        binary_io = BytesIO(binary_data)

        try:
            result = to_markdown(binary_io)
            # Should fallback to text and likely fail gracefully
        except (All2MdError, UnicodeDecodeError):
            pass  # Expected for binary data

        # Test RTF detection
        rtf_content = b"{\\rtf1\\ansi\\hello world}"
        rtf_io = BytesIO(rtf_content)

        # Should detect as RTF from magic bytes
        try:
            result = to_markdown(rtf_io)
            # RTF conversion may fail due to missing dependencies, but detection should work
        except ImportError:
            pass  # Expected if pyth not installed

    def test_mime_type_fallback_enhancement(self):
        """Test that MIME type detection works as fallback."""
        # Create a file with unusual extension but recognizable MIME type
        pdf_like_file = self.temp_dir / "document.weird"
        # Write something that would trigger PDF MIME detection
        pdf_like_file.write_text("Not really PDF but has PDF extension mapping")

        # The comprehensive detection should try multiple methods
        try:
            result = to_markdown(str(pdf_like_file))
            # Should fall back to text since it's not really a PDF
            assert isinstance(result, str)
        except Exception:
            pass  # Expected for invalid format data
