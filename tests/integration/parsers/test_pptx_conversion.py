"""Integration tests for PPTX to Markdown conversion."""

import pytest
from pptx import Presentation
from pptx.util import Inches

from all2md import to_ast, to_markdown
from all2md.ast.nodes import Document


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_basic(tmp_path):
    """Test basic PPTX to Markdown conversion."""
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Test Presentation"
    subtitle.text = "A simple test presentation"

    # Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Main Content"
    tf = body.text_frame
    tf.text = "First bullet point"

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Test Presentation" in result
    assert "Main Content" in result
    assert "First bullet point" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_bullet_lists(tmp_path):
    """Test PPTX with bullet lists conversion."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Bullet List Example"
    tf = body.text_frame
    tf.text = "First bullet"

    p = tf.add_paragraph()
    p.text = "Second bullet"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Third bullet"
    p.level = 0

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Bullet List Example" in result
    assert "First bullet" in result
    assert "Second bullet" in result
    assert "Third bullet" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_nested_bullets(tmp_path):
    """Test PPTX with nested bullet lists conversion."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Nested Bullets"
    tf = body.text_frame
    tf.text = "Main point 1"

    p = tf.add_paragraph()
    p.text = "Sub-point 1.1"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Sub-point 1.2"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Main point 2"
    p.level = 0

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Main point 1" in result
    assert "Sub-point 1.1" in result
    assert "Sub-point 1.2" in result
    assert "Main point 2" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_multiple_slides(tmp_path):
    """Test PPTX with multiple slides conversion."""
    prs = Presentation()

    # Slide 1
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Slide 1 Title"

    # Slide 2
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Slide 2 Title"
    body.text_frame.text = "Slide 2 content"

    # Slide 3
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Slide 3 Title"
    body.text_frame.text = "Slide 3 content"

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Slide 1 Title" in result
    assert "Slide 2 Title" in result
    assert "Slide 3 Title" in result
    assert "Slide 2 content" in result
    assert "Slide 3 content" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_text_boxes(tmp_path):
    """Test PPTX with text boxes conversion."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(4)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "This is a text box"

    p = text_frame.add_paragraph()
    p.text = "Second paragraph in text box"

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "This is a text box" in result
    assert "Second paragraph in text box" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_tables(tmp_path):
    """Test PPTX with tables conversion."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add table
    rows = 3
    cols = 3
    left = Inches(1)
    top = Inches(1)
    width = Inches(5)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set header row
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(0, 2).text = "City"

    # Set data rows
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"
    table.cell(1, 2).text = "NYC"

    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "25"
    table.cell(2, 2).text = "LA"

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Name" in result
    assert "Alice" in result
    assert "Bob" in result
    assert "NYC" in result
    assert "LA" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_section_headers(tmp_path):
    """Test PPTX with section header slides."""
    prs = Presentation()

    # Section header slide
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Section 1: Introduction"

    # Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Content Slide"
    body.text_frame.text = "Section content"

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Section 1: Introduction" in result
    assert "Content Slide" in result
    assert "Section content" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_notes(tmp_path):
    """Test PPTX with speaker notes conversion."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Slide with Notes"

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are speaker notes for the slide."

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    # Check that slide title is present
    assert "Slide with Notes" in result
    # Speaker notes may or may not be included depending on parser options


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_empty_presentation(tmp_path):
    """Test empty PPTX conversion."""
    prs = Presentation()

    pptx_file = tmp_path / "empty.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    # Should complete without error, may produce minimal output
    assert isinstance(result, str)


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_blank_slides(tmp_path):
    """Test PPTX with blank slides."""
    prs = Presentation()

    # Add a blank slide
    blank_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_layout)

    # Add another blank slide
    prs.slides.add_slide(blank_layout)

    pptx_file = tmp_path / "blank.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    # Should complete without error
    assert isinstance(result, str)


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_mixed_layouts(tmp_path):
    """Test PPTX with mixed slide layouts."""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Title Slide"

    # Title and content
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Content Slide"

    # Section header
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "Section Header"

    # Two content
    if len(prs.slide_layouts) > 3:
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        slide.shapes.title.text = "Two Content Slide"

    pptx_file = tmp_path / "mixed.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Title Slide" in result
    assert "Content Slide" in result
    assert "Section Header" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_ast_conversion(tmp_path):
    """Test PPTX to AST conversion pipeline."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "AST Test Slide"
    body.text_frame.text = "Testing AST conversion"

    pptx_file = tmp_path / "test.pptx"
    prs.save(str(pptx_file))

    doc = to_ast(pptx_file)

    # Verify AST structure
    assert isinstance(doc, Document)
    assert doc.children is not None
    assert len(doc.children) > 0

    # Verify content through markdown conversion
    result = to_markdown(pptx_file)
    assert "AST Test Slide" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_unicode_content(tmp_path):
    """Test PPTX with Unicode characters conversion."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Unicode Test \U0001f600"
    tf = body.text_frame
    tf.text = "Chinese: \U00004e2d\U00006587"
    p = tf.add_paragraph()
    p.text = "Greek: \U00000391\U000003b1"
    p = tf.add_paragraph()
    p.text = "Math: \U0000221e \U000000b1"

    pptx_file = tmp_path / "unicode.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Unicode Test" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_long_content(tmp_path):
    """Test PPTX with long text content."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Long Content Test"
    tf = body.text_frame

    # Add many paragraphs
    for i in range(20):
        if i == 0:
            tf.text = f"Paragraph {i + 1}: This is a long paragraph with lots of text content."
        else:
            p = tf.add_paragraph()
            p.text = f"Paragraph {i + 1}: This is a long paragraph with lots of text content."
            p.level = 0

    pptx_file = tmp_path / "long.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Long Content Test" in result
    assert "Paragraph 1:" in result
    assert "Paragraph 20:" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_special_characters(tmp_path):
    """Test PPTX with special characters."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Special Characters: & < > \" '"
    body.text_frame.text = "Content with *asterisks* and _underscores_"

    pptx_file = tmp_path / "special.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Special Characters" in result
    assert "asterisks" in result
    assert "underscores" in result


@pytest.mark.integration
@pytest.mark.pptx
def test_pptx_to_markdown_numbered_lists(tmp_path):
    """Test PPTX with numbered lists conversion."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Numbered List"
    tf = body.text_frame
    tf.text = "First item"

    for i in range(2, 5):
        p = tf.add_paragraph()
        p.text = f"Item {i}"
        p.level = 0

    pptx_file = tmp_path / "numbered.pptx"
    prs.save(str(pptx_file))

    result = to_markdown(pptx_file)

    assert "Numbered List" in result
    assert "First item" in result
    assert "Item 2" in result
