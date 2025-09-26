"""Advanced tests for PPTX chart handling edge cases."""

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from all2md.converters.pptx2markdown import _process_shape
from all2md.options import PptxOptions
from tests.utils import assert_markdown_valid, cleanup_test_dir, create_test_temp_dir


class TestPptxChartsAdvanced:
    """Test complex chart scenarios in PPTX presentations."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_multiple_series_column_chart(self):
        """Test column chart with multiple data series."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        slide.shapes.title.text = "Multiple Series Chart"

        # Create chart data with multiple series
        chart_data = ChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Revenue', (10.5, 15.2, 12.8, 18.9))
        chart_data.add_series('Expenses', (8.2, 12.1, 10.5, 14.3))
        chart_data.add_series('Profit', (2.3, 3.1, 2.3, 4.6))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain all series and categories
        assert "Revenue" in markdown
        assert "Expenses" in markdown
        assert "Profit" in markdown
        assert "Q1" in markdown
        assert "Q4" in markdown

        # Should include all data values
        assert "10.5" in markdown
        assert "18.9" in markdown

    def test_pie_chart_with_percentages(self):
        """Test pie chart data representation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Market Share"

        chart_data = ChartData()
        chart_data.categories = ['Company A', 'Company B', 'Company C', 'Others']
        chart_data.add_series('Market Share', (35, 25, 20, 20))

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain categories and values
        assert "Company A" in markdown
        assert "Company B" in markdown
        assert "Company C" in markdown
        assert "Others" in markdown
        assert "35" in markdown
        assert "25" in markdown

    def test_line_chart_with_trends(self):
        """Test line chart with multiple trend lines."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Trend Analysis"

        chart_data = ChartData()
        chart_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
        chart_data.add_series('Sales', (10, 12, 8, 15, 18, 22))
        chart_data.add_series('Target', (12, 14, 16, 18, 20, 22))
        chart_data.add_series('Previous Year', (8, 10, 12, 11, 14, 16))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain all series and time periods
        assert "Sales" in markdown
        assert "Target" in markdown
        assert "Previous Year" in markdown
        assert "Jan" in markdown
        assert "Jun" in markdown

    def test_stacked_column_chart(self):
        """Test stacked column chart representation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Stacked Performance"

        chart_data = ChartData()
        chart_data.categories = ['Team A', 'Team B', 'Team C']
        chart_data.add_series('Completed', (80, 65, 90))
        chart_data.add_series('In Progress', (15, 25, 8))
        chart_data.add_series('Not Started', (5, 10, 2))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should represent stacked data appropriately
        assert "Completed" in markdown
        assert "In Progress" in markdown
        assert "Not Started" in markdown
        assert "Team A" in markdown
        assert "Team C" in markdown

    def test_bar_chart_horizontal(self):
        """Test horizontal bar chart."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Horizontal Comparison"

        chart_data = ChartData()
        chart_data.categories = ['Product 1', 'Product 2', 'Product 3', 'Product 4']
        chart_data.add_series('Sales', (25, 40, 30, 35))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain products and values
        assert "Product 1" in markdown
        assert "Product 4" in markdown
        assert "25" in markdown
        assert "40" in markdown

    def test_area_chart(self):
        """Test area chart data representation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Area Chart Analysis"

        chart_data = ChartData()
        chart_data.categories = ['Week 1', 'Week 2', 'Week 3', 'Week 4']
        chart_data.add_series('Visitors', (1000, 1200, 1100, 1400))
        chart_data.add_series('Conversions', (50, 65, 48, 72))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain area chart data
        assert "Visitors" in markdown
        assert "Conversions" in markdown
        assert "Week 1" in markdown
        assert "1000" in markdown
        assert "1400" in markdown

    def test_scatter_plot_chart(self):
        """Test scatter plot chart."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Scatter Plot Analysis"

        # Scatter charts need XyChartData, not ChartData with categories
        from pptx.chart.data import XyChartData

        chart_data = XyChartData()
        series = chart_data.add_series('Data Points')
        series.add_data_point(10, 5)  # Point 1
        series.add_data_point(20, 15)  # Point 2
        series.add_data_point(15, 12)  # Point 3
        series.add_data_point(25, 20)  # Point 4

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain scatter data
        assert "Data Points" in markdown
        # Should contain the x,y values
        assert "10" in markdown
        assert "5" in markdown
        assert "25" in markdown
        assert "20" in markdown

    def test_combination_chart(self):
        """Test combination chart with mixed chart types."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Combination Chart"

        # Create base chart
        chart_data = ChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Volume', (100, 150, 120, 180))
        chart_data.add_series('Revenue %', (5.2, 6.8, 4.9, 7.1))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should contain both series with different scales
        assert "Volume" in markdown
        assert "Revenue %" in markdown
        assert "100" in markdown
        assert "5.2" in markdown

    def test_chart_with_missing_data(self):
        """Test chart handling with missing or zero data points."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Chart with Gaps"

        chart_data = ChartData()
        chart_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May']
        # Include zeros and potentially None values
        chart_data.add_series('Sales', (10, 0, 15, 0, 20))
        chart_data.add_series('Returns', (1, 2, 0, 1, 0))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should handle zero values
        assert "Sales" in markdown
        assert "Returns" in markdown
        assert "0" in markdown
        assert "20" in markdown

    def test_chart_with_long_labels(self):
        """Test chart with very long category or series labels."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Long Labels Chart"

        chart_data = ChartData()
        chart_data.categories = [
            'Very Long Category Name That Exceeds Normal Length',
            'Another Extremely Long Category Name',
            'Short',
            'Medium Length Category'
        ]
        chart_data.add_series(
            'Extremely Long Series Name That Might Cause Formatting Issues',
            (25, 40, 35, 30)
        )

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should handle long labels
        assert "Very Long Category Name" in markdown
        assert "Extremely Long Series Name" in markdown
        assert "Short" in markdown

    def test_chart_with_special_characters(self):
        """Test chart with special characters in labels and data."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Special Characters Chart"

        chart_data = ChartData()
        chart_data.categories = ['Café & Restaurant', 'H₂O Solutions', 'α-Beta Corp', '100% Organic']
        chart_data.add_series('Revenue (€)', (1000.50, 2500.75, 1800.25, 3200.00))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        # Should handle special characters
        assert "Café & Restaurant" in markdown
        assert "H₂O Solutions" in markdown
        assert "α-Beta Corp" in markdown
        assert "100% Organic" in markdown
        assert "Revenue (€)" in markdown

    def test_multiple_charts_on_slide(self):
        """Test slide with multiple charts."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Multiple Charts"

        # First chart
        chart_data1 = ChartData()
        chart_data1.categories = ['A', 'B', 'C']
        chart_data1.add_series('Series 1', (10, 20, 15))

        chart1 = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1.5), Inches(4), Inches(3),
            chart_data1
        )

        # Second chart
        chart_data2 = ChartData()
        chart_data2.categories = ['X', 'Y', 'Z']
        chart_data2.add_series('Series 2', (5, 15, 25))

        chart2 = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(5.5), Inches(1.5), Inches(4), Inches(3),
            chart_data2
        )

        options = PptxOptions(attachment_mode="alt_text")

        # Process each chart separately
        markdown1 = _process_shape(chart1, options)
        markdown2 = _process_shape(chart2, options)

        assert_markdown_valid(markdown1)
        assert_markdown_valid(markdown2)

        # Each should contain their respective data
        assert "Series 1" in markdown1
        assert "Series 2" in markdown2
        assert "A" in markdown1
        assert "X" in markdown2

    def test_chart_error_handling(self):
        """Test error handling with malformed chart data."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Create chart with minimal data
        chart_data = ChartData()
        chart_data.categories = ['Single']
        chart_data.add_series('Data', (42,))

        x, y, cx, cy = Inches(1), Inches(2), Inches(4), Inches(3)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

        options = PptxOptions(attachment_mode="alt_text")

        # Should handle minimal data gracefully
        markdown = _process_shape(chart_shape, options)
        assert_markdown_valid(markdown)

        assert "Single" in markdown
        assert "42" in markdown
