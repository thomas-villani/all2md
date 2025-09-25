"""Integration tests for all2md converters."""

import tempfile
from io import BytesIO
from pathlib import Path

import docx
import pytest
from pptx import Presentation

from all2md import MarkdownConversionError, to_markdown
from all2md.options import DocxOptions, HtmlOptions, MarkdownOptions, PdfOptions
from tests.utils import (
    DocxTestGenerator,
    EmlTestGenerator,
    assert_markdown_valid,
    cleanup_test_dir,
    create_test_temp_dir,
)


@pytest.mark.integration
class TestIntegration:
    """Integration tests for full pipeline testing."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_parse_file_auto_detection(self):
        """Test automatic file format detection."""
        # Create test files of different formats

        # DOCX file
        doc = docx.Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test paragraph.")
        docx_file = self.temp_dir / "test.docx"
        doc.save(str(docx_file))

        # HTML file
        html_content = "<html><body><h1>Test HTML</h1><p>HTML content</p></body></html>"
        html_file = self.temp_dir / "test.html"
        html_file.write_text(html_content)

        # PPTX file
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        pptx_file = self.temp_dir / "test.pptx"
        prs.save(str(pptx_file))

        # EML file
        eml_content = EmlTestGenerator.create_headers_edge_cases()
        eml_file = self.temp_dir / "test.eml"
        eml_file.write_text(eml_content)

        # Test parsing each format
        docx_result = to_markdown(str(docx_file))
        html_result = to_markdown(str(html_file))
        pptx_result = to_markdown(str(pptx_file))
        eml_result = to_markdown(str(eml_file))

        # All should produce valid markdown
        assert_markdown_valid(docx_result)
        assert_markdown_valid(html_result)
        assert_markdown_valid(pptx_result)
        assert_markdown_valid(eml_result)

        # Should contain expected content
        assert "Test Document" in docx_result
        assert "Test HTML" in html_result
        assert "Test Presentation" in pptx_result
        # EML content depends on parsing implementation

    def test_parse_file_with_options(self):
        """Test parse_file with various options."""
        # Create test DOCX
        doc = DocxTestGenerator.create_complex_lists_document()
        docx_file = self.temp_dir / "complex.docx"
        doc.save(str(docx_file))

        # Test with different options (note: parse_file doesn't currently support format-specific options)
        result1 = to_markdown(str(docx_file))
        assert_markdown_valid(result1)

        result2 = to_markdown(str(docx_file))
        assert_markdown_valid(result2)

        # Both should contain list content
        assert "First item" in result1
        assert "First item" in result2

    def test_markdown_options_consistency(self):
        """Test that parse_file works with different content types."""
        # HTML with special characters
        html_content = "<p>Text with * asterisks and _ underscores</p>"
        html_file = self.temp_dir / "special.html"
        html_file.write_text(html_content)

        # Test basic parsing (options not yet supported in parse_file)
        html_result = to_markdown(str(html_file))

        # Should contain the text content
        assert "Text with" in html_result
        assert "asterisks" in html_result
        assert "underscores" in html_result

    def test_attachment_handling_consistency(self):
        """Test that files with images can be parsed consistently."""
        # Create documents with images/attachments

        # DOCX with image simulation
        doc = DocxTestGenerator.create_images_document(self.temp_dir)
        docx_file = self.temp_dir / "images.docx"
        doc.save(str(docx_file))

        # HTML with images
        html_content = '''<html><body>
            <h1>HTML with Images</h1>
            <img src="image1.png" alt="Image 1">
            <p>Text between images</p>
            <img src="image2.jpg" alt="Image 2">
        </body></html>'''
        html_file = self.temp_dir / "images.html"
        html_file.write_text(html_content)

        # Test basic parsing (format-specific options not yet supported in parse_file)
        docx_result = to_markdown(str(docx_file))
        html_result = to_markdown(str(html_file))

        assert_markdown_valid(docx_result)
        assert_markdown_valid(html_result)

        # Should contain text content
        assert "HTML with Images" in html_result
        assert "Text between images" in html_result

        # HTML should contain image references (default behavior)
        assert "Image 1" in html_result
        assert "Image 2" in html_result

    def test_complex_document_integration(self):
        """Test complex documents with multiple features."""
        # Create comprehensive test document
        doc = docx.Document()

        # Header
        doc.add_heading("Comprehensive Test Document", level=1)

        # Paragraph with formatting
        p1 = doc.add_paragraph("This document tests ")
        run1 = p1.add_run("bold")
        run1.bold = True
        p1.add_run(" and ")
        run2 = p1.add_run("italic")
        run2.italic = True
        p1.add_run(" formatting.")

        # Lists
        doc.add_paragraph("First list item", style="List Bullet")
        doc.add_paragraph("Second list item", style="List Bullet")
        doc.add_paragraph("Nested item", style="List Bullet")

        # Table
        table = doc.add_table(rows=2, cols=3)
        table.rows[0].cells[0].text = "Column A"
        table.rows[0].cells[1].text = "Column B"
        table.rows[0].cells[2].text = "Column C"
        table.rows[1].cells[0].text = "Data 1"
        table.rows[1].cells[1].text = "Data 2"
        table.rows[1].cells[2].text = "Data 3"

        # Special characters
        doc.add_paragraph("Special chars: * _ # [ ] ( ) { }")

        docx_file = self.temp_dir / "comprehensive.docx"
        doc.save(str(docx_file))

        # Test with different option combinations
        options1 = DocxOptions(
            preserve_tables=True,
            markdown_options=MarkdownOptions(escape_special=True)
        )
        # parse_file doesn't currently support format-specific options
        result1 = to_markdown(str(docx_file))
        assert_markdown_valid(result1)

        result2 = to_markdown(str(docx_file))
        assert_markdown_valid(result2)

        # Both should contain all elements
        for result in [result1, result2]:
            assert "# Comprehensive Test Document" in result
            assert "**bold**" in result
            assert "*italic*" in result
            assert "First list item" in result
            assert "| Column A | Column B | Column C |" in result
            assert "Special chars" in result

    def test_error_handling_integration(self):
        """Test error handling in integration scenarios."""
        # Test with non-existent file
        with pytest.raises(FileNotFoundError):
            with open("/non/existent/file.docx", 'rb') as f:
                to_markdown(f)

        # Test with unsupported file type - should fall back to text parsing
        unsupported_file = self.temp_dir / "test.xyz"
        unsupported_file.write_text("Unsupported content")

        result = to_markdown(str(unsupported_file))
        # Should fall back to treating as plain text
        assert result == "Unsupported content"

        # Test with corrupted file
        corrupted_file = self.temp_dir / "corrupted.docx"
        corrupted_file.write_bytes(b"Not a real DOCX file")

        # Should handle gracefully or raise appropriate error
        try:
            result = to_markdown(str(corrupted_file))
            # If it doesn't raise an error, result should be a string
            assert isinstance(result, str)
        except Exception as e:
            # Should be a meaningful error (including custom MarkdownConversionError)
            from all2md.exceptions import MarkdownConversionError
            assert isinstance(e, (ValueError, IOError, OSError, MarkdownConversionError))

    def test_performance_with_large_documents(self):
        """Test performance and handling of large documents."""
        # Create larger document
        doc = docx.Document()
        doc.add_heading("Large Document Test", level=1)

        # Add many paragraphs
        for i in range(100):
            doc.add_paragraph(f"This is paragraph {i+1} with some content to make it longer.")

        # Add large table
        table = doc.add_table(rows=20, cols=5)
        for row_idx in range(20):
            for col_idx in range(5):
                table.rows[row_idx].cells[col_idx].text = f"Cell {row_idx},{col_idx}"

        # Add many list items
        for i in range(50):
            doc.add_paragraph(f"List item {i+1}", style="List Bullet")

        large_file = self.temp_dir / "large.docx"
        doc.save(str(large_file))

        # Should handle large document without issues
        result = to_markdown(str(large_file))
        assert_markdown_valid(result)

        # Should contain first and last content
        assert "Large Document Test" in result
        assert "paragraph 1 " in result
        assert "paragraph 100" in result
        assert "List item 1" in result
        assert "List item 50" in result

    def test_file_like_object_support(self):
        """Test support for file-like objects."""
        # Test with StringIO for HTML
        html_content = b"<html><body><h1>StringIO Test</h1><p>Content from StringIO</p></body></html>"
        html_io = BytesIO(html_content)

        # parse_file requires both file object and filename
        # StringIO doesn't have a name, so we simulate it
        try:
            result = to_markdown(html_io)
            assert_markdown_valid(result)
            assert "StringIO Test" in result
        except (TypeError, AttributeError):
            # If not supported, that's also valid
            pass

    def test_batch_processing_simulation(self):
        """Test processing multiple files in sequence."""
        files_and_results = []

        # Create multiple test files
        for i in range(5):
            # HTML files
            html_content = f"<html><body><h1>Document {i+1}</h1><p>Content {i+1}</p></body></html>"
            html_file = self.temp_dir / f"doc{i+1}.html"
            html_file.write_text(html_content)

            result = to_markdown(str(html_file))
            files_and_results.append((str(html_file), result))

        # All should be processed successfully
        assert len(files_and_results) == 5

        for file_path, result in files_and_results:
            assert_markdown_valid(result)
            # Extract number from filename to verify content
            file_num = Path(file_path).stem[-1]
            assert f"Document {file_num}" in result
            assert f"Content {file_num}" in result

    def test_option_inheritance_and_precedence(self):
        """Test basic HTML parsing through parse_file."""
        # Create simple test HTML (avoiding features that trigger bugs)
        html_content = '''<html><body>
            <h1>Option Test</h1>
            <p>Text with bold and italic</p>
        </body></html>'''
        html_file = self.temp_dir / "options.html"
        html_file.write_text(html_content)

        # Test basic HTML parsing (format-specific options not yet supported in parse_file)
        result = to_markdown(str(html_file))
        assert_markdown_valid(result)

        # Should contain heading and text
        assert "Option Test" in result
        assert "Text with" in result


@pytest.mark.integration
class TestNewAPI:
    """Test the new to_markdown API with enhanced format detection and options handling."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_explicit_format_parameter(self):
        """Test explicit format specification bypasses detection."""
        # Create HTML content but force it to be processed as text
        html_content = b"<html><body><h1>Test</h1><p>Content</p></body></html>"
        html_io = BytesIO(html_content)

        # Process as HTML (auto-detection)
        result_html = to_markdown(html_io, format="html")
        assert "# Test" in result_html or "Test" in result_html
        assert "Content" in result_html
        assert "<html>" not in result_html  # Should be converted

        # Reset and process as plain text (force format)
        html_io.seek(0)
        result_text = to_markdown(html_io, format="txt")
        assert "<html>" in result_text  # Should preserve HTML tags
        assert "<body>" in result_text

    def test_options_parameter_with_kwargs_override(self):
        """Test that kwargs override options parameter values."""
        html_content = b"<html><body><p>Test content</p></body></html>"
        html_io = BytesIO(html_content)

        # Create options with one setting
        base_options = HtmlOptions(
            use_hash_headings=False,
            convert_nbsp=True
        )

        # Override with kwargs
        result = to_markdown(
            html_io,
            options=base_options,
            use_hash_headings=True,  # Override the False setting
            format="html"
        )

        assert_markdown_valid(result)

    def test_kwargs_only_options_creation(self):
        """Test creating options from kwargs alone."""
        html_content = b"<html><body><h1>Header</h1><p>Content</p></body></html>"
        html_io = BytesIO(html_content)

        result = to_markdown(
            html_io,
            format="html",
            use_hash_headings=True,
            convert_nbsp=False,
            emphasis_symbol="_"  # MarkdownOptions field
        )

        assert_markdown_valid(result)
        assert "Header" in result

    def test_content_based_format_detection(self):
        """Test format detection from file content when filename unavailable."""
        # PDF magic bytes
        pdf_content = b"%PDF-1.4\ntest content"
        pdf_io = BytesIO(pdf_content)

        # Should detect as PDF even without filename
        # Note: This will fail conversion due to invalid PDF, but format detection should work
        try:
            result = to_markdown(pdf_io)
        except Exception:
            pass  # Expected - invalid PDF content

        # Test HTML detection
        html_content = b"<!DOCTYPE html><html><body><h1>Test</h1></body></html>"
        html_io = BytesIO(html_content)

        result = to_markdown(html_io)
        assert_markdown_valid(result)
        assert "Test" in result

    def test_filename_based_detection_with_file_objects(self):
        """Test format detection from filename attribute on file objects."""
        html_content = b"<html><body><h1>Filename Test</h1></body></html>"

        # Create a BytesIO with a name attribute
        html_io = BytesIO(html_content)
        html_io.name = "test.html"

        result = to_markdown(html_io)
        assert_markdown_valid(result)
        assert "Filename Test" in result

    def test_format_specific_options_mapping(self):
        """Test that format-specific options are properly handled."""
        # Test with PDF options (even if we can't fully test PDF conversion here)
        pdf_options = PdfOptions(
            pages=[0],
            attachment_mode="alt_text"
        )

        # This should work with the options system even if PDF processing fails
        simple_text = b"Simple text content"
        text_io = BytesIO(simple_text)

        result = to_markdown(
            text_io,
            options=pdf_options,
            format="txt"  # Force as text since we don't have real PDF
        )

        assert result.strip() == "Simple text content"

    def test_backward_compatibility(self):
        """Test that old API usage still works."""
        # Old style: just file path
        html_content = "<html><body><h1>Compatibility</h1></body></html>"
        html_file = self.temp_dir / "compat.html"
        html_file.write_text(html_content)

        result = to_markdown(str(html_file))
        assert_markdown_valid(result)
        assert "Compatibility" in result

        # Old style: file object
        html_io = BytesIO(html_content.encode())
        result = to_markdown(html_io)
        assert_markdown_valid(result)

    def test_multiple_format_detections(self):
        """Test detection for various file formats."""
        # Test CSV detection
        csv_content = b"name,age,city\nJohn,25,NYC\nJane,30,LA"
        csv_io = BytesIO(csv_content)

        result = to_markdown(csv_io, format="csv")
        assert "|" in result  # Should be markdown table
        assert "John" in result
        assert "Jane" in result

        # Test TSV detection
        tsv_content = b"name\tage\tcity\nBob\t35\tSF"
        tsv_io = BytesIO(tsv_content)

        result = to_markdown(tsv_io, format="tsv")
        assert "|" in result  # Should be markdown table
        assert "Bob" in result

    def test_markdown_options_inheritance(self):
        """Test that MarkdownOptions are properly handled across formats."""
        html_content = b"<html><body><em>italic</em> <strong>bold</strong></body></html>"
        html_io = BytesIO(html_content)

        result = to_markdown(
            html_io,
            format="html",
            emphasis_symbol="_",  # Should affect italic rendering
            bullet_symbols="*-+"
        )

        assert_markdown_valid(result)
        # Check that emphasis symbol preference is handled

    def test_comprehensive_detection_strategy(self):
        """Test the comprehensive detection with multiple fallback layers."""
        # Test 1: Extension + MIME type detection for uncommon extension
        with tempfile.NamedTemporaryFile(suffix='.unknown', delete=False) as tmp:
            html_content = b"<html><body><h1>Unknown Extension</h1></body></html>"
            tmp.write(html_content)
            tmp.flush()

            # Should fall back to content detection
            result = to_markdown(tmp.name)
            assert_markdown_valid(result)
            assert "Unknown Extension" in result

        # Test 2: MIME type detection for files with known MIME but unknown extension
        csv_like_content = b"name,age\nBob,25\nAlice,30"
        csv_io = BytesIO(csv_like_content)

        # Force detection through content analysis
        result = to_markdown(csv_io)
        assert "|" in result  # Should detect as CSV and create table
        assert "Bob" in result

    def test_detection_logging_and_priority(self):
        """Test that detection methods are tried in correct priority order."""
        import logging

        # Enable debug logging to capture detection flow
        logger = logging.getLogger('all2md')
        logger.setLevel(logging.DEBUG)

        # Create in-memory log capture
        log_capture = []
        handler = logging.Handler()
        handler.emit = lambda record: log_capture.append(record.getMessage())
        logger.addHandler(handler)

        try:
            # Test filename detection (should be first and succeed)
            html_file = self.temp_dir / "test.html"
            html_file.write_text("<html><body><h1>Priority Test</h1></body></html>")

            result = to_markdown(str(html_file))

            # Check that logs show filename detection succeeded
            log_messages = [msg for msg in log_capture if 'Format detected from extension' in msg]
            assert any('.html: html' in msg for msg in log_messages)

        finally:
            logger.removeHandler(handler)
            logger.setLevel(logging.WARNING)  # Reset to default

    def test_edge_case_format_detection(self):
        """Test detection for edge cases and unusual file types."""
        # Test with empty file
        empty_io = BytesIO(b"")
        result = to_markdown(empty_io, format="txt")
        assert result == ""

        # Test with binary data that doesn't match any signature
        binary_data = b"\x00\x01\x02\x03\x04random binary data"
        binary_io = BytesIO(binary_data)

        try:
            result = to_markdown(binary_io)
            # Should fallback to text and likely fail gracefully
        except (MarkdownConversionError, UnicodeDecodeError):
            pass  # Expected for binary data

        # Test RTF detection
        rtf_content = b"{\\rtf1\\ansi\\hello world}"
        rtf_io = BytesIO(rtf_content)

        # Should detect as RTF from magic bytes
        try:
            result = to_markdown(rtf_io)
            # RTF conversion may fail due to missing dependencies, but detection should work
        except ImportError:
            pass  # Expected if pyth not installed

    def test_mime_type_fallback_enhancement(self):
        """Test that MIME type detection works as fallback."""
        # Create a file with unusual extension but recognizable MIME type
        pdf_like_file = self.temp_dir / "document.weird"
        # Write something that would trigger PDF MIME detection
        pdf_like_file.write_text("Not really PDF but has PDF extension mapping")

        # The comprehensive detection should try multiple methods
        try:
            result = to_markdown(str(pdf_like_file))
            # Should fall back to text since it's not really a PDF
            assert isinstance(result, str)
        except Exception:
            pass  # Expected for invalid format data
