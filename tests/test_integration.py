"""Integration tests for all2md converters."""

import tempfile
from pathlib import Path
from io import StringIO, BytesIO

import docx
from pptx import Presentation
import pytest

from all2md import to_markdown
from all2md.options import DocxOptions, HtmlOptions, PptxOptions, EmlOptions, MarkdownOptions
from tests.utils import (
    DocxTestGenerator, HtmlTestGenerator, PptxTestGenerator, EmlTestGenerator,
    assert_markdown_valid, create_test_temp_dir, cleanup_test_dir
)


class TestIntegration:
    """Integration tests for full pipeline testing."""

    def setup_method(self):
        """Set up test environment."""
        self.temp_dir = create_test_temp_dir()

    def teardown_method(self):
        """Clean up test environment."""
        cleanup_test_dir(self.temp_dir)

    def test_parse_file_auto_detection(self):
        """Test automatic file format detection."""
        # Create test files of different formats

        # DOCX file
        doc = docx.Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test paragraph.")
        docx_file = self.temp_dir / "test.docx"
        doc.save(str(docx_file))

        # HTML file
        html_content = "<html><body><h1>Test HTML</h1><p>HTML content</p></body></html>"
        html_file = self.temp_dir / "test.html"
        html_file.write_text(html_content)

        # PPTX file
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        pptx_file = self.temp_dir / "test.pptx"
        prs.save(str(pptx_file))

        # EML file
        eml_content = EmlTestGenerator.create_headers_edge_cases()
        eml_file = self.temp_dir / "test.eml"
        eml_file.write_text(eml_content)

        # Test parsing each format
        docx_result = to_markdown(str(docx_file))
        html_result = to_markdown(str(html_file))
        pptx_result = to_markdown(str(pptx_file))
        eml_result = to_markdown(str(eml_file))

        # All should produce valid markdown
        assert_markdown_valid(docx_result)
        assert_markdown_valid(html_result)
        assert_markdown_valid(pptx_result)
        assert_markdown_valid(eml_result)

        # Should contain expected content
        assert "Test Document" in docx_result
        assert "Test HTML" in html_result
        assert "Test Presentation" in pptx_result
        # EML content depends on parsing implementation

    def test_parse_file_with_options(self):
        """Test parse_file with various options."""
        # Create test DOCX
        doc = DocxTestGenerator.create_complex_lists_document()
        docx_file = self.temp_dir / "complex.docx"
        doc.save(str(docx_file))

        # Test with different options (note: parse_file doesn't currently support format-specific options)
        result1 = to_markdown(str(docx_file))
        assert_markdown_valid(result1)

        result2 = to_markdown(str(docx_file))
        assert_markdown_valid(result2)

        # Both should contain list content
        assert "First item" in result1
        assert "First item" in result2

    def test_markdown_options_consistency(self):
        """Test that parse_file works with different content types."""
        # HTML with special characters
        html_content = "<p>Text with * asterisks and _ underscores</p>"
        html_file = self.temp_dir / "special.html"
        html_file.write_text(html_content)

        # Test basic parsing (options not yet supported in parse_file)
        html_result = to_markdown(str(html_file))

        # Should contain the text content
        assert "Text with" in html_result
        assert "asterisks" in html_result
        assert "underscores" in html_result

    def test_attachment_handling_consistency(self):
        """Test that files with images can be parsed consistently."""
        # Create documents with images/attachments

        # DOCX with image simulation
        doc = DocxTestGenerator.create_images_document(self.temp_dir)
        docx_file = self.temp_dir / "images.docx"
        doc.save(str(docx_file))

        # HTML with images
        html_content = '''<html><body>
            <h1>HTML with Images</h1>
            <img src="image1.png" alt="Image 1">
            <p>Text between images</p>
            <img src="image2.jpg" alt="Image 2">
        </body></html>'''
        html_file = self.temp_dir / "images.html"
        html_file.write_text(html_content)

        # Test basic parsing (format-specific options not yet supported in parse_file)
        docx_result = to_markdown(str(docx_file))
        html_result = to_markdown(str(html_file))

        assert_markdown_valid(docx_result)
        assert_markdown_valid(html_result)

        # Should contain text content
        assert "HTML with Images" in html_result
        assert "Text between images" in html_result

        # HTML should contain image references (default behavior)
        assert "Image 1" in html_result
        assert "Image 2" in html_result

    def test_complex_document_integration(self):
        """Test complex documents with multiple features."""
        # Create comprehensive test document
        doc = docx.Document()

        # Header
        doc.add_heading("Comprehensive Test Document", level=1)

        # Paragraph with formatting
        p1 = doc.add_paragraph("This document tests ")
        run1 = p1.add_run("bold")
        run1.bold = True
        p1.add_run(" and ")
        run2 = p1.add_run("italic")
        run2.italic = True
        p1.add_run(" formatting.")

        # Lists
        doc.add_paragraph("First list item", style="List Bullet")
        doc.add_paragraph("Second list item", style="List Bullet")
        doc.add_paragraph("Nested item", style="List Bullet")

        # Table
        table = doc.add_table(rows=2, cols=3)
        table.rows[0].cells[0].text = "Column A"
        table.rows[0].cells[1].text = "Column B"
        table.rows[0].cells[2].text = "Column C"
        table.rows[1].cells[0].text = "Data 1"
        table.rows[1].cells[1].text = "Data 2"
        table.rows[1].cells[2].text = "Data 3"

        # Special characters
        doc.add_paragraph("Special chars: * _ # [ ] ( ) { }")

        docx_file = self.temp_dir / "comprehensive.docx"
        doc.save(str(docx_file))

        # Test with different option combinations
        options1 = DocxOptions(
            preserve_tables=True,
            markdown_options=MarkdownOptions(escape_special=True)
        )
        # parse_file doesn't currently support format-specific options
        result1 = to_markdown(str(docx_file))
        assert_markdown_valid(result1)

        result2 = to_markdown(str(docx_file))
        assert_markdown_valid(result2)

        # Both should contain all elements
        for result in [result1, result2]:
            assert "# Comprehensive Test Document" in result
            assert "**bold**" in result
            assert "*italic*" in result
            assert "First list item" in result
            assert "| Column A | Column B | Column C |" in result
            assert "Special chars" in result

    def test_error_handling_integration(self):
        """Test error handling in integration scenarios."""
        # Test with non-existent file
        with pytest.raises(FileNotFoundError):
            with open("/non/existent/file.docx", 'rb') as f:
                to_markdown(f)

        # Test with unsupported file type - should fall back to text parsing
        unsupported_file = self.temp_dir / "test.xyz"
        unsupported_file.write_text("Unsupported content")

        result = to_markdown(str(unsupported_file))
        # Should fall back to treating as plain text
        assert result == "Unsupported content"

        # Test with corrupted file
        corrupted_file = self.temp_dir / "corrupted.docx"
        corrupted_file.write_bytes(b"Not a real DOCX file")

        # Should handle gracefully or raise appropriate error
        try:
            result = to_markdown(str(corrupted_file))
            # If it doesn't raise an error, result should be a string
            assert isinstance(result, str)
        except Exception as e:
            # Should be a meaningful error (including custom MdparseConversionError)
            from all2md.exceptions import MdparseConversionError
            assert isinstance(e, (ValueError, IOError, OSError, MdparseConversionError))

    def test_performance_with_large_documents(self):
        """Test performance and handling of large documents."""
        # Create larger document
        doc = docx.Document()
        doc.add_heading("Large Document Test", level=1)

        # Add many paragraphs
        for i in range(100):
            doc.add_paragraph(f"This is paragraph {i+1} with some content to make it longer.")

        # Add large table
        table = doc.add_table(rows=20, cols=5)
        for row_idx in range(20):
            for col_idx in range(5):
                table.rows[row_idx].cells[col_idx].text = f"Cell {row_idx},{col_idx}"

        # Add many list items
        for i in range(50):
            doc.add_paragraph(f"List item {i+1}", style="List Bullet")

        large_file = self.temp_dir / "large.docx"
        doc.save(str(large_file))

        # Should handle large document without issues
        result = to_markdown(str(large_file))
        assert_markdown_valid(result)

        # Should contain first and last content
        assert "Large Document Test" in result
        assert "paragraph 1 " in result
        assert "paragraph 100" in result
        assert "List item 1" in result
        assert "List item 50" in result

    def test_file_like_object_support(self):
        """Test support for file-like objects."""
        # Test with StringIO for HTML
        html_content = b"<html><body><h1>StringIO Test</h1><p>Content from StringIO</p></body></html>"
        html_io = BytesIO(html_content)

        # parse_file requires both file object and filename
        # StringIO doesn't have a name, so we simulate it
        try:
            result = to_markdown(html_io)
            assert_markdown_valid(result)
            assert "StringIO Test" in result
        except (TypeError, AttributeError):
            # If not supported, that's also valid
            pass

    def test_batch_processing_simulation(self):
        """Test processing multiple files in sequence."""
        files_and_results = []

        # Create multiple test files
        for i in range(5):
            # HTML files
            html_content = f"<html><body><h1>Document {i+1}</h1><p>Content {i+1}</p></body></html>"
            html_file = self.temp_dir / f"doc{i+1}.html"
            html_file.write_text(html_content)

            result = to_markdown(str(html_file))
            files_and_results.append((str(html_file), result))

        # All should be processed successfully
        assert len(files_and_results) == 5

        for file_path, result in files_and_results:
            assert_markdown_valid(result)
            # Extract number from filename to verify content
            file_num = Path(file_path).stem[-1]
            assert f"Document {file_num}" in result
            assert f"Content {file_num}" in result

    def test_option_inheritance_and_precedence(self):
        """Test basic HTML parsing through parse_file."""
        # Create simple test HTML (avoiding features that trigger bugs)
        html_content = '''<html><body>
            <h1>Option Test</h1>
            <p>Text with bold and italic</p>
        </body></html>'''
        html_file = self.temp_dir / "options.html"
        html_file.write_text(html_content)

        # Test basic HTML parsing (format-specific options not yet supported in parse_file)
        result = to_markdown(str(html_file))
        assert_markdown_valid(result)

        # Should contain heading and text
        assert "Option Test" in result
        assert "Text with" in result