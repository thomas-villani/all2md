#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/unit/test_pptx_ast.py
"""Unit tests for PPTX to AST converter.

Tests cover:
- PPTX slide to AST document conversion
- Slide title and content extraction
- Text formatting (bold, italic, underline, etc.)
- Table structure conversion
- Image and chart detection
- Text frame processing
- Shape handling
- Run formatting preservation
"""

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from all2md.ast import (
    CodeBlock,
    Document,
    Emphasis,
    Heading,
    List,
    Paragraph,
    Strong,
    Table,
    Text,
    Underline,
)
from all2md.ast.transforms import extract_nodes
from all2md.options import PptxOptions
from all2md.parsers.pptx import PptxToAstConverter


@pytest.mark.unit
class TestBasicSlides:
    """Tests for basic slide conversion."""

    def test_empty_presentation(self) -> None:
        """Test converting empty presentation."""
        prs = Presentation()

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        assert isinstance(ast_doc, Document)
        assert len(ast_doc.children) == 0

    def test_single_slide_with_title(self) -> None:
        """Test converting slide with title."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Slide Title"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have heading for title (level 2 by default)
        assert len(ast_doc.children) >= 1
        assert isinstance(ast_doc.children[0], Heading)
        assert ast_doc.children[0].level == 2
        assert isinstance(ast_doc.children[0].content[0], Text)
        assert ast_doc.children[0].content[0].content == "Slide Title"

    def test_multiple_slides(self) -> None:
        """Test converting multiple slides."""
        prs = Presentation()
        layout = prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Slide 1"

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Slide 2"

        slide3 = prs.slides.add_slide(layout)
        slide3.shapes.title.text = "Slide 3"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have headings for all 3 slides
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        assert len(headings) >= 3

    def test_slide_with_title_and_content(self) -> None:
        """Test slide with title and content placeholder."""
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        title.text = "Test Title"

        # Add content to content placeholder
        # Note: Short content is detected as list by heuristics
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.text = "This is content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have heading and list (short text detected as list item)
        assert len(ast_doc.children) >= 2
        assert isinstance(ast_doc.children[0], Heading)
        assert any(isinstance(child, List) for child in ast_doc.children)

    def test_slide_without_title(self) -> None:
        """Test slide with no title (blank slide)."""
        prs = Presentation()
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        # Add text box (short text detected as list item by heuristics)
        left = top = Inches(1)
        width = height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = "Text content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have list (short text detected as list item, no heading since no title)
        assert len(ast_doc.children) >= 1
        assert any(isinstance(child, List) for child in ast_doc.children)


@pytest.mark.unit
class TestTextFormatting:
    """Tests for text formatting in slides."""

    def test_bold_text(self) -> None:
        """Test bold text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Bold text"
        run.font.bold = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Short text detected as list item - check within list
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        list_node = list_nodes[0]
        assert len(list_node.items) >= 1
        # List item contains paragraph with strong formatting
        item_para = list_node.items[0].children[0]
        assert isinstance(item_para.content[0], Strong)
        assert item_para.content[0].content[0].content == "Bold text"

    def test_italic_text(self) -> None:
        """Test italic text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Italic text"
        run.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        assert isinstance(item_para.content[0], Emphasis)
        assert item_para.content[0].content[0].content == "Italic text"

    def test_underline_text(self) -> None:
        """Test underline text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Underlined text"
        run.font.underline = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        assert isinstance(item_para.content[0], Underline)

    def test_multiple_formatting(self) -> None:
        """Test text with multiple formatting applied."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Bold and italic"
        run.font.bold = True
        run.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should be nested: Emphasis -> Strong -> Text
        assert isinstance(item_para.content[0], Emphasis)
        assert isinstance(item_para.content[0].content[0], Strong)
        assert item_para.content[0].content[0].content[0].content == "Bold and italic"

    def test_mixed_formatting_runs(self) -> None:
        """Test paragraph with multiple runs having different formatting."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        run1 = p.add_run()
        run1.text = "Normal "

        run2 = p.add_run()
        run2.text = "bold "
        run2.font.bold = True

        run3 = p.add_run()
        run3.text = "normal "

        run4 = p.add_run()
        run4.text = "italic"
        run4.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should have 4 inline nodes
        assert len(item_para.content) == 4
        assert isinstance(item_para.content[0], Text)
        assert isinstance(item_para.content[1], Strong)
        assert isinstance(item_para.content[2], Text)
        assert isinstance(item_para.content[3], Emphasis)


@pytest.mark.unit
class TestTextFrames:
    """Tests for text frame processing."""

    def test_simple_text_frame(self) -> None:
        """Test simple text frame conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = "Simple text"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Short text detected as list
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        item_para = list_nodes[0].items[0].children[0]
        assert item_para.content[0].content == "Simple text"

    def test_multiple_paragraphs_in_text_frame(self) -> None:
        """Test text frame with multiple paragraphs."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame

        p1 = tf.paragraphs[0]
        p1.text = "First paragraph"

        p2 = tf.add_paragraph()
        p2.text = "Second paragraph"

        p3 = tf.add_paragraph()
        p3.text = "Third paragraph"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # All detected as list items
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        assert len(list_nodes[0].items) >= 3

    def test_empty_text_frame_skipped(self) -> None:
        """Test that empty text frames are skipped."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add empty text box
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = ""

        # Add non-empty text box
        textbox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        textbox2.text_frame.text = "Content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should only have 1 paragraph (empty one skipped)
        para_nodes = [child for child in ast_doc.children if isinstance(child, Paragraph)]
        assert len(para_nodes) == 1
        assert para_nodes[0].content[0].content == "Content"


@pytest.mark.unit
class TestTables:
    """Tests for table conversion."""

    def test_simple_table(self) -> None:
        """Test simple table conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add table
        rows, cols = 2, 2
        left = top = Inches(1)
        width = Inches(4)
        height = Inches(2)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set cell values
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "R1C1"
        table.cell(1, 1).text = "R1C2"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have one Table node
        table_nodes = [child for child in ast_doc.children if isinstance(child, Table)]
        assert len(table_nodes) == 1
        table_node = table_nodes[0]

        # Check header
        assert table_node.header is not None
        assert len(table_node.header.cells) == 2

        # Check data rows
        assert len(table_node.rows) == 1
        assert len(table_node.rows[0].cells) == 2

    def test_table_with_formatted_cells(self) -> None:
        """Test table with formatted cell content."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table

        # Add bold text to header
        cell = table.cell(0, 0)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Bold Header"
        run.font.bold = True

        # Add normal text to data cell
        table.cell(1, 0).text = "Normal text"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        table_nodes = [child for child in ast_doc.children if isinstance(child, Table)]
        table_node = table_nodes[0]
        # Check that formatting is preserved
        header_cell = table_node.header.cells[0]
        assert isinstance(header_cell.content[0], Strong)

    def test_large_table(self) -> None:
        """Test larger table conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Create 4x5 table
        table = slide.shapes.add_table(5, 4, Inches(1), Inches(1), Inches(6), Inches(3)).table

        # Fill with data
        for row in range(5):
            for col in range(4):
                table.cell(row, col).text = f"R{row}C{col}"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        table_nodes = [child for child in ast_doc.children if isinstance(child, Table)]
        table_node = table_nodes[0]

        # Check structure
        assert table_node.header is not None
        assert len(table_node.header.cells) == 4
        assert len(table_node.rows) == 4  # 5 total - 1 header = 4 data rows


@pytest.mark.unit
class TestImages:
    """Tests for image handling."""

    def test_image_detection(self) -> None:
        """Test that images are detected and converted."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Note: Adding actual image requires file, so we just verify basic functionality
        # In real usage, the converter would detect picture shapes
        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Just verify it doesn't crash with no images
        assert isinstance(ast_doc, Document)


@pytest.mark.unit
class TestBulletLists:
    """Tests for bullet list conversion."""

    def test_bullet_list(self) -> None:
        """Test bullet list conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Item 1"

        p2 = tf.add_paragraph()
        p2.text = "Item 2"
        p2.level = 0

        p3 = tf.add_paragraph()
        p3.text = "Item 3"
        p3.level = 0

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should detect list
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        # Note: List detection depends on implementation details
        # This test may need adjustment based on how bullets are detected

    def test_nested_list(self) -> None:
        """Test nested bullet list."""
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Item 1"

        p2 = tf.add_paragraph()
        p2.text = "Item 1.1"
        p2.level = 1  # Nested

        p3 = tf.add_paragraph()
        p3.text = "Item 2"
        p3.level = 0

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have nested structure
        # Implementation-dependent test


@pytest.mark.unit
class TestComplexStructures:
    """Tests for complex slide structures."""

    def test_mixed_content_slide(self) -> None:
        """Test slide with mixed content types."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(0.5))
        title_box.text_frame.text = "Slide Title"

        # Add paragraph
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
        text_box.text_frame.text = "Some content"

        # Add table
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1.5)).table
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have multiple node types
        assert len(ast_doc.children) >= 2
        assert any(isinstance(child, Paragraph) for child in ast_doc.children)
        assert any(isinstance(child, Table) for child in ast_doc.children)

    def test_slide_with_notes(self) -> None:
        """Test that slide notes are handled appropriately."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have title at minimum
        assert len(ast_doc.children) >= 1

    def test_multiple_shapes_per_slide(self) -> None:
        """Test slide with multiple text boxes."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add 3 text boxes (short text detected as list items)
        for i in range(3):
            textbox = slide.shapes.add_textbox(
                Inches(1), Inches(1 + i * 1.5), Inches(3), Inches(1)
            )
            textbox.text_frame.text = f"Text box {i + 1}"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Short text detected as list items - should have list with 3 items
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        total_items = sum(len(list_node.items) for list_node in list_nodes)
        assert total_items >= 3


@pytest.mark.unit
class TestRunGrouping:
    """Tests for run grouping optimization."""

    def test_runs_with_same_formatting_grouped(self) -> None:
        """Test that consecutive runs with same formatting are grouped."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        p = textbox.text_frame.paragraphs[0]

        run1 = p.add_run()
        run1.text = "Bold "
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = "text"
        run2.font.bold = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should be grouped into single Strong node
        assert len(item_para.content) == 1
        assert isinstance(item_para.content[0], Strong)
        # Check text is combined
        text_content = item_para.content[0].content[0].content
        assert text_content == "Bold text"

    def test_runs_with_different_formatting_separate(self) -> None:
        """Test that runs with different formatting are kept separate."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        p = textbox.text_frame.paragraphs[0]

        run1 = p.add_run()
        run1.text = "Bold"
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = "Italic"
        run2.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should have 2 separate inline nodes
        assert len(item_para.content) == 2
        assert isinstance(item_para.content[0], Strong)
        assert isinstance(item_para.content[1], Emphasis)


@pytest.mark.unit
class TestSlideOptions:
    """Tests for slide-specific options."""

    def test_slide_numbers_disabled(self) -> None:
        """Test that slide numbers can be disabled."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        options = PptxOptions(include_slide_numbers=False)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should convert without slide number markers
        assert isinstance(ast_doc, Document)
        # Check that title doesn't have slide number prefix
        heading = ast_doc.children[0]
        assert isinstance(heading, Heading)
        assert "Slide 1:" not in heading.content[0].content

    def test_slide_numbers_enabled(self) -> None:
        """Test that slide numbers can be enabled."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        options = PptxOptions(include_slide_numbers=True)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should convert with slide number markers
        assert isinstance(ast_doc, Document)
        # Check that title includes slide number prefix
        heading = ast_doc.children[0]
        assert isinstance(heading, Heading)
        assert "Slide 1:" in heading.content[0].content


@pytest.mark.unit
class TestChartHandling:
    """Tests for chart conversion modes."""

    @staticmethod
    def _build_chart_presentation(chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED) -> Presentation:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = ChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Series 1", (1, 2))
        slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)
        return prs

    def test_mermaid_mode_produces_codeblock(self) -> None:
        prs = self._build_chart_presentation()
        converter = PptxToAstConverter(PptxOptions(charts_mode="mermaid"))
        ast_doc = converter.convert_to_ast(prs)

        code_blocks = list(extract_nodes(ast_doc, CodeBlock))
        tables = list(extract_nodes(ast_doc, Table))

        assert code_blocks, "Expected Mermaid code block when charts_mode='mermaid'"
        assert code_blocks[0].language == "mermaid"
        assert "xychart-beta" in code_blocks[0].content
        assert not tables

    def test_both_mode_produces_table_and_codeblock(self) -> None:
        prs = self._build_chart_presentation()
        converter = PptxToAstConverter(PptxOptions(charts_mode="both"))
        ast_doc = converter.convert_to_ast(prs)

        code_blocks = list(extract_nodes(ast_doc, CodeBlock))
        tables = list(extract_nodes(ast_doc, Table))

        assert code_blocks and tables
        assert code_blocks[0].language == "mermaid"

    def test_default_mode_produces_table_only(self) -> None:
        prs = self._build_chart_presentation()
        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        code_blocks = list(extract_nodes(ast_doc, CodeBlock))
        tables = list(extract_nodes(ast_doc, Table))

        assert tables
        assert not code_blocks
