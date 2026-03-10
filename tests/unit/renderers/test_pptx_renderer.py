#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/unit/test_pptx_renderer.py
"""Unit tests for PptxRenderer.

Tests cover:
- Rendering all node types to PPTX
- Slide splitting strategies
- Slide layout and formatting
- Edge cases and options

Note: These tests require python-pptx to be installed.
PPTX content verification is limited as we mainly test structure.
"""

from io import BytesIO

import pytest

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from all2md.ast import (
    Code,
    CodeBlock,
    Document,
    Emphasis,
    Heading,
    Image,
    List,
    ListItem,
    Paragraph,
    Strong,
    Table,
    TableCell,
    TableRow,
    Text,
    ThematicBreak,
)
from all2md.options import PptxRendererOptions

if PPTX_AVAILABLE:
    from all2md.renderers.pptx import PptxRenderer

pytestmark = pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")


@pytest.mark.unit
@pytest.mark.pptx
class TestBasicRendering:
    """Tests for basic PPTX rendering."""

    def test_render_empty_document(self, tmp_path):
        """Test rendering an empty document."""
        doc = Document()
        renderer = PptxRenderer()
        output_file = tmp_path / "empty.pptx"
        renderer.render(doc, output_file)

        # Verify file was created
        assert output_file.exists()
        assert output_file.stat().st_size > 0

        # Verify it's a valid PPTX
        prs = Presentation(str(output_file))
        assert prs is not None

    def test_render_to_file_path(self, tmp_path):
        """Test rendering to file path string."""
        doc = Document(children=[Paragraph(content=[Text(content="Test content")])])
        renderer = PptxRenderer()
        output_file = tmp_path / "test.pptx"
        renderer.render(doc, str(output_file))

        assert output_file.exists()

    def test_render_to_path_object(self, tmp_path):
        """Test rendering to Path object."""
        doc = Document(children=[Paragraph(content=[Text(content="Test content")])])
        renderer = PptxRenderer()
        output_file = tmp_path / "test.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_to_bytes_io(self):
        """Test rendering to BytesIO object."""
        doc = Document(children=[Paragraph(content=[Text(content="Test content")])])
        renderer = PptxRenderer()
        buffer = BytesIO()
        renderer.render(doc, buffer)

        # Verify data was written
        assert buffer.tell() > 0
        buffer.seek(0)
        data = buffer.read()
        assert len(data) > 0

    def test_render_single_slide(self, tmp_path):
        """Test rendering document to single slide."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide Title")]),
                Paragraph(content=[Text(content="Slide content")]),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "single.pptx"
        renderer.render(doc, output_file)

        # Verify slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestSplittingStrategies:
    """Tests for slide splitting strategies."""

    def test_split_by_separator(self, tmp_path):
        """Test splitting slides using separator mode."""
        doc = Document(
            children=[
                Paragraph(content=[Text(content="Slide 1 content")]),
                ThematicBreak(),
                Paragraph(content=[Text(content="Slide 2 content")]),
                ThematicBreak(),
                Paragraph(content=[Text(content="Slide 3 content")]),
            ]
        )

        options = PptxRendererOptions(slide_split_mode="separator")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_separator.pptx"
        renderer.render(doc, output_file)

        # Verify slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3

    def test_split_by_heading(self, tmp_path):
        """Test splitting slides using heading mode."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide 1")]),
                Paragraph(content=[Text(content="Content 1")]),
                Heading(level=2, content=[Text(content="Slide 2")]),
                Paragraph(content=[Text(content="Content 2")]),
                Heading(level=2, content=[Text(content="Slide 3")]),
                Paragraph(content=[Text(content="Content 3")]),
            ]
        )

        options = PptxRendererOptions(slide_split_mode="heading", slide_split_heading_level=2)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_heading.pptx"
        renderer.render(doc, output_file)

        # Verify slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3

    def test_split_auto_prefers_separator(self, tmp_path):
        """Test auto mode prefers separator when available."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="H2")]),
                Paragraph(content=[Text(content="Content 1")]),
                ThematicBreak(),
                Paragraph(content=[Text(content="Content 2")]),
            ]
        )

        options = PptxRendererOptions(slide_split_mode="auto")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_auto.pptx"
        renderer.render(doc, output_file)

        # Should split by separator, creating 2 slides
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 2

    def test_split_auto_fallback_to_heading(self, tmp_path):
        """Test auto mode falls back to headings when no separators."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide 1")]),
                Paragraph(content=[Text(content="Content 1")]),
                Heading(level=2, content=[Text(content="Slide 2")]),
                Paragraph(content=[Text(content="Content 2")]),
            ]
        )

        options = PptxRendererOptions(slide_split_mode="auto")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_auto_heading.pptx"
        renderer.render(doc, output_file)

        # Should split by H2, creating 2 slides
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 2

    def test_split_respects_heading_level(self, tmp_path):
        """Test splitting respects heading level setting."""
        doc = Document(
            children=[
                Heading(level=1, content=[Text(content="H1")]),
                Paragraph(content=[Text(content="Content 1")]),
                Heading(level=2, content=[Text(content="H2")]),
                Paragraph(content=[Text(content="Content 2")]),
                Heading(level=2, content=[Text(content="H2 Again")]),
                Paragraph(content=[Text(content="Content 3")]),
            ]
        )

        # Split on H2
        options = PptxRendererOptions(slide_split_mode="heading", slide_split_heading_level=2)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_h2.pptx"
        renderer.render(doc, output_file)

        # Should have 3 slides (H1+content, H2, H2)
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3


@pytest.mark.unit
@pytest.mark.pptx
class TestSlideTitles:
    """Tests for slide title generation."""

    def test_use_heading_as_title(self, tmp_path):
        """Test using heading text as slide title."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="My Slide Title")]),
                Paragraph(content=[Text(content="Content")]),
            ]
        )

        options = PptxRendererOptions(slide_split_mode="heading", use_heading_as_slide_title=True)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "heading_titles.pptx"
        renderer.render(doc, output_file)

        # Verify slide has title
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        assert slide.shapes.title.text == "My Slide Title"

    def test_disable_heading_titles(self, tmp_path):
        """Test disabling heading-based titles."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Heading Text")]),
                Paragraph(content=[Text(content="Content")]),
            ]
        )

        options = PptxRendererOptions(slide_split_mode="heading", use_heading_as_slide_title=False)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "no_heading_titles.pptx"
        renderer.render(doc, output_file)

        # Verify slide has no title (empty)
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        assert slide.shapes.title.text == ""


@pytest.mark.unit
@pytest.mark.pptx
class TestContentRendering:
    """Tests for rendering various content types."""

    def test_render_paragraphs(self, tmp_path):
        """Test rendering paragraphs."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Paragraph(content=[Text(content="First paragraph")]),
                Paragraph(content=[Text(content="Second paragraph")]),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "paragraphs.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_formatted_text(self, tmp_path):
        """Test rendering text with formatting."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Formatting")]),
                Paragraph(
                    content=[
                        Text(content="Normal "),
                        Strong(content=[Text(content="bold")]),
                        Text(content=" and "),
                        Emphasis(content=[Text(content="italic")]),
                        Text(content=" text."),
                    ]
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "formatted.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_code(self, tmp_path):
        """Test rendering inline code and code blocks."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Code")]),
                Paragraph(content=[Text(content="Inline "), Code(content="code"), Text(content=" here.")]),
                CodeBlock(content='def hello():\n    print("Hi")', language="python"),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "code.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_lists(self, tmp_path):
        """Test rendering bullet and numbered lists."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Lists")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Bullet 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Bullet 2")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Bullet 3")])]),
                    ],
                ),
                List(
                    ordered=True,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Number 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Number 2")])]),
                    ],
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "lists.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_tables(self, tmp_path):
        """Test rendering tables."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Table")]),
                Table(
                    header=TableRow(
                        cells=[TableCell(content=[Text(content="Name")]), TableCell(content=[Text(content="Value")])]
                    ),
                    rows=[
                        TableRow(
                            cells=[TableCell(content=[Text(content="Alpha")]), TableCell(content=[Text(content="1")])]
                        ),
                        TableRow(
                            cells=[TableCell(content=[Text(content="Beta")]), TableCell(content=[Text(content="2")])]
                        ),
                    ],
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "table.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()
        # Verify table was created
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        # Check that a table shape exists (basic verification)
        has_table = any(hasattr(shape, "table") for shape in slide.shapes)
        assert has_table


@pytest.mark.unit
@pytest.mark.pptx
class TestOptions:
    """Tests for rendering options."""

    def test_default_font_size(self, tmp_path):
        """Test default font size option."""
        doc = Document(children=[Paragraph(content=[Text(content="Content")])])

        options = PptxRendererOptions(default_font_size=24)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "font_size.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_title_font_size(self, tmp_path):
        """Test title font size option."""
        doc = Document(children=[Heading(level=2, content=[Text(content="Title")])])

        options = PptxRendererOptions(title_font_size=36)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "title_size.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_default_font(self, tmp_path):
        """Test default font option."""
        doc = Document(children=[Paragraph(content=[Text(content="Content")])])

        options = PptxRendererOptions(default_font="Arial")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "font.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()


@pytest.mark.unit
@pytest.mark.pptx
class TestComplexPresentation:
    """Tests for rendering complex presentations."""

    def test_render_multi_slide_presentation(self, tmp_path):
        """Test rendering multi-slide presentation with various content."""
        doc = Document(
            children=[
                Heading(level=1, content=[Text(content="Presentation Title")]),
                Paragraph(content=[Text(content="Introduction slide")]),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Bullet Points")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Point 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Point 2")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Point 3")])]),
                    ],
                ),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Code Example")]),
                CodeBlock(content='print("Hello")', language="python"),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Data Table")]),
                Table(
                    header=TableRow(
                        cells=[TableCell(content=[Text(content="Item")]), TableCell(content=[Text(content="Count")])]
                    ),
                    rows=[
                        TableRow(
                            cells=[TableCell(content=[Text(content="A")]), TableCell(content=[Text(content="10")])]
                        )
                    ],
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "complex.pptx"
        renderer.render(doc, output_file)

        # Verify presentation created with correct slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 4

    def test_render_nested_content(self, tmp_path):
        """Test rendering nested and complex content structures."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title with "), Strong(content=[Text(content="bold")])]),
                Paragraph(
                    content=[
                        Text(content="Mix of "),
                        Strong(content=[Text(content="bold")]),
                        Text(content=", "),
                        Emphasis(content=[Text(content="italic")]),
                        Text(content=", and "),
                        Code(content="code"),
                        Text(content="."),
                    ]
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "nested.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()


@pytest.mark.unit
@pytest.mark.pptx
class TestSpeakerNotes:
    """Tests for speaker notes rendering."""

    def test_render_with_speaker_notes(self, tmp_path):
        """Test rendering slide with speaker notes."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide Title")]),
                Paragraph(content=[Text(content="Slide content")]),
                Heading(level=3, content=[Text(content="Speaker Notes")]),
                Paragraph(content=[Text(content="These are speaker notes")]),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "with_notes.pptx"
        renderer.render(doc, output_file)

        # Verify notes were added
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "speaker notes" in notes_text.lower()

    def test_render_without_notes_when_disabled(self, tmp_path):
        """Test that notes are not rendered when include_notes=False."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide Title")]),
                Paragraph(content=[Text(content="Slide content")]),
                Heading(level=3, content=[Text(content="Speaker Notes")]),
                Paragraph(content=[Text(content="These should not appear")]),
            ]
        )

        options = PptxRendererOptions(include_notes=False)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "no_notes.pptx"
        renderer.render(doc, output_file)

        # Verify notes were not added (or are empty)
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        # Notes should be empty or only contain the "Speaker Notes" heading as regular content
        assert "These should not appear" not in notes_text

    def test_render_slide_without_notes_section(self, tmp_path):
        """Test rendering slide without speaker notes section."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide Title")]),
                Paragraph(content=[Text(content="Slide content")]),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "no_notes_section.pptx"
        renderer.render(doc, output_file)

        # Should not crash, notes should be empty
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        assert notes_text == ""

    def test_render_notes_with_formatting(self, tmp_path):
        """Test rendering speaker notes with formatting."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide Title")]),
                Paragraph(content=[Text(content="Slide content")]),
                Heading(level=3, content=[Text(content="Speaker Notes")]),
                Paragraph(
                    content=[
                        Strong(content=[Text(content="Bold")]),
                        Text(content=" and "),
                        Emphasis(content=[Text(content="italic")]),
                        Text(content=" text"),
                    ]
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "formatted_notes.pptx"
        renderer.render(doc, output_file)

        # Verify notes contain the text
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Bold" in notes_text
        assert "italic" in notes_text

    def test_render_multiple_slides_with_notes(self, tmp_path):
        """Test rendering multiple slides each with their own notes."""
        doc = Document(
            children=[
                # Slide 1
                Heading(level=2, content=[Text(content="Slide 1")]),
                Paragraph(content=[Text(content="Content 1")]),
                Heading(level=3, content=[Text(content="Speaker Notes")]),
                Paragraph(content=[Text(content="Notes for slide 1")]),
                ThematicBreak(),
                # Slide 2
                Heading(level=2, content=[Text(content="Slide 2")]),
                Paragraph(content=[Text(content="Content 2")]),
                Heading(level=3, content=[Text(content="Speaker Notes")]),
                Paragraph(content=[Text(content="Notes for slide 2")]),
                ThematicBreak(),
                # Slide 3 without notes
                Heading(level=2, content=[Text(content="Slide 3")]),
                Paragraph(content=[Text(content="Content 3")]),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "multi_notes.pptx"
        renderer.render(doc, output_file)

        # Verify each slide has correct notes
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3

        notes1 = prs.slides[0].notes_slide.notes_text_frame.text
        assert "Notes for slide 1" in notes1

        notes2 = prs.slides[1].notes_slide.notes_text_frame.text
        assert "Notes for slide 2" in notes2

        notes3 = prs.slides[2].notes_slide.notes_text_frame.text.strip()
        assert notes3 == ""  # No notes for slide 3


@pytest.mark.unit
@pytest.mark.pptx
class TestForceTextboxBullets:
    """Tests for force_textbox_bullets option."""

    def test_force_textbox_bullets_enabled_by_default(self, tmp_path):
        """Test that force_textbox_bullets is enabled by default."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide with List")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Item 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Item 2")])]),
                    ],
                ),
            ]
        )

        renderer = PptxRenderer()
        output_file = tmp_path / "bullets_default.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

        # Verify PPTX was created successfully
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_force_textbox_bullets_enabled_explicit(self, tmp_path):
        """Test force_textbox_bullets explicitly enabled."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide with List")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Item 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Item 2")])]),
                    ],
                ),
            ]
        )

        options = PptxRendererOptions(force_textbox_bullets=True)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "bullets_enabled.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

        # Verify PPTX structure
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

        # Check that the slide has text content
        slide = prs.slides[0]
        text_found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                if "Item" in shape.text_frame.text:
                    text_found = True
        assert text_found

    def test_force_textbox_bullets_disabled(self, tmp_path):
        """Test force_textbox_bullets disabled for strict templates."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide with List")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Item 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Item 2")])]),
                    ],
                ),
            ]
        )

        options = PptxRendererOptions(force_textbox_bullets=False)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "bullets_disabled.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

        # Verify PPTX structure
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

        # Content should still be rendered, just without OOXML bullet manipulation
        slide = prs.slides[0]
        text_found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                if "Item" in shape.text_frame.text:
                    text_found = True
        assert text_found

    def test_force_textbox_bullets_with_nested_lists(self, tmp_path):
        """Test force_textbox_bullets with nested lists."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Nested Lists")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(
                            children=[
                                Paragraph(content=[Text(content="Item 1")]),
                                List(
                                    ordered=False,
                                    items=[
                                        ListItem(children=[Paragraph(content=[Text(content="Subitem 1.1")])]),
                                        ListItem(children=[Paragraph(content=[Text(content="Subitem 1.2")])]),
                                    ],
                                ),
                            ]
                        ),
                        ListItem(children=[Paragraph(content=[Text(content="Item 2")])]),
                    ],
                ),
            ]
        )

        # Test with bullets enabled
        options = PptxRendererOptions(force_textbox_bullets=True)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "nested_bullets_enabled.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_force_textbox_bullets_ordered_lists_unaffected(self, tmp_path):
        """Test that force_textbox_bullets only affects unordered lists."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Ordered List")]),
                List(
                    ordered=True,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="First")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Second")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Third")])]),
                    ],
                ),
            ]
        )

        # Ordered lists use manual numbering regardless of this option
        options_enabled = PptxRendererOptions(force_textbox_bullets=True)
        options_disabled = PptxRendererOptions(force_textbox_bullets=False)

        # Test with enabled
        renderer = PptxRenderer(options_enabled)
        output_enabled = tmp_path / "ordered_enabled.pptx"
        renderer.render(doc, output_enabled)
        assert output_enabled.exists()

        # Test with disabled
        renderer = PptxRenderer(options_disabled)
        output_disabled = tmp_path / "ordered_disabled.pptx"
        renderer.render(doc, output_disabled)
        assert output_disabled.exists()

        # Both should have manual numbering
        prs_enabled = Presentation(str(output_enabled))
        prs_disabled = Presentation(str(output_disabled))
        assert len(prs_enabled.slides) == 1
        assert len(prs_disabled.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestRenderToBytes:
    """Tests for render_to_bytes method."""

    def test_render_to_bytes_simple(self):
        """Test render_to_bytes returns valid PPTX bytes."""
        doc = Document(children=[Paragraph(content=[Text(content="Byte content")])])
        renderer = PptxRenderer()
        result = renderer.render_to_bytes(doc)

        assert isinstance(result, bytes)
        assert len(result) > 0

        # Verify it's a valid PPTX by loading it
        from io import BytesIO

        prs = Presentation(BytesIO(result))
        assert len(prs.slides) >= 1


@pytest.mark.unit
@pytest.mark.pptx
class TestCommentModes:
    """Tests for comment rendering modes."""

    def test_comment_mode_ignore(self, tmp_path):
        """Test that comments are ignored when comment_mode is 'ignore'."""
        from all2md.ast import Comment

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Paragraph(content=[Text(content="Visible text")]),
                Comment(content="This should be ignored", metadata={"author": "Test"}),
            ]
        )
        options = PptxRendererOptions(comment_mode="ignore")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "comment_ignore.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1
        # Verify comment is not present in slide text
        slide = prs.slides[0]
        all_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                all_text += shape.text_frame.text
        assert "This should be ignored" not in all_text

    def test_comment_mode_visible(self, tmp_path):
        """Test that comments are visible when comment_mode is 'visible'."""
        from all2md.ast import Comment

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Comment(
                    content="Visible comment",
                    metadata={"author": "Author", "date": "2025-01-20"},
                ),
            ]
        )
        options = PptxRendererOptions(comment_mode="visible")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "comment_visible.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestInlineComments:
    """Tests for inline comment rendering in PPTX."""

    def test_inline_comment_ignore(self, tmp_path):
        """Test inline comment is ignored when comment_mode is 'ignore'."""
        from all2md.ast import CommentInline

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Paragraph(
                    content=[
                        Text(content="Before "),
                        CommentInline(content="ignored", metadata={}),
                        Text(content=" after"),
                    ]
                ),
            ]
        )
        options = PptxRendererOptions(comment_mode="ignore")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "inline_ignore.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_inline_comment_visible(self, tmp_path):
        """Test inline comment is visible when comment_mode is 'visible'."""
        from all2md.ast import CommentInline

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Paragraph(
                    content=[
                        Text(content="Text "),
                        CommentInline(
                            content="visible inline",
                            metadata={"author": "Test Author", "label": "1"},
                        ),
                    ]
                ),
            ]
        )
        options = PptxRendererOptions(comment_mode="visible")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "inline_visible.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()


@pytest.mark.unit
@pytest.mark.pptx
class TestMathRendering:
    """Tests for math rendering in PPTX."""

    def test_inline_math(self, tmp_path):
        """Test inline math rendering."""
        from all2md.ast import MathInline

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Math Example")]),
                Paragraph(
                    content=[
                        Text(content="The formula is "),
                        MathInline(content="E = mc^2", notation="latex"),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "math_inline.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1
        # Math should be rendered as text
        slide = prs.slides[0]
        all_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                all_text += shape.text_frame.text
        assert "E = mc^2" in all_text

    def test_block_math(self, tmp_path):
        """Test block math rendering."""
        from all2md.ast import MathBlock

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Math Block")]),
                MathBlock(content="\\int_0^\\infty e^{-x^2} dx", notation="latex"),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "math_block.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestDefinitionListRendering:
    """Tests for definition list rendering in PPTX."""

    def test_simple_definition_list(self, tmp_path):
        """Test simple definition list rendering."""
        from all2md.ast import DefinitionDescription, DefinitionList, DefinitionTerm

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Glossary")]),
                DefinitionList(
                    items=[
                        (
                            DefinitionTerm(content=[Text(content="API")]),
                            [DefinitionDescription(content=[Text(content="Application Programming Interface")])],
                        ),
                        (
                            DefinitionTerm(content=[Text(content="SDK")]),
                            [DefinitionDescription(content=[Text(content="Software Development Kit")])],
                        ),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "deflist.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1
        # Definition list should be rendered as text
        slide = prs.slides[0]
        all_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                all_text += shape.text_frame.text
        assert "API" in all_text
        assert "SDK" in all_text


@pytest.mark.unit
@pytest.mark.pptx
class TestInlineFormatting:
    """Additional tests for inline formatting in PPTX."""

    def test_underline(self, tmp_path):
        """Test underline rendering."""
        from all2md.ast import Underline

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Formatting")]),
                Paragraph(content=[Underline(content=[Text(content="underlined text")])]),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "underline.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_strikethrough(self, tmp_path):
        """Test strikethrough rendering."""
        from all2md.ast import Strikethrough

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Formatting")]),
                Paragraph(content=[Strikethrough(content=[Text(content="deleted text")])]),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "strikethrough.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_superscript(self, tmp_path):
        """Test superscript rendering."""
        from all2md.ast import Superscript

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Math")]),
                Paragraph(
                    content=[
                        Text(content="E = mc"),
                        Superscript(content=[Text(content="2")]),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "superscript.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_subscript(self, tmp_path):
        """Test subscript rendering."""
        from all2md.ast import Subscript

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Chemistry")]),
                Paragraph(
                    content=[
                        Text(content="H"),
                        Subscript(content=[Text(content="2")]),
                        Text(content="O"),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "subscript.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestBlockquoteRendering:
    """Tests for blockquote rendering in PPTX."""

    def test_blockquote(self, tmp_path):
        """Test blockquote rendering."""
        from all2md.ast import BlockQuote

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Quote")]),
                BlockQuote(children=[Paragraph(content=[Text(content="This is a quoted text")])]),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "blockquote.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_nested_blockquote(self, tmp_path):
        """Test nested blockquote rendering."""
        from all2md.ast import BlockQuote

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Nested Quote")]),
                BlockQuote(
                    children=[
                        Paragraph(content=[Text(content="Outer quote")]),
                        BlockQuote(children=[Paragraph(content=[Text(content="Inner quote")])]),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "nested_blockquote.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestFootnoteRendering:
    """Tests for footnote rendering in PPTX."""

    def test_footnote_reference(self, tmp_path):
        """Test footnote reference rendering."""
        from all2md.ast import FootnoteDefinition, FootnoteReference

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Content")]),
                Paragraph(
                    content=[
                        Text(content="Some text"),
                        FootnoteReference(identifier="1"),
                    ]
                ),
                FootnoteDefinition(identifier="1", content=[Text(content="Footnote content")]),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "footnote.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestThematicBreakWithContent:
    """Tests for thematic break with surrounding content."""

    def test_thematic_break_creates_slide(self, tmp_path):
        """Test that content after thematic break goes to new slide."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide 1")]),
                Paragraph(content=[Text(content="Content 1")]),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Slide 2")]),
                Paragraph(content=[Text(content="Content 2")]),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "thematic_slides.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 2


@pytest.mark.unit
@pytest.mark.pptx
class TestTableCellContent:
    """Tests for various table cell content types."""

    def test_table_cell_with_formatting(self, tmp_path):
        """Test table cell with formatted content."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Table")]),
                Table(
                    header=TableRow(
                        cells=[
                            TableCell(content=[Strong(content=[Text(content="Bold Header")])]),
                            TableCell(content=[Emphasis(content=[Text(content="Italic Header")])]),
                        ]
                    ),
                    rows=[
                        TableRow(
                            cells=[
                                TableCell(content=[Text(content="Normal")]),
                                TableCell(content=[Code(content="code")]),
                            ]
                        )
                    ],
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "table_formatted.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_table_with_empty_cells(self, tmp_path):
        """Test table with empty cells."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Table")]),
                Table(
                    rows=[
                        TableRow(
                            cells=[
                                TableCell(content=[Text(content="A")]),
                                TableCell(content=[]),
                            ]
                        ),
                        TableRow(
                            cells=[
                                TableCell(content=[]),
                                TableCell(content=[Text(content="D")]),
                            ]
                        ),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "table_empty_cells.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestLinkRendering:
    """Tests for link rendering in PPTX."""

    def test_simple_link(self, tmp_path):
        """Test simple link rendering."""
        from all2md.ast import Link

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Links")]),
                Paragraph(
                    content=[
                        Text(content="Visit "),
                        Link(url="https://example.com", content=[Text(content="Example")]),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "link.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_link_with_title(self, tmp_path):
        """Test link with title attribute."""
        from all2md.ast import Link

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Links")]),
                Paragraph(
                    content=[
                        Link(url="https://example.com", title="Example Site", content=[Text(content="Click here")]),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "link_title.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestLineBreakRendering:
    """Tests for line break rendering in PPTX."""

    def test_hard_line_break(self, tmp_path):
        """Test hard line break rendering."""
        from all2md.ast import LineBreak

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Paragraph(
                    content=[
                        Text(content="Line 1"),
                        LineBreak(soft=False),
                        Text(content="Line 2"),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "line_break.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

    def test_soft_line_break(self, tmp_path):
        """Test soft line break rendering."""
        from all2md.ast import LineBreak

        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Title")]),
                Paragraph(
                    content=[
                        Text(content="Line 1"),
                        LineBreak(soft=True),
                        Text(content="Line 2"),
                    ]
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "soft_line_break.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestContentPlaceholderFallback:
    """Tests for content placeholder fallback behavior."""

    def test_slide_without_title(self, tmp_path):
        """Test slide content without a heading/title."""
        doc = Document(
            children=[
                Paragraph(content=[Text(content="Content without heading")]),
                List(
                    ordered=False,
                    items=[
                        ListItem(children=[Paragraph(content=[Text(content="Item 1")])]),
                        ListItem(children=[Paragraph(content=[Text(content="Item 2")])]),
                    ],
                ),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "no_title.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestEmptySlideHandling:
    """Tests for handling empty slides."""

    def test_empty_slide_content_between_separators(self, tmp_path):
        """Test empty content between thematic breaks."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide 1")]),
                Paragraph(content=[Text(content="Content 1")]),
                ThematicBreak(),
                ThematicBreak(),
                Heading(level=2, content=[Text(content="Slide 3")]),
                Paragraph(content=[Text(content="Content 3")]),
            ]
        )
        renderer = PptxRenderer()
        output_file = tmp_path / "empty_between.pptx"
        renderer.render(doc, output_file)

        prs = Presentation(str(output_file))
        # Empty slide should be created
        assert len(prs.slides) >= 2


# ======================================================================
# Flow Layout Tests
# ======================================================================


def _make_simple_table() -> Table:
    """Helper: create a 2×2 table."""
    return Table(
        header=TableRow(cells=[TableCell(content=[Text(content="H1")]), TableCell(content=[Text(content="H2")])]),
        rows=[TableRow(cells=[TableCell(content=[Text(content="A")]), TableCell(content=[Text(content="B")])])],
    )


def _make_test_image(tmp_path) -> tuple[Image, str]:
    """Helper: create a small PNG on disk and return (Image node, path)."""
    import base64

    png_bytes = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/w8AAn8B9FpQHLwAAAAASUVORK5CYII="
    )
    img_path = tmp_path / "test.png"
    img_path.write_bytes(png_bytes)
    return Image(url=str(img_path), alt_text="test"), str(img_path)


def _is_title_placeholder(shape) -> bool:
    """Check if a shape is the title placeholder (idx 0)."""
    try:
        return shape.placeholder_format.idx == 0
    except (ValueError, AttributeError):
        return False


def _is_placeholder(shape) -> bool:
    """Check if a shape is any placeholder."""
    try:
        _ = shape.placeholder_format.idx
        return True
    except (ValueError, AttributeError):
        return False


def _shapes_by_top(slide):
    """Return slide shapes sorted by top position (inches)."""
    shapes = list(slide.shapes)
    shapes.sort(key=lambda s: s.top)
    return shapes


def _shape_top_inches(shape) -> float:
    from pptx.util import Emu

    return shape.top / Emu(914400)


def _shape_bottom_inches(shape) -> float:
    from pptx.util import Emu

    return (shape.top + shape.height) / Emu(914400)


@pytest.mark.unit
@pytest.mark.pptx
class TestFlowLayout:
    """Verify the flow layout engine places elements without vertical overlap."""

    def test_text_then_table_no_overlap(self, tmp_path):
        """Table top should be >= text bottom."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                Paragraph(content=[Text(content="Some text before the table.")]),
                _make_simple_table(),
            ]
        )
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        sorted_shapes = _shapes_by_top(slide)

        # Find textbox and table shapes (skip title placeholder)
        non_title = [s for s in sorted_shapes if not _is_title_placeholder(s)]
        assert len(non_title) >= 2
        text_shape = non_title[0]
        table_shape = non_title[1]
        assert _shape_top_inches(table_shape) >= _shape_bottom_inches(text_shape) - 0.01

    def test_text_then_image_no_overlap(self, tmp_path):
        """Image top should be >= text bottom."""
        img_node, _ = _make_test_image(tmp_path)
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                Paragraph(content=[Text(content="Text before image.")]),
                img_node,
            ]
        )
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        non_title = [s for s in _shapes_by_top(slide) if not _is_title_placeholder(s)]
        assert len(non_title) >= 2
        assert _shape_top_inches(non_title[1]) >= _shape_bottom_inches(non_title[0]) - 0.01

    def test_table_then_text_no_overlap(self, tmp_path):
        """Text top should be >= table bottom."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                _make_simple_table(),
                Paragraph(content=[Text(content="Text after table.")]),
            ]
        )
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        non_title = [s for s in _shapes_by_top(slide) if not _is_title_placeholder(s)]
        assert len(non_title) >= 2
        assert _shape_top_inches(non_title[1]) >= _shape_bottom_inches(non_title[0]) - 0.01

    def test_text_table_text_sequential(self, tmp_path):
        """Three blocks should be vertically sequential."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                Paragraph(content=[Text(content="Before.")]),
                _make_simple_table(),
                Paragraph(content=[Text(content="After.")]),
            ]
        )
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        non_title = [s for s in _shapes_by_top(slide) if not _is_title_placeholder(s)]
        assert len(non_title) >= 3
        for i in range(len(non_title) - 1):
            assert _shape_top_inches(non_title[i + 1]) >= _shape_bottom_inches(non_title[i]) - 0.01

    def test_multiple_tables_sequential(self, tmp_path):
        """Second table should be below the first."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                _make_simple_table(),
                _make_simple_table(),
            ]
        )
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        non_title = [s for s in _shapes_by_top(slide) if not _is_title_placeholder(s)]
        assert len(non_title) >= 2
        assert _shape_top_inches(non_title[1]) >= _shape_bottom_inches(non_title[0]) - 0.01

    def test_flow_layout_creates_separate_textboxes(self, tmp_path):
        """Text before and after a table should produce 2 separate textbox shapes."""
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                Paragraph(content=[Text(content="Before.")]),
                _make_simple_table(),
                Paragraph(content=[Text(content="After.")]),
            ]
        )
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        # Count textbox shapes (those with text_frame that are NOT placeholders and NOT tables)
        textboxes = [s for s in slide.shapes if s.has_text_frame and not s.has_table and not _is_placeholder(s)]
        assert len(textboxes) >= 2


@pytest.mark.unit
@pytest.mark.pptx
class TestWidescreenDefaults:
    """Verify widescreen (16:9) slide dimensions."""

    def test_default_slide_dimensions(self, tmp_path):
        """Default should be 13.333 x 7.5 inches."""
        doc = Document(children=[Paragraph(content=[Text(content="Hello")])])
        renderer = PptxRenderer()
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        from pptx.util import Emu

        width_in = prs.slide_width / Emu(914400)
        height_in = prs.slide_height / Emu(914400)
        assert abs(width_in - 13.333) < 0.01
        assert abs(height_in - 7.5) < 0.01

    def test_custom_slide_dimensions(self, tmp_path):
        """Custom dimensions should be applied."""
        opts = PptxRendererOptions(slide_width=10.0, slide_height=5.625)
        renderer = PptxRenderer(opts)
        doc = Document(children=[Paragraph(content=[Text(content="Hello")])])
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        from pptx.util import Emu

        width_in = prs.slide_width / Emu(914400)
        height_in = prs.slide_height / Emu(914400)
        assert abs(width_in - 10.0) < 0.01
        assert abs(height_in - 5.625) < 0.01

    def test_template_preserves_dimensions(self, tmp_path):
        """When a template is provided, its dimensions should be preserved."""
        # Create a template with 4:3 dimensions
        from pptx.util import Inches

        template_prs = Presentation()
        template_prs.slide_width = Inches(10.0)
        template_prs.slide_height = Inches(7.5)
        template_path = tmp_path / "template.pptx"
        template_prs.save(str(template_path))

        opts = PptxRendererOptions(template_path=str(template_path))
        renderer = PptxRenderer(opts)
        doc = Document(children=[Paragraph(content=[Text(content="Hello")])])
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        from pptx.util import Emu

        width_in = prs.slide_width / Emu(914400)
        # Template was 10.0", renderer should NOT override it
        assert abs(width_in - 10.0) < 0.01


@pytest.mark.unit
@pytest.mark.pptx
class TestTextHeightEstimation:
    """Test the _estimate_text_block_height helper."""

    def test_short_paragraph_height(self):
        """Short paragraph should return at least MIN_TEXT_HEIGHT."""
        from all2md.constants import DEFAULT_PPTX_MIN_TEXT_HEIGHT

        renderer = PptxRenderer()
        nodes = [Paragraph(content=[Text(content="Short")])]
        height = renderer._estimate_text_block_height(nodes, 12.0)
        assert height >= DEFAULT_PPTX_MIN_TEXT_HEIGHT

    def test_long_paragraph_wraps(self):
        """Long paragraph should produce a taller estimate than a short one."""
        renderer = PptxRenderer()
        short_nodes = [Paragraph(content=[Text(content="Short")])]
        long_text = "This is a much longer paragraph. " * 20
        long_nodes = [Paragraph(content=[Text(content=long_text)])]
        h_short = renderer._estimate_text_block_height(short_nodes, 5.0)
        h_long = renderer._estimate_text_block_height(long_nodes, 5.0)
        assert h_long > h_short

    def test_code_block_multiline(self):
        """Code block height should scale with line count."""
        renderer = PptxRenderer()
        code_2 = [CodeBlock(content="a\nb")]
        code_10 = [CodeBlock(content="\n".join(["line"] * 10))]
        h2 = renderer._estimate_text_block_height(code_2, 12.0)
        h10 = renderer._estimate_text_block_height(code_10, 12.0)
        assert h10 > h2

    def test_empty_returns_minimum(self):
        """Empty node list should return minimum height."""
        from all2md.constants import DEFAULT_PPTX_MIN_TEXT_HEIGHT

        renderer = PptxRenderer()
        height = renderer._estimate_text_block_height([], 12.0)
        # With no nodes total_lines is 0, so height = max(0, MIN) = MIN
        assert height == pytest.approx(DEFAULT_PPTX_MIN_TEXT_HEIGHT, abs=0.01)


@pytest.mark.unit
@pytest.mark.pptx
class TestFlowLayoutBackwardCompat:
    """Verify that use_flow_layout=False preserves old fixed-position behavior."""

    def test_use_flow_layout_false(self, tmp_path):
        """With flow layout disabled, tables should use the fixed table_top position."""
        opts = PptxRendererOptions(use_flow_layout=False, table_top=2.0)
        renderer = PptxRenderer(opts)
        doc = Document(
            children=[
                Heading(level=2, content=[Text(content="Slide")]),
                Paragraph(content=[Text(content="Some text.")]),
                _make_simple_table(),
            ]
        )
        out = tmp_path / "t.pptx"
        renderer.render(doc, out)

        prs = Presentation(str(out))
        slide = prs.slides[0]
        # Find the table shape
        table_shapes = [s for s in slide.shapes if s.has_table]
        assert len(table_shapes) == 1
        assert abs(_shape_top_inches(table_shapes[0]) - 2.0) < 0.01
