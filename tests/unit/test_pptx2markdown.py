import base64

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from all2md.converters.pptx2markdown import (
    _process_paragraph_format,
    _process_run_format,
    _process_shape,
    _process_table,
    _process_text_frame,
    create_test_presentation,
    pptx_to_markdown,
)
from all2md.utils.attachments import extract_pptx_image_data


@pytest.mark.unit
def test_process_paragraph_format():
    DummyPara = type("DummyPara", (), {"level": 2})
    assert _process_paragraph_format(DummyPara) == [("  " * 2, "")]


@pytest.mark.unit
def test_process_run_format():
    font = type("Font", (), {"bold": True, "italic": False, "underline": True})()
    run = type("Run", (), {"font": font})()
    # Test with empty text to see formatting structure
    result = _process_run_format(run, "test")
    assert "**test**" in result  # bold formatting
    assert "<u>" in result and "</u>" in result  # underline formatting (HTML mode default)


@pytest.mark.unit
def test_process_text_frame_formatting():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    left, top, width, height = Inches(1), Inches(1), Inches(4), Inches(1)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.add_paragraph()
    p.text = "Hello"
    run = p.runs[0]
    run.font.bold = True
    run.font.italic = True
    run.font.underline = True
    result = _process_text_frame(tf)
    # With new formatting, underline uses HTML tags by default
    expected_patterns = ["***Hello***", "<u>", "</u>"]  # bold+italic wrapped in underline HTML
    for pattern in expected_patterns:
        assert pattern in result


@pytest.mark.unit
def test_process_table():
    prs = create_test_presentation()
    slide = prs.slides[2]
    table = slide.shapes[-1].table
    md_table = _process_table(table)
    expected = "\n".join(
        [
            "| Header 1 | Header 2 | Header 3 |",
            "| --- | --- | --- |",
            "| Cell 1,0 | Cell 1,1 | Cell 1,2 |",
            "| Cell 2,0 | Cell 2,1 | Cell 2,2 |",
        ]
    )
    assert md_table == expected


@pytest.mark.unit
def test_extract_image_data():
    png_b64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAAWgmWQ0AAAAASUVORK5CYII="
    png_bytes = base64.b64decode(png_b64)
    DummyImage = type("DummyImage", (), {})()
    DummyImage.blob = png_bytes
    DummyShape = type("DummyShape", (), {})()
    DummyShape.image = DummyImage
    data = extract_pptx_image_data(DummyShape)
    assert data == png_bytes  # Should return raw bytes now
    assert isinstance(data, bytes)


@pytest.mark.unit
def test_process_shape_picture():
    png_b64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAAWgmWQ0AAAAASUVORK5CYII="
    png_bytes = base64.b64decode(png_b64)
    dummy_img = type("Img", (), {})()
    dummy_img.blob = png_bytes
    shape = type(
        "S", (), {"has_text_frame": False, "shape_type": MSO_SHAPE_TYPE.PICTURE, "alt_text": "alt", "image": dummy_img}
    )()
    from all2md.options import PptxOptions

    options_alt = PptxOptions(attachment_mode="alt_text")
    no_md = _process_shape(shape, options_alt)
    assert no_md == "![alt]"
    options_base64 = PptxOptions(attachment_mode="base64")
    yes_md = _process_shape(shape, options_base64)
    assert yes_md.startswith("![alt](data:image/") and yes_md.endswith(")")


@pytest.mark.unit
def test_process_shape_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    chart_data = ChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Series 1", (1.0, 2.0, 3.0))
    x, y, cx, cy = Inches(1), Inches(1), Inches(4), Inches(3)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    from all2md.options import PptxOptions

    options = PptxOptions(attachment_mode="alt_text")
    md_chart = _process_shape(graphic_frame, options)
    expected = "\n".join(["| Category | A | B | C |", "| --- | --- | --- | --- |", "| Series 1 | 1.0 | 2.0 | 3.0 |"])
    assert md_chart == expected


@pytest.mark.unit
def test_pptx_to_markdown_default():
    prs = create_test_presentation()
    # tmpfile = NamedTemporaryFile("wb", suffix=".pptx", delete=False)
    # prs.save(tmpfile.name)

    from all2md.options import PptxOptions

    options = PptxOptions(attachment_mode="alt_text", include_slide_numbers=True)
    output = pptx_to_markdown(prs, options=options)
    assert "# Slide 1: Test PowerPoint Presentation" in output
    assert "First Level Bullet" in output
    assert "| Header 1 | Header 2 | Header 3 |" in output
    assert output.strip().endswith("---")


@pytest.mark.unit
def test_pptx_to_markdown_no_slide_numbers():
    prs = create_test_presentation()
    from all2md.options import PptxOptions

    options = PptxOptions(attachment_mode="alt_text", include_slide_numbers=False)
    output = pptx_to_markdown(prs, options=options)
    assert "# Test PowerPoint Presentation" in output
    assert "Slide 1:" not in output
