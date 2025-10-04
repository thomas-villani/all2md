#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/unit/test_pptx_renderer.py
"""Unit tests for PptxRenderer.

Tests cover:
- Rendering all node types to PPTX
- Slide splitting strategies
- Slide layout and formatting
- Edge cases and options

Note: These tests require python-pptx to be installed.
PPTX content verification is limited as we mainly test structure.
"""

import pytest
from pathlib import Path
from io import BytesIO

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from all2md.options import PptxRendererOptions
from all2md.ast import (
    Document,
    Heading,
    Paragraph,
    Text,
    ThematicBreak,
    Strong,
    Emphasis,
    Code,
    List,
    ListItem,
    CodeBlock,
    Table,
    TableCell,
    TableRow,
)

if PPTX_AVAILABLE:
    from all2md.renderers.pptx import PptxRenderer


pytestmark = pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")


@pytest.mark.unit
@pytest.mark.pptx
class TestBasicRendering:
    """Tests for basic PPTX rendering."""

    def test_render_empty_document(self, tmp_path):
        """Test rendering an empty document."""
        doc = Document()
        renderer = PptxRenderer()
        output_file = tmp_path / "empty.pptx"
        renderer.render(doc, output_file)

        # Verify file was created
        assert output_file.exists()
        assert output_file.stat().st_size > 0

        # Verify it's a valid PPTX
        prs = Presentation(str(output_file))
        assert prs is not None

    def test_render_to_file_path(self, tmp_path):
        """Test rendering to file path string."""
        doc = Document(children=[
            Paragraph(content=[Text(content="Test content")])
        ])
        renderer = PptxRenderer()
        output_file = tmp_path / "test.pptx"
        renderer.render(doc, str(output_file))

        assert output_file.exists()

    def test_render_to_path_object(self, tmp_path):
        """Test rendering to Path object."""
        doc = Document(children=[
            Paragraph(content=[Text(content="Test content")])
        ])
        renderer = PptxRenderer()
        output_file = tmp_path / "test.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_to_bytes_io(self):
        """Test rendering to BytesIO object."""
        doc = Document(children=[
            Paragraph(content=[Text(content="Test content")])
        ])
        renderer = PptxRenderer()
        buffer = BytesIO()
        renderer.render(doc, buffer)

        # Verify data was written
        assert buffer.tell() > 0
        buffer.seek(0)
        data = buffer.read()
        assert len(data) > 0

    def test_render_single_slide(self, tmp_path):
        """Test rendering document to single slide."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Slide Title")]),
            Paragraph(content=[Text(content="Slide content")])
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "single.pptx"
        renderer.render(doc, output_file)

        # Verify slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1


@pytest.mark.unit
@pytest.mark.pptx
class TestSplittingStrategies:
    """Tests for slide splitting strategies."""

    def test_split_by_separator(self, tmp_path):
        """Test splitting slides using separator mode."""
        doc = Document(children=[
            Paragraph(content=[Text(content="Slide 1 content")]),
            ThematicBreak(),
            Paragraph(content=[Text(content="Slide 2 content")]),
            ThematicBreak(),
            Paragraph(content=[Text(content="Slide 3 content")])
        ])

        options = PptxRendererOptions(slide_split_mode="separator")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_separator.pptx"
        renderer.render(doc, output_file)

        # Verify slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3

    def test_split_by_heading(self, tmp_path):
        """Test splitting slides using heading mode."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Slide 1")]),
            Paragraph(content=[Text(content="Content 1")]),
            Heading(level=2, content=[Text(content="Slide 2")]),
            Paragraph(content=[Text(content="Content 2")]),
            Heading(level=2, content=[Text(content="Slide 3")]),
            Paragraph(content=[Text(content="Content 3")])
        ])

        options = PptxRendererOptions(
            slide_split_mode="heading",
            slide_split_heading_level=2
        )
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_heading.pptx"
        renderer.render(doc, output_file)

        # Verify slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3

    def test_split_auto_prefers_separator(self, tmp_path):
        """Test auto mode prefers separator when available."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="H2")]),
            Paragraph(content=[Text(content="Content 1")]),
            ThematicBreak(),
            Paragraph(content=[Text(content="Content 2")])
        ])

        options = PptxRendererOptions(slide_split_mode="auto")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_auto.pptx"
        renderer.render(doc, output_file)

        # Should split by separator, creating 2 slides
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 2

    def test_split_auto_fallback_to_heading(self, tmp_path):
        """Test auto mode falls back to headings when no separators."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Slide 1")]),
            Paragraph(content=[Text(content="Content 1")]),
            Heading(level=2, content=[Text(content="Slide 2")]),
            Paragraph(content=[Text(content="Content 2")])
        ])

        options = PptxRendererOptions(slide_split_mode="auto")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_auto_heading.pptx"
        renderer.render(doc, output_file)

        # Should split by H2, creating 2 slides
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 2

    def test_split_respects_heading_level(self, tmp_path):
        """Test splitting respects heading level setting."""
        doc = Document(children=[
            Heading(level=1, content=[Text(content="H1")]),
            Paragraph(content=[Text(content="Content 1")]),
            Heading(level=2, content=[Text(content="H2")]),
            Paragraph(content=[Text(content="Content 2")]),
            Heading(level=2, content=[Text(content="H2 Again")]),
            Paragraph(content=[Text(content="Content 3")])
        ])

        # Split on H2
        options = PptxRendererOptions(
            slide_split_mode="heading",
            slide_split_heading_level=2
        )
        renderer = PptxRenderer(options)
        output_file = tmp_path / "split_h2.pptx"
        renderer.render(doc, output_file)

        # Should have 3 slides (H1+content, H2, H2)
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3


@pytest.mark.unit
@pytest.mark.pptx
class TestSlideTitles:
    """Tests for slide title generation."""

    def test_use_heading_as_title(self, tmp_path):
        """Test using heading text as slide title."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="My Slide Title")]),
            Paragraph(content=[Text(content="Content")])
        ])

        options = PptxRendererOptions(
            slide_split_mode="heading",
            use_heading_as_slide_title=True
        )
        renderer = PptxRenderer(options)
        output_file = tmp_path / "heading_titles.pptx"
        renderer.render(doc, output_file)

        # Verify slide has title
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        assert slide.shapes.title.text == "My Slide Title"

    def test_disable_heading_titles(self, tmp_path):
        """Test disabling heading-based titles."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Heading Text")]),
            Paragraph(content=[Text(content="Content")])
        ])

        options = PptxRendererOptions(
            slide_split_mode="heading",
            use_heading_as_slide_title=False
        )
        renderer = PptxRenderer(options)
        output_file = tmp_path / "no_heading_titles.pptx"
        renderer.render(doc, output_file)

        # Verify slide has no title (empty)
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        assert slide.shapes.title.text == ""


@pytest.mark.unit
@pytest.mark.pptx
class TestContentRendering:
    """Tests for rendering various content types."""

    def test_render_paragraphs(self, tmp_path):
        """Test rendering paragraphs."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Title")]),
            Paragraph(content=[Text(content="First paragraph")]),
            Paragraph(content=[Text(content="Second paragraph")])
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "paragraphs.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_formatted_text(self, tmp_path):
        """Test rendering text with formatting."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Formatting")]),
            Paragraph(content=[
                Text(content="Normal "),
                Strong(content=[Text(content="bold")]),
                Text(content=" and "),
                Emphasis(content=[Text(content="italic")]),
                Text(content=" text.")
            ])
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "formatted.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_code(self, tmp_path):
        """Test rendering inline code and code blocks."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Code")]),
            Paragraph(content=[
                Text(content="Inline "),
                Code(content="code"),
                Text(content=" here.")
            ]),
            CodeBlock(content='def hello():\n    print("Hi")', language="python")
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "code.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_lists(self, tmp_path):
        """Test rendering bullet and numbered lists."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Lists")]),
            List(ordered=False, items=[
                ListItem(children=[Paragraph(content=[Text(content="Bullet 1")])]),
                ListItem(children=[Paragraph(content=[Text(content="Bullet 2")])]),
                ListItem(children=[Paragraph(content=[Text(content="Bullet 3")])])
            ]),
            List(ordered=True, items=[
                ListItem(children=[Paragraph(content=[Text(content="Number 1")])]),
                ListItem(children=[Paragraph(content=[Text(content="Number 2")])])
            ])
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "lists.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_render_tables(self, tmp_path):
        """Test rendering tables."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Table")]),
            Table(
                header=TableRow(cells=[
                    TableCell(content=[Text(content="Name")]),
                    TableCell(content=[Text(content="Value")])
                ]),
                rows=[
                    TableRow(cells=[
                        TableCell(content=[Text(content="Alpha")]),
                        TableCell(content=[Text(content="1")])
                    ]),
                    TableRow(cells=[
                        TableCell(content=[Text(content="Beta")]),
                        TableCell(content=[Text(content="2")])
                    ])
                ]
            )
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "table.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()
        # Verify table was created
        prs = Presentation(str(output_file))
        slide = prs.slides[0]
        # Check that a table shape exists (basic verification)
        has_table = any(hasattr(shape, 'table') for shape in slide.shapes)
        assert has_table


@pytest.mark.unit
@pytest.mark.pptx
class TestOptions:
    """Tests for rendering options."""

    def test_default_font_size(self, tmp_path):
        """Test default font size option."""
        doc = Document(children=[
            Paragraph(content=[Text(content="Content")])
        ])

        options = PptxRendererOptions(default_font_size=24)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "font_size.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_title_font_size(self, tmp_path):
        """Test title font size option."""
        doc = Document(children=[
            Heading(level=2, content=[Text(content="Title")])
        ])

        options = PptxRendererOptions(title_font_size=36)
        renderer = PptxRenderer(options)
        output_file = tmp_path / "title_size.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()

    def test_default_font(self, tmp_path):
        """Test default font option."""
        doc = Document(children=[
            Paragraph(content=[Text(content="Content")])
        ])

        options = PptxRendererOptions(default_font="Arial")
        renderer = PptxRenderer(options)
        output_file = tmp_path / "font.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()


@pytest.mark.unit
@pytest.mark.pptx
class TestComplexPresentation:
    """Tests for rendering complex presentations."""

    def test_render_multi_slide_presentation(self, tmp_path):
        """Test rendering multi-slide presentation with various content."""
        doc = Document(children=[
            Heading(level=1, content=[Text(content="Presentation Title")]),
            Paragraph(content=[Text(content="Introduction slide")]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Bullet Points")]),
            List(ordered=False, items=[
                ListItem(children=[Paragraph(content=[Text(content="Point 1")])]),
                ListItem(children=[Paragraph(content=[Text(content="Point 2")])]),
                ListItem(children=[Paragraph(content=[Text(content="Point 3")])])
            ]),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Code Example")]),
            CodeBlock(content='print("Hello")', language="python"),
            ThematicBreak(),
            Heading(level=2, content=[Text(content="Data Table")]),
            Table(
                header=TableRow(cells=[
                    TableCell(content=[Text(content="Item")]),
                    TableCell(content=[Text(content="Count")])
                ]),
                rows=[
                    TableRow(cells=[
                        TableCell(content=[Text(content="A")]),
                        TableCell(content=[Text(content="10")])
                    ])
                ]
            )
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "complex.pptx"
        renderer.render(doc, output_file)

        # Verify presentation created with correct slide count
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 4

    def test_render_nested_content(self, tmp_path):
        """Test rendering nested and complex content structures."""
        doc = Document(children=[
            Heading(level=2, content=[
                Text(content="Title with "),
                Strong(content=[Text(content="bold")])
            ]),
            Paragraph(content=[
                Text(content="Mix of "),
                Strong(content=[Text(content="bold")]),
                Text(content=", "),
                Emphasis(content=[Text(content="italic")]),
                Text(content=", and "),
                Code(content="code"),
                Text(content=".")
            ])
        ])

        renderer = PptxRenderer()
        output_file = tmp_path / "nested.pptx"
        renderer.render(doc, output_file)

        assert output_file.exists()
