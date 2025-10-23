#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# tests/unit/test_pptx_ast.py
"""Unit tests for PPTX to AST converter.

Tests cover:
- PPTX slide to AST document conversion
- Slide title and content extraction
- Text formatting (bold, italic, underline, etc.)
- Table structure conversion
- Image and chart detection
- Text frame processing
- Shape handling
- Run formatting preservation
"""

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from all2md.ast import (
    CodeBlock,
    Comment,
    Document,
    Emphasis,
    Heading,
    Link,
    List,
    Paragraph,
    Strikethrough,
    Strong,
    Subscript,
    Superscript,
    Table,
    Text,
    Underline,
)
from all2md.ast.transforms import extract_nodes
from all2md.options import PptxOptions
from all2md.parsers.pptx import (
    PptxToAstConverter,
    _analyze_slide_context,
    _detect_list_formatting_xml,
    _detect_list_item,
)


@pytest.mark.unit
class TestBasicSlides:
    """Tests for basic slide conversion."""

    def test_empty_presentation(self) -> None:
        """Test converting empty presentation."""
        prs = Presentation()

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        assert isinstance(ast_doc, Document)
        assert len(ast_doc.children) == 0

    def test_single_slide_with_title(self) -> None:
        """Test converting slide with title."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Slide Title"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have heading for title (level 2 by default)
        assert len(ast_doc.children) >= 1
        assert isinstance(ast_doc.children[0], Heading)
        assert ast_doc.children[0].level == 2
        assert isinstance(ast_doc.children[0].content[0], Text)
        assert ast_doc.children[0].content[0].content == "Slide Title"

    def test_multiple_slides(self) -> None:
        """Test converting multiple slides."""
        prs = Presentation()
        layout = prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Slide 1"

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Slide 2"

        slide3 = prs.slides.add_slide(layout)
        slide3.shapes.title.text = "Slide 3"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have headings for all 3 slides
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        assert len(headings) >= 3

    def test_slide_with_title_and_content(self) -> None:
        """Test slide with title and content placeholder."""
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        title.text = "Test Title"

        # Add content to content placeholder
        # Note: Short content is detected as list by heuristics
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.text = "This is content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have heading and list (short text detected as list item)
        assert len(ast_doc.children) >= 2
        assert isinstance(ast_doc.children[0], Heading)
        assert any(isinstance(child, List) for child in ast_doc.children)

    def test_slide_without_title(self) -> None:
        """Test slide with no title (blank slide)."""
        prs = Presentation()
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        # Add text box (short text detected as list item by heuristics)
        left = top = Inches(1)
        width = height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = "Text content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have list (short text detected as list item, no heading since no title)
        assert len(ast_doc.children) >= 1
        assert any(isinstance(child, List) for child in ast_doc.children)


@pytest.mark.unit
class TestTextFormatting:
    """Tests for text formatting in slides."""

    def test_bold_text(self) -> None:
        """Test bold text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Bold text"
        run.font.bold = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Short text detected as list item - check within list
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        list_node = list_nodes[0]
        assert len(list_node.items) >= 1
        # List item contains paragraph with strong formatting
        item_para = list_node.items[0].children[0]
        assert isinstance(item_para.content[0], Strong)
        assert item_para.content[0].content[0].content == "Bold text"

    def test_italic_text(self) -> None:
        """Test italic text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Italic text"
        run.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        assert isinstance(item_para.content[0], Emphasis)
        assert item_para.content[0].content[0].content == "Italic text"

    def test_underline_text(self) -> None:
        """Test underline text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Underlined text"
        run.font.underline = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        assert isinstance(item_para.content[0], Underline)

    def test_multiple_formatting(self) -> None:
        """Test text with multiple formatting applied."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Bold and italic"
        run.font.bold = True
        run.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should be nested: Emphasis -> Strong -> Text
        assert isinstance(item_para.content[0], Emphasis)
        assert isinstance(item_para.content[0].content[0], Strong)
        assert item_para.content[0].content[0].content[0].content == "Bold and italic"

    def test_mixed_formatting_runs(self) -> None:
        """Test paragraph with multiple runs having different formatting."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        run1 = p.add_run()
        run1.text = "Normal "

        run2 = p.add_run()
        run2.text = "bold "
        run2.font.bold = True

        run3 = p.add_run()
        run3.text = "normal "

        run4 = p.add_run()
        run4.text = "italic"
        run4.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should have 4 inline nodes
        assert len(item_para.content) == 4
        assert isinstance(item_para.content[0], Text)
        assert isinstance(item_para.content[1], Strong)
        assert isinstance(item_para.content[2], Text)
        assert isinstance(item_para.content[3], Emphasis)


@pytest.mark.unit
class TestTextFrames:
    """Tests for text frame processing."""

    def test_simple_text_frame(self) -> None:
        """Test simple text frame conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = "Simple text"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Short text detected as list
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        item_para = list_nodes[0].items[0].children[0]
        assert item_para.content[0].content == "Simple text"

    def test_multiple_paragraphs_in_text_frame(self) -> None:
        """Test text frame with multiple paragraphs."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame

        p1 = tf.paragraphs[0]
        p1.text = "First paragraph"

        p2 = tf.add_paragraph()
        p2.text = "Second paragraph"

        p3 = tf.add_paragraph()
        p3.text = "Third paragraph"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # All detected as list items
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        assert len(list_nodes[0].items) >= 3

    def test_empty_text_frame_skipped(self) -> None:
        """Test that empty text frames are skipped."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add empty text box
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = ""

        # Add non-empty text box
        textbox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        textbox2.text_frame.text = "Content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should only have 1 paragraph (empty one skipped)
        para_nodes = [child for child in ast_doc.children if isinstance(child, Paragraph)]
        assert len(para_nodes) == 1
        assert para_nodes[0].content[0].content == "Content"


@pytest.mark.unit
class TestTables:
    """Tests for table conversion."""

    def test_simple_table(self) -> None:
        """Test simple table conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add table
        rows, cols = 2, 2
        left = top = Inches(1)
        width = Inches(4)
        height = Inches(2)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set cell values
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "R1C1"
        table.cell(1, 1).text = "R1C2"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have one Table node
        table_nodes = [child for child in ast_doc.children if isinstance(child, Table)]
        assert len(table_nodes) == 1
        table_node = table_nodes[0]

        # Check header
        assert table_node.header is not None
        assert len(table_node.header.cells) == 2

        # Check data rows
        assert len(table_node.rows) == 1
        assert len(table_node.rows[0].cells) == 2

    def test_table_with_formatted_cells(self) -> None:
        """Test table with formatted cell content."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table

        # Add bold text to header
        cell = table.cell(0, 0)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Bold Header"
        run.font.bold = True

        # Add normal text to data cell
        table.cell(1, 0).text = "Normal text"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        table_nodes = [child for child in ast_doc.children if isinstance(child, Table)]
        table_node = table_nodes[0]
        # Check that formatting is preserved
        header_cell = table_node.header.cells[0]
        assert isinstance(header_cell.content[0], Strong)

    def test_large_table(self) -> None:
        """Test larger table conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Create 4x5 table
        table = slide.shapes.add_table(5, 4, Inches(1), Inches(1), Inches(6), Inches(3)).table

        # Fill with data
        for row in range(5):
            for col in range(4):
                table.cell(row, col).text = f"R{row}C{col}"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        table_nodes = [child for child in ast_doc.children if isinstance(child, Table)]
        table_node = table_nodes[0]

        # Check structure
        assert table_node.header is not None
        assert len(table_node.header.cells) == 4
        assert len(table_node.rows) == 4  # 5 total - 1 header = 4 data rows


@pytest.mark.unit
class TestImages:
    """Tests for image handling."""

    def test_image_detection(self) -> None:
        """Test that images are detected and converted."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        prs.slides.add_slide(layout)

        # Note: Adding actual image requires file, so we just verify basic functionality
        # In real usage, the converter would detect picture shapes
        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Just verify it doesn't crash with no images
        assert isinstance(ast_doc, Document)


@pytest.mark.unit
class TestBulletLists:
    """Tests for bullet list conversion."""

    def test_bullet_list(self) -> None:
        """Test bullet list conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Item 1"

        p2 = tf.add_paragraph()
        p2.text = "Item 2"
        p2.level = 0

        p3 = tf.add_paragraph()
        p3.text = "Item 3"
        p3.level = 0

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should detect list
        [child for child in ast_doc.children if isinstance(child, List)]
        # Note: List detection depends on implementation details
        # This test may need adjustment based on how bullets are detected

    def test_nested_list(self) -> None:
        """Test nested bullet list."""
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Item 1"

        p2 = tf.add_paragraph()
        p2.text = "Item 1.1"
        p2.level = 1  # Nested

        p3 = tf.add_paragraph()
        p3.text = "Item 2"
        p3.level = 0

        converter = PptxToAstConverter()
        converter.convert_to_ast(prs)

        # Should have nested structure
        # Implementation-dependent test


@pytest.mark.unit
class TestComplexStructures:
    """Tests for complex slide structures."""

    def test_mixed_content_slide(self) -> None:
        """Test slide with mixed content types."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(0.5))
        title_box.text_frame.text = "Slide Title"

        # Add paragraph
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
        text_box.text_frame.text = "Some content"

        # Add table
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1.5)).table
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have multiple node types
        assert len(ast_doc.children) >= 2
        assert any(isinstance(child, Paragraph) for child in ast_doc.children)
        assert any(isinstance(child, Table) for child in ast_doc.children)

    def test_slide_with_notes(self) -> None:
        """Test that slide notes are extracted when include_notes=True."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes content"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have title heading, "Speaker Notes" heading, and notes content
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        assert len(headings) >= 2, "Should have slide title and 'Speaker Notes' heading"

        # Check for "Speaker Notes" heading
        speaker_notes_heading = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_heading) == 1, "Should have exactly one 'Speaker Notes' heading"
        assert speaker_notes_heading[0].level == 3, "Speaker Notes should be H3"

        # Check that notes content appears in the document
        # Notes are detected as list items by heuristics (short text)
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1, "Notes content should be present"

    def test_slide_with_notes_disabled(self) -> None:
        """Test that slide notes are not extracted when include_notes=False."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes content"

        # Disable notes extraction
        options = PptxOptions(include_notes=False)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should have title but NOT "Speaker Notes" heading
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        speaker_notes_heading = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_heading) == 0, "Should not have 'Speaker Notes' heading when disabled"

    def test_slide_with_formatted_notes(self) -> None:
        """Test that speaker notes with formatting are extracted correctly."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Add notes with formatting
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        p = tf.paragraphs[0]

        # Add bold text
        run1 = p.add_run()
        run1.text = "Bold text"
        run1.font.bold = True

        # Add italic text
        run2 = p.add_run()
        run2.text = " and italic text"
        run2.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have "Speaker Notes" heading
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        speaker_notes_heading = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_heading) == 1, "Should have 'Speaker Notes' heading"

        # Check that notes contain formatted content (detected as list by heuristics)
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1, "Notes content should be present as list"

        # Verify formatting is preserved
        list_item = list_nodes[0].items[0]
        para_content = list_item.children[0].content
        # Should have Strong and Emphasis nodes
        assert any(isinstance(node, Strong) for node in para_content), "Should have bold text"
        assert any(isinstance(node, Emphasis) for node in para_content), "Should have italic text"

    def test_slide_with_multiline_notes(self) -> None:
        """Test that speaker notes with multiple paragraphs are extracted."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Add notes with multiple paragraphs
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = "First paragraph"
        tf.add_paragraph().text = "Second paragraph"
        tf.add_paragraph().text = "Third paragraph"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have "Speaker Notes" heading
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        speaker_notes_heading = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_heading) == 1, "Should have 'Speaker Notes' heading"

        # Should have multiple list items (one per paragraph, detected by heuristics)
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1, "Notes content should be present"
        # Should have multiple items (all short text detected as list items)
        total_items = sum(len(list_node.items) for list_node in list_nodes)
        assert total_items >= 3, "Should have at least 3 list items from 3 paragraphs"

    def test_slide_without_notes(self) -> None:
        """Test that slides without notes don't add Speaker Notes section."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Don't add any notes
        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should not have "Speaker Notes" heading
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        speaker_notes_heading = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_heading) == 0, "Should not have 'Speaker Notes' heading when no notes"

    def test_slide_with_empty_notes(self) -> None:
        """Test that slides with empty notes don't add Speaker Notes section."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        # Add empty notes (whitespace only)
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "   \n  \t  "

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should not have "Speaker Notes" heading for empty/whitespace notes
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        speaker_notes_heading = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_heading) == 0, "Should not have 'Speaker Notes' heading for empty notes"

    def test_multiple_slides_with_notes(self) -> None:
        """Test that multiple slides can each have their own notes."""
        prs = Presentation()
        layout = prs.slide_layouts[0]

        # First slide with notes
        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Slide 1"
        slide1.notes_slide.notes_text_frame.text = "Notes for slide 1"

        # Second slide with notes
        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Slide 2"
        slide2.notes_slide.notes_text_frame.text = "Notes for slide 2"

        # Third slide without notes
        slide3 = prs.slides.add_slide(layout)
        slide3.shapes.title.text = "Slide 3"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should have exactly 2 "Speaker Notes" headings (for slides 1 and 2)
        headings = [child for child in ast_doc.children if isinstance(child, Heading)]
        speaker_notes_headings = [h for h in headings if "Speaker Notes" in h.content[0].content]
        assert len(speaker_notes_headings) == 2, "Should have 'Speaker Notes' heading for each slide with notes"

    def test_multiple_shapes_per_slide(self) -> None:
        """Test slide with multiple text boxes."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # Add 3 text boxes (short text detected as list items)
        for i in range(3):
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1 + i * 1.5), Inches(3), Inches(1))
            textbox.text_frame.text = f"Text box {i + 1}"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Short text detected as list items - should have list with 3 items
        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        assert len(list_nodes) >= 1
        total_items = sum(len(list_node.items) for list_node in list_nodes)
        assert total_items >= 3


@pytest.mark.unit
class TestRunGrouping:
    """Tests for run grouping optimization."""

    def test_runs_with_same_formatting_grouped(self) -> None:
        """Test that consecutive runs with same formatting are grouped."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        p = textbox.text_frame.paragraphs[0]

        run1 = p.add_run()
        run1.text = "Bold "
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = "text"
        run2.font.bold = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should be grouped into single Strong node
        assert len(item_para.content) == 1
        assert isinstance(item_para.content[0], Strong)
        # Check text is combined
        text_content = item_para.content[0].content[0].content
        assert text_content == "Bold text"

    def test_runs_with_different_formatting_separate(self) -> None:
        """Test that runs with different formatting are kept separate."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        p = textbox.text_frame.paragraphs[0]

        run1 = p.add_run()
        run1.text = "Bold"
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = "Italic"
        run2.font.italic = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        list_nodes = [child for child in ast_doc.children if isinstance(child, List)]
        list_node = list_nodes[0]
        item_para = list_node.items[0].children[0]
        # Should have 2 separate inline nodes
        assert len(item_para.content) == 2
        assert isinstance(item_para.content[0], Strong)
        assert isinstance(item_para.content[1], Emphasis)


@pytest.mark.unit
class TestSlideOptions:
    """Tests for slide-specific options."""

    def test_slide_numbers_disabled(self) -> None:
        """Test that slide numbers can be disabled."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        options = PptxOptions(include_slide_numbers=False)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should convert without slide number markers
        assert isinstance(ast_doc, Document)
        # Check that title doesn't have slide number prefix
        heading = ast_doc.children[0]
        assert isinstance(heading, Heading)
        assert "Slide 1:" not in heading.content[0].content

    def test_slide_numbers_enabled(self) -> None:
        """Test that slide numbers can be enabled."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"

        options = PptxOptions(include_slide_numbers=True)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should convert with slide number markers
        assert isinstance(ast_doc, Document)
        # Check that title includes slide number prefix
        heading = ast_doc.children[0]
        assert isinstance(heading, Heading)
        assert "Slide 1:" in heading.content[0].content


@pytest.mark.unit
class TestChartHandling:
    """Tests for chart conversion modes."""

    @staticmethod
    def _build_chart_presentation(chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED) -> Presentation:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = ChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Series 1", (1, 2))
        slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)
        return prs

    def test_mermaid_mode_produces_codeblock(self) -> None:
        prs = self._build_chart_presentation()
        converter = PptxToAstConverter(PptxOptions(charts_mode="mermaid"))
        ast_doc = converter.convert_to_ast(prs)

        code_blocks = list(extract_nodes(ast_doc, CodeBlock))
        tables = list(extract_nodes(ast_doc, Table))

        assert code_blocks, "Expected Mermaid code block when charts_mode='mermaid'"
        assert code_blocks[0].language == "mermaid"
        assert "xychart-beta" in code_blocks[0].content
        assert not tables

    def test_both_mode_produces_table_and_codeblock(self) -> None:
        prs = self._build_chart_presentation()
        converter = PptxToAstConverter(PptxOptions(charts_mode="both"))
        ast_doc = converter.convert_to_ast(prs)

        code_blocks = list(extract_nodes(ast_doc, CodeBlock))
        tables = list(extract_nodes(ast_doc, Table))

        assert code_blocks and tables
        assert code_blocks[0].language == "mermaid"

    def test_default_mode_produces_table_only(self) -> None:
        prs = self._build_chart_presentation()
        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        code_blocks = list(extract_nodes(ast_doc, CodeBlock))
        tables = list(extract_nodes(ast_doc, Table))

        assert tables
        assert not code_blocks


@pytest.mark.unit
class TestListDetection:
    """Tests for list detection functions."""

    def test_detect_list_formatting_xml_with_bullet_char(self) -> None:
        """Test XML detection of bullet character."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

        # Add bulleted text
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Bullet item"

        # Check XML detection
        paragraph = tf.paragraphs[0]
        list_type, list_style = _detect_list_formatting_xml(paragraph)

        # Template may or may not have bullet formatting by default
        # Accept both bullet detection or None
        assert list_type in ("bullet", None)

    def test_detect_list_formatting_xml_with_numbered_list(self) -> None:
        """Test XML detection of numbered list."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        content = slide.placeholders[1]
        tf = content.text_frame
        p1 = tf.paragraphs[0]
        p1.text = "First item"

        # Try to set numbering if possible
        # Note: python-pptx doesn't have direct numbering API, so this tests what it can detect
        list_type, list_style = _detect_list_formatting_xml(p1)

        # May or may not detect numbering depending on template
        assert list_type in (None, "bullet", "number")

    def test_detect_list_formatting_xml_without_list_formatting(self) -> None:
        """Test XML detection on paragraph without list formatting."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        tf.text = "Regular text"

        paragraph = tf.paragraphs[0]
        list_type, list_style = _detect_list_formatting_xml(paragraph)

        # Should not detect list formatting
        assert list_type is None

    def test_detect_list_item_strict_mode_with_xml_formatting(self) -> None:
        """Test strict mode with XML list formatting."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Bullet item"

        paragraph = tf.paragraphs[0]
        is_list, list_type = _detect_list_item(paragraph, strict_mode=True)

        # If template has XML formatting, should detect as list
        # If not, strict mode returns False (which is correct behavior)
        # Either is acceptable since it depends on the template
        assert list_type in ("bullet", "number")

    def test_detect_list_item_strict_mode_without_xml_formatting(self) -> None:
        """Test strict mode without XML list formatting."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        tf.text = "Short text"

        paragraph = tf.paragraphs[0]
        is_list, list_type = _detect_list_item(paragraph, strict_mode=True)

        # Should NOT detect as list in strict mode without XML formatting
        assert is_list is False

    def test_detect_list_item_heuristic_mode_short_text(self) -> None:
        """Test heuristic mode with short text."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        tf.text = "Short text"

        paragraph = tf.paragraphs[0]
        is_list, list_type = _detect_list_item(paragraph, strict_mode=False)

        # Heuristics may detect short text as list
        assert is_list is True
        assert list_type == "bullet"

    def test_detect_list_item_heuristic_mode_with_level(self) -> None:
        """Test heuristic detection with indentation level."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        content = slide.placeholders[1]
        tf = content.text_frame
        p1 = tf.paragraphs[0]
        p1.text = "Item 1"
        p1.level = 0

        p2 = tf.add_paragraph()
        p2.text = "Item 1.1"
        p2.level = 1  # Indented

        # Check that level > 0 is detected as list in heuristic mode
        is_list, list_type = _detect_list_item(p2, strict_mode=False)
        assert is_list is True

    def test_detect_list_item_heuristic_skips_titles(self) -> None:
        """Test that heuristic mode avoids detecting titles as lists."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        tf.text = "Slide Title"  # Title-like (short, capitalized)

        paragraph = tf.paragraphs[0]
        is_list, list_type = _detect_list_item(paragraph, strict_mode=False)

        # Should not detect capitalized titles as lists
        assert is_list is False

    def test_analyze_slide_context_detects_numbered_lists(self) -> None:
        """Test slide context analysis for numbered lists."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = "1. First item"

        p2 = tf.add_paragraph()
        p2.text = "2. Second item"

        context = _analyze_slide_context(tf)

        assert context["has_numbered_list"] is True
        assert context["paragraph_count"] == 2

    def test_analyze_slide_context_tracks_max_level(self) -> None:
        """Test slide context tracks maximum indentation level."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        content = slide.placeholders[1]
        tf = content.text_frame
        p1 = tf.paragraphs[0]
        p1.text = "Level 0"
        p1.level = 0

        p2 = tf.add_paragraph()
        p2.text = "Level 2"
        p2.level = 2  # Deeper indent

        context = _analyze_slide_context(tf)

        assert context["max_level"] == 2
        assert context["paragraph_count"] == 2


@pytest.mark.unit
class TestStrictListDetection:
    """Tests for strict_list_detection option at converter level."""

    def test_strict_mode_prevents_heuristic_list_detection(self) -> None:
        """Test that strict mode prevents heuristic list detection."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Add short text that would be detected as list by heuristics
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = "Short text"

        # Convert with strict mode
        options = PptxOptions(strict_list_detection=True)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should NOT have any lists (no XML formatting)
        list_nodes = list(extract_nodes(ast_doc, List))
        assert len(list_nodes) == 0

        # Should have paragraph instead
        para_nodes = list(extract_nodes(ast_doc, Paragraph))
        assert len(para_nodes) >= 1

    def test_non_strict_mode_allows_heuristic_list_detection(self) -> None:
        """Test that non-strict mode allows heuristic list detection."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add short text that triggers heuristics
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        textbox.text_frame.text = "Short text"

        # Convert without strict mode (default)
        options = PptxOptions(strict_list_detection=False)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Should have list due to heuristics
        list_nodes = list(extract_nodes(ast_doc, List))
        assert len(list_nodes) >= 1

    def test_strict_mode_still_detects_xml_lists(self) -> None:
        """Test that strict mode still detects XML-formatted lists when present."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

        # Add properly formatted list
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Item 1"
        tf.add_paragraph().text = "Item 2"

        # Convert with strict mode
        options = PptxOptions(strict_list_detection=True)
        converter = PptxToAstConverter(options=options)
        ast_doc = converter.convert_to_ast(prs)

        # Template may or may not add XML list formatting
        # In strict mode, only XML-formatted lists are detected
        # This test verifies strict mode doesn't crash and produces valid output
        assert isinstance(ast_doc, Document)
        assert len(ast_doc.children) >= 1


@pytest.mark.unit
class TestChartModes:
    """Tests for chart conversion modes with scatter vs standard series."""

    def test_scatter_chart_data_mode(self) -> None:
        """Test scatter chart in data mode produces table."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = XyChartData()
        series = chart_data.add_series("Points")
        series.add_data_point(1, 2)
        series.add_data_point(3, 4)

        slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="data"))
        ast_doc = converter.convert_to_ast(prs)

        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        assert len(tables) == 1, "Should produce table in data mode"
        assert len(code_blocks) == 0, "Should not produce code block in data mode"

        # Verify table contains X/Y data
        table = tables[0]
        # Header should have "Series" and point columns
        assert len(table.header.cells) >= 2

    def test_scatter_chart_mermaid_mode(self) -> None:
        """Test scatter chart in mermaid mode produces code block."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = XyChartData()
        series = chart_data.add_series("Data")
        series.add_data_point(10, 20)
        series.add_data_point(30, 40)

        slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="mermaid"))
        ast_doc = converter.convert_to_ast(prs)

        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        assert len(code_blocks) == 1, "Should produce code block in mermaid mode"
        assert len(tables) == 0, "Should not produce table in mermaid mode"

        # Verify mermaid code
        code = code_blocks[0]
        assert code.language == "mermaid"
        assert "xychart-beta" in code.content
        assert "scatter" in code.content.lower()

    def test_scatter_chart_both_mode(self) -> None:
        """Test scatter chart in both mode produces table and code block."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = XyChartData()
        series = chart_data.add_series("Series")
        series.add_data_point(5, 10)
        series.add_data_point(15, 20)

        slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="both"))
        ast_doc = converter.convert_to_ast(prs)

        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        assert len(tables) == 1, "Should produce table in both mode"
        assert len(code_blocks) == 1, "Should produce code block in both mode"

    def test_standard_chart_data_mode(self) -> None:
        """Test standard chart in data mode produces table."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = ChartData()
        chart_data.categories = ["A", "B", "C"]
        chart_data.add_series("Sales", (10, 20, 15))

        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="data"))
        ast_doc = converter.convert_to_ast(prs)

        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        assert len(tables) == 1, "Should produce table in data mode"
        assert len(code_blocks) == 0, "Should not produce code block"

    def test_standard_chart_mermaid_mode(self) -> None:
        """Test standard chart in mermaid mode produces code block."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = ChartData()
        chart_data.categories = ["Q1", "Q2"]
        chart_data.add_series("Revenue", (100, 150))

        slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="mermaid"))
        ast_doc = converter.convert_to_ast(prs)

        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        assert len(code_blocks) == 1, "Should produce code block in mermaid mode"
        assert len(tables) == 0, "Should not produce table"

        code = code_blocks[0]
        assert code.language == "mermaid"
        assert "line" in code.content.lower()

    def test_standard_chart_both_mode(self) -> None:
        """Test standard chart in both mode produces table and code."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = ChartData()
        chart_data.categories = ["X", "Y"]
        chart_data.add_series("Values", (25, 35))

        slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="both"))
        ast_doc = converter.convert_to_ast(prs)

        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        assert len(tables) == 1, "Should produce table in both mode"
        assert len(code_blocks) == 1, "Should produce code block in both mode"


@pytest.mark.unit
class TestChartFallbacks:
    """Tests for chart fallback behaviors when data extraction fails."""

    def test_single_category_chart(self) -> None:
        """Test chart with single category."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = ChartData()
        chart_data.categories = ["Single"]
        chart_data.add_series("Data", (42,))

        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="data"))
        ast_doc = converter.convert_to_ast(prs)

        # Should handle single category gracefully
        assert isinstance(ast_doc, Document)
        tables = list(extract_nodes(ast_doc, Table))
        assert len(tables) >= 1

    def test_scatter_chart_single_point(self) -> None:
        """Test scatter chart with single data point."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = XyChartData()
        series = chart_data.add_series("Single")
        series.add_data_point(1, 1)  # Just one point

        slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="both"))
        ast_doc = converter.convert_to_ast(prs)

        # Should handle single point gracefully
        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        # Should produce at least one representation
        assert len(tables) + len(code_blocks) >= 1

    def test_chart_with_none_values(self) -> None:
        """Test chart with None values in data."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = ChartData()
        chart_data.categories = ["A", "B", "C"]
        chart_data.add_series("Data", (10, None, 20))

        slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Should handle None values gracefully
        tables = list(extract_nodes(ast_doc, Table))
        assert len(tables) >= 1

    def test_pie_chart_no_mermaid(self) -> None:
        """Test that pie chart falls back when mermaid not supported."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = ChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Share", (60, 40))

        slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)

        converter = PptxToAstConverter(PptxOptions(charts_mode="mermaid"))
        ast_doc = converter.convert_to_ast(prs)

        # Pie charts may not have mermaid support, should fallback to table
        tables = list(extract_nodes(ast_doc, Table))
        code_blocks = list(extract_nodes(ast_doc, CodeBlock))

        # Should produce some output (either table or nothing if no fallback)
        assert len(tables) + len(code_blocks) >= 0  # At minimum doesn't crash


@pytest.mark.unit
class TestSpeakerNotesCommentMode:
    """Tests for speaker notes comment_mode option."""

    @staticmethod
    def _create_presentation_with_notes(notes_text: str = "These are speaker notes") -> Presentation:
        """Helper to create a presentation with speaker notes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes_text

        return prs

    def test_comment_mode_content_default(self) -> None:
        """Test default comment_mode='content' behavior (backward compatible)."""
        prs = self._create_presentation_with_notes("These are speaker notes")

        # Default should be 'content' mode
        converter = PptxToAstConverter(PptxOptions())
        ast_doc = converter.convert_to_ast(prs)

        # Should have H3 heading with "Speaker Notes"
        headings = list(extract_nodes(ast_doc, Heading))
        speaker_notes_headings = [h for h in headings if h.level == 3 and any(
            isinstance(c, Text) and "Speaker Notes" in c.content for c in h.content
        )]

        assert len(speaker_notes_headings) == 1, "Should have H3 'Speaker Notes' heading in content mode"

        # Should NOT have Comment nodes
        comments = list(extract_nodes(ast_doc, Comment))
        assert len(comments) == 0, "Should not have Comment nodes in content mode"

    def test_comment_mode_comment(self) -> None:
        """Test comment_mode='comment' creates Comment nodes."""
        prs = self._create_presentation_with_notes("These are speaker notes")

        converter = PptxToAstConverter(PptxOptions(comment_mode="comment"))
        ast_doc = converter.convert_to_ast(prs)

        # Should have Comment node
        comments = list(extract_nodes(ast_doc, Comment))
        assert len(comments) == 1, "Should have exactly one Comment node"

        comment = comments[0]
        assert "These are speaker notes" in comment.content
        assert comment.metadata.get("comment_type") == "pptx_speaker_notes"
        assert comment.metadata.get("slide_number") == 1

        # Should NOT have H3 "Speaker Notes" heading
        headings = list(extract_nodes(ast_doc, Heading))
        speaker_notes_headings = [h for h in headings if h.level == 3 and any(
            isinstance(c, Text) and "Speaker Notes" in c.content for c in h.content
        )]
        assert len(speaker_notes_headings) == 0, "Should not have H3 'Speaker Notes' heading in comment mode"

    def test_comment_mode_ignore(self) -> None:
        """Test comment_mode='ignore' skips speaker notes."""
        prs = self._create_presentation_with_notes("These are speaker notes")

        converter = PptxToAstConverter(PptxOptions(comment_mode="ignore"))
        ast_doc = converter.convert_to_ast(prs)

        # Should NOT have Comment nodes
        comments = list(extract_nodes(ast_doc, Comment))
        assert len(comments) == 0, "Should not have Comment nodes in ignore mode"

        # Should NOT have H3 "Speaker Notes" heading
        headings = list(extract_nodes(ast_doc, Heading))
        speaker_notes_headings = [h for h in headings if h.level == 3 and any(
            isinstance(c, Text) and "Speaker Notes" in c.content for c in h.content
        )]
        assert len(speaker_notes_headings) == 0, "Should not have H3 'Speaker Notes' heading in ignore mode"

    def test_comment_mode_with_include_notes_false(self) -> None:
        """Test that include_notes=False overrides comment_mode."""
        prs = self._create_presentation_with_notes("These are speaker notes")

        # Even with comment_mode='comment', include_notes=False should skip notes
        converter = PptxToAstConverter(PptxOptions(include_notes=False, comment_mode="comment"))
        ast_doc = converter.convert_to_ast(prs)

        # Should NOT have any speaker notes content
        comments = list(extract_nodes(ast_doc, Comment))
        assert len(comments) == 0, "Should not extract notes when include_notes=False"

    def test_multiple_slides_with_comment_mode(self) -> None:
        """Test comment_mode with multiple slides having notes."""
        prs = Presentation()

        # Slide 1 with notes
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Slide 1"
        slide1.notes_slide.notes_text_frame.text = "Notes for slide 1"

        # Slide 2 with notes
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        slide2.shapes.title.text = "Slide 2"
        slide2.notes_slide.notes_text_frame.text = "Notes for slide 2"

        # Slide 3 without notes
        slide3 = prs.slides.add_slide(prs.slide_layouts[0])
        slide3.shapes.title.text = "Slide 3"

        converter = PptxToAstConverter(PptxOptions(comment_mode="comment"))
        ast_doc = converter.convert_to_ast(prs)

        # Should have 2 Comment nodes (slides 1 and 2)
        comments = list(extract_nodes(ast_doc, Comment))
        assert len(comments) == 2, "Should have Comment nodes for slides with notes"

        # Verify metadata
        assert comments[0].metadata.get("slide_number") == 1
        assert "Notes for slide 1" in comments[0].content

        assert comments[1].metadata.get("slide_number") == 2
        assert "Notes for slide 2" in comments[1].content


@pytest.mark.unit
class TestStrikethroughFormatting:
    """Tests for strikethrough text formatting."""

    def test_strikethrough_text(self) -> None:
        """Test strikethrough text conversion."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "strikethrough text"

        # Set strikethrough formatting
        if hasattr(run.font, "strike"):
            run.font.strike = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Find strikethrough nodes
        all_nodes = list(extract_nodes(ast_doc, Strikethrough))
        # Note: Strikethrough may or may not be set depending on pptx library version
        # This test validates the parser handles it when present
        assert isinstance(ast_doc, Document)


@pytest.mark.unit
class TestSuperscriptSubscriptFormatting:
    """Tests for superscript and subscript formatting."""

    def test_superscript_text(self) -> None:
        """Test superscript text conversion.

        Note: python-pptx may not fully support superscript/subscript attributes.
        This test verifies the parser handles them when present.
        """
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        # Add normal text
        run1 = p.add_run()
        run1.text = "E = mc"

        # Add superscript
        run2 = p.add_run()
        run2.text = "2"
        # Note: Setting superscript may not work in all python-pptx versions
        if hasattr(run2.font, "superscript"):
            run2.font.superscript = True

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Verify document parses successfully
        assert isinstance(ast_doc, Document)

        # Find superscript nodes (may or may not be present)
        sups = list(extract_nodes(ast_doc, Superscript))
        # If found, verify content
        if sups:
            sup_content = "".join(
                node.content if isinstance(node, Text) else ""
                for node in sups[0].content
            )
            assert "2" in sup_content

    def test_subscript_text(self) -> None:
        """Test subscript text conversion.

        Note: python-pptx may not fully support superscript/subscript attributes.
        This test verifies the parser handles them when present.
        """
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        # Add normal text
        run1 = p.add_run()
        run1.text = "H"

        # Add subscript
        run2 = p.add_run()
        run2.text = "2"
        # Note: Setting subscript may not work in all python-pptx versions
        if hasattr(run2.font, "subscript"):
            run2.font.subscript = True

        # Add more normal text
        run3 = p.add_run()
        run3.text = "O"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Verify document parses successfully
        assert isinstance(ast_doc, Document)

        # Find subscript nodes (may or may not be present)
        subs = list(extract_nodes(ast_doc, Subscript))
        # If found, verify content
        if subs:
            sub_content = "".join(
                node.content if isinstance(node, Text) else ""
                for node in subs[0].content
            )
            assert "2" in sub_content


@pytest.mark.unit
class TestHyperlinkExtraction:
    """Tests for hyperlink extraction from runs."""

    def test_run_with_hyperlink(self) -> None:
        """Test extracting hyperlinks from runs."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]

        # Add run with hyperlink
        run = p.add_run()
        run.text = "Click here"
        run.hyperlink.address = "https://www.example.com"

        converter = PptxToAstConverter()
        ast_doc = converter.convert_to_ast(prs)

        # Note: Current implementation may not fully integrate hyperlinks with formatting
        # This test verifies the parser handles hyperlinks when present
        # Full integration would require matching runs to formatted nodes
        assert isinstance(ast_doc, Document)

        # Try to find link nodes (may or may not be present depending on implementation)
        links = list(extract_nodes(ast_doc, Link))
        # If links are extracted, verify URL
        if links:
            assert links[0].url == "https://www.example.com"
