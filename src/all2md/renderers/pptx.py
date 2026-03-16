#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# src/all2md/renderers/pptx.py
"""PPTX rendering from AST.

This module provides the PptxRenderer class which converts AST nodes
to PowerPoint presentations. The renderer uses python-pptx to generate
.pptx files with proper slide layouts and formatting.

The rendering process splits the AST into slides using configurable
strategies (separator-based or heading-based), converts each slide's
content to PowerPoint shapes, and assembles a complete presentation.

"""

from __future__ import annotations

import logging
import os
import tempfile
from io import BytesIO
from pathlib import Path
from typing import IO, TYPE_CHECKING, Any, Union
from urllib.parse import urlparse

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide
    from pptx.text.text import TextFrame

from all2md.ast.nodes import (
    BlockQuote,
    Code,
    CodeBlock,
    Comment,
    CommentInline,
    DefinitionDescription,
    DefinitionList,
    DefinitionTerm,
    Document,
    Emphasis,
    FootnoteDefinition,
    FootnoteReference,
    Heading,
    HTMLBlock,
    HTMLInline,
    Image,
    LineBreak,
    Link,
    List,
    ListItem,
    MathBlock,
    MathInline,
    Node,
    Paragraph,
    Strikethrough,
    Strong,
    Subscript,
    Superscript,
    Table,
    TableCell,
    TableRow,
    Text,
    ThematicBreak,
    Underline,
)
from all2md.ast.visitors import NodeVisitor
from all2md.constants import (
    DEFAULT_PPTX_LINE_HEIGHT_FACTOR,
    DEFAULT_PPTX_MIN_TEXT_HEIGHT,
    DEPS_PPTX_RENDER,
)
from all2md.exceptions import RenderingError
from all2md.options.pptx import PptxRendererOptions
from all2md.renderers._split_utils import (
    auto_split_ast,
    extract_heading_text,
    split_ast_by_heading,
    split_ast_by_separator,
)
from all2md.renderers.base import BaseRenderer
from all2md.utils.decorators import requires_dependencies
from all2md.utils.images import decode_base64_image_to_file
from all2md.utils.network_security import fetch_image_securely, is_network_disabled

logger = logging.getLogger(__name__)


class PptxRenderer(NodeVisitor, BaseRenderer):
    """Render AST nodes to PPTX format.

    This class converts an AST document into a PowerPoint presentation
    using python-pptx. It splits the document into slides based on
    configured strategy and generates proper slide layouts and content.

    Parameters
    ----------
    options : PptxRendererOptions or None, default = None
        PPTX rendering options

    Examples
    --------
    Basic usage:

        >>> from all2md.options.pptx import PptxRendererOptions
        >>> from all2md.ast import Document, Heading, Paragraph, Text
        >>> from all2md.renderers.pptx import PptxRenderer
        >>> doc = Document(children=[
        ...     Heading(level=2, content=[Text(content="Slide 1")]),
        ...     Paragraph(content=[Text(content="Content here")])
        ... ])
        >>> options = PptxRendererOptions()
        >>> renderer = PptxRenderer(options)
        >>> renderer.render(doc, "output.pptx")


    """

    def __init__(self, options: PptxRendererOptions | None = None):
        """Initialize the PPTX renderer with options."""
        BaseRenderer._validate_options_type(options, PptxRendererOptions, "pptx")
        options = options or PptxRendererOptions()
        BaseRenderer.__init__(self, options)
        self.options: PptxRendererOptions = options

        # Rendering state
        self._current_textbox: TextFrame | None = None
        self._current_paragraph: Any = None
        self._current_slide: Any = None  # Current Slide object (for speaker notes routing)
        self._using_placeholder: bool = False  # True when rendering into a template placeholder
        self._list_ordered_stack: list[bool] = []  # Track ordered/unordered at each level
        self._list_item_counters: list[int] = []  # Track item number at each level for ordered lists
        self._temp_files: list[str] = []  # Track temp files for cleanup

    @requires_dependencies("pptx_render", DEPS_PPTX_RENDER)
    def render(self, doc: Document, output: Union[str, Path, IO[bytes]]) -> None:
        """Render the AST to a PPTX file.

        Parameters
        ----------
        doc : Document
            AST Document node to render
        output : str, Path, or IO[bytes]
            Output destination (file path or file-like object)

        Raises
        ------
        DependencyError
            If python-pptx is not installed
        RenderingError
            If PPTX generation fails

        """
        from pptx import Presentation
        from pptx.util import Inches, Pt

        # Store imports
        self._Inches = Inches
        self._Pt = Pt

        # Create presentation
        if self.options.template_path:
            prs = Presentation(self.options.template_path)
        else:
            prs = Presentation()
            # Set widescreen dimensions (default 16:9)
            prs.slide_width = Inches(self.options.slide_width)  # type: ignore[misc,assignment]
            prs.slide_height = Inches(self.options.slide_height)  # type: ignore[misc,assignment]

        # Set creator metadata if configured
        if self.options.creator:
            prs.core_properties.last_modified_by = self.options.creator

        # Split document into slides
        slides_data = self._split_into_slides(doc)

        # Create slides
        for idx, (heading, content_nodes) in enumerate(slides_data, start=1):
            self._create_slide(prs, heading, content_nodes, is_first=(idx == 1))

        # Save presentation
        try:
            if isinstance(output, (str, Path)):
                prs.save(str(output))
            else:
                # For file-like objects, save to BytesIO first
                buffer = BytesIO()
                prs.save(buffer)
                buffer.seek(0)
                output.write(buffer.read())
        except Exception as e:
            raise RenderingError(
                f"Failed to write PPTX file: {e!r}", rendering_stage="rendering", original_error=e
            ) from e
        finally:
            # Clean up temporary files
            for temp_file in self._temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.unlink(temp_file)
                except Exception:
                    pass  # Ignore cleanup errors

    @requires_dependencies("pptx_render", DEPS_PPTX_RENDER)
    def render_to_bytes(self, doc: Document) -> bytes:
        """Render the AST to PPTX bytes.

        Returns
        -------
        bytes
            PPTX file content as bytes

        """
        buffer = BytesIO()
        self.render(doc, buffer)
        return buffer.getvalue()

    def _split_into_slides(self, doc: Document) -> list[tuple[Heading | None, list[Node]]]:
        """Split AST document into slides based on configured strategy.

        Parameters
        ----------
        doc : Document
            AST document to split

        Returns
        -------
        list of tuple[Heading or None, list of Node]
            List of (heading, content_nodes) tuples

        """
        split_mode = self.options.slide_split_mode

        if split_mode == "separator":
            # Split on ThematicBreak nodes
            separator_chunks = split_ast_by_separator(doc)
            return [(None, chunk) for chunk in separator_chunks]

        elif split_mode == "heading":
            # Split on heading level
            return split_ast_by_heading(doc, heading_level=self.options.slide_split_heading_level)

        else:  # "auto"
            # Auto-detect best strategy
            return auto_split_ast(doc, heading_level=self.options.slide_split_heading_level)

    def _create_slide(
        self, prs: "Presentation", heading: Heading | None, content_nodes: list[Node], is_first: bool = False
    ) -> "Slide":
        """Create a slide with content.

        Parameters
        ----------
        prs : Presentation
            PowerPoint presentation object
        heading : Heading or None
            Slide heading (becomes title if use_heading_as_slide_title=True)
        content_nodes : list of Node
            AST nodes to render on slide
        is_first : bool, default False
            Whether this is the first slide (uses title slide layout)

        Returns
        -------
        Slide
            Created slide object

        """
        # Determine layout
        if is_first:
            layout_name = self.options.title_slide_layout
        else:
            layout_name = self.options.default_layout

        # Find layout (fallback to index 0 if not found)
        layout = None
        for slide_layout in prs.slide_layouts:
            if slide_layout.name == layout_name:
                layout = slide_layout
                break

        if layout is None:
            layout = prs.slide_layouts[0]  # Use first layout as fallback

        # Create slide
        slide = prs.slides.add_slide(layout)

        # Separate slide content from speaker notes (if enabled)
        if self.options.include_notes:
            slide_content, notes_content = self._extract_speaker_notes(content_nodes)
        else:
            slide_content = content_nodes
            notes_content = []

        # Set title
        title = ""
        nodes_to_render = slide_content

        if self.options.use_heading_as_slide_title:
            if heading:
                # Heading provided by splitting strategy (heading mode)
                title = extract_heading_text(heading)
            elif slide_content and isinstance(slide_content[0], Heading):
                # Extract first heading from content (separator mode)
                title = extract_heading_text(slide_content[0])
                # Don't render the heading again in content
                nodes_to_render = slide_content[1:]

        # Add title to slide if it has a title placeholder
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Render content nodes
        self._render_slide_content(slide, nodes_to_render, is_title_slide=is_first)

        # Add speaker notes if present
        if notes_content:
            self._write_speaker_notes(slide, notes_content)

        return slide

    def _render_slide_content(self, slide: "Slide", nodes: list[Node], *, is_title_slide: bool = False) -> None:
        """Render AST nodes as slide content.

        Dispatches to flow layout (default) or legacy fixed-position layout
        based on ``use_flow_layout`` option.  Title slides always use the
        fixed-position path so that template placeholders (e.g. subtitle)
        are preserved with their built-in styles.

        Parameters
        ----------
        slide : Slide
            Slide to add content to
        nodes : list of Node
            AST nodes to render
        is_title_slide : bool, default False
            Whether this is a title slide (forces fixed layout to preserve
            template subtitle placeholder and built-in styles).

        """
        self._current_slide = slide
        if self.options.use_flow_layout and not is_title_slide:
            self._render_slide_content_flow(slide, nodes)
        else:
            self._render_slide_content_fixed(slide, nodes)

    # ------------------------------------------------------------------
    # Legacy fixed-position layout
    # ------------------------------------------------------------------

    def _render_slide_content_fixed(self, slide: "Slide", nodes: list[Node]) -> None:
        """Render AST nodes using legacy fixed-position layout.

        All text nodes share a single text frame, and tables/images are placed
        at their configured fixed positions (``table_top``, ``image_top``).

        """
        # Try to find content placeholder with robust detection
        content_placeholder = None

        # Strategy 1: Look for text frame placeholders that aren't the title
        for shape in slide.placeholders:
            if not hasattr(shape, "text_frame"):
                continue

            # Skip title placeholder (usually idx 0)
            try:
                if shape.placeholder_format.idx == 0:
                    continue
            except Exception:
                pass

            # Try to find body/content placeholder by type
            try:
                # PP_PLACEHOLDER type 2 is body/content
                if hasattr(shape.placeholder_format, "type") and shape.placeholder_format.type == 2:
                    content_placeholder = shape
                    break
            except Exception:
                pass

            # Fallback: Use first non-title placeholder with text frame
            if content_placeholder is None and shape.has_text_frame:
                content_placeholder = shape

        if content_placeholder:
            # Use placeholder text frame (inherits template styles)
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear any default text
            self._current_textbox = text_frame
            self._using_placeholder = True
        else:
            # Create a text box for content
            from pptx.util import Inches

            opts = self.options
            content_width = opts.slide_width - opts.content_margin_left - opts.content_margin_right
            content_height = opts.slide_height - opts.content_margin_top - opts.content_margin_bottom
            left = Inches(opts.content_margin_left)
            top = Inches(opts.content_margin_top)
            width = Inches(content_width)
            height = Inches(content_height)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            self._current_textbox = textbox.text_frame
            self._using_placeholder = False

        # Render nodes
        for node in nodes:
            if isinstance(node, Table):
                self._render_table(slide, node)
            elif isinstance(node, Image):
                self._render_image(slide, node)
            else:
                # Render to text frame
                node.accept(self)

    # ------------------------------------------------------------------
    # Flow layout engine
    # ------------------------------------------------------------------

    @staticmethod
    def _group_into_layout_blocks(nodes: list[Node]) -> list[tuple[str, list[Node]]]:
        """Group consecutive nodes into layout blocks.

        Consecutive non-Table/non-Image nodes are merged into ``("text", [nodes])``
        blocks.  Table and Image nodes become standalone
        ``("table", [table])`` / ``("image", [image])`` blocks.

        Parameters
        ----------
        nodes : list of Node
            Flat list of AST nodes for a single slide.

        Returns
        -------
        list of tuple[str, list[Node]]
            Each element is ``(block_type, nodes)`` where *block_type* is one
            of ``"text"``, ``"table"``, or ``"image"``.

        """
        blocks: list[tuple[str, list[Node]]] = []
        text_acc: list[Node] = []

        for node in nodes:
            if isinstance(node, Table):
                if text_acc:
                    blocks.append(("text", text_acc))
                    text_acc = []
                blocks.append(("table", [node]))
            elif isinstance(node, Image):
                if text_acc:
                    blocks.append(("text", text_acc))
                    text_acc = []
                blocks.append(("image", [node]))
            else:
                text_acc.append(node)

        if text_acc:
            blocks.append(("text", text_acc))

        return blocks

    def _estimate_text_block_height(self, nodes: list[Node], available_width_inches: float) -> float:
        """Estimate the height of a text block in inches.

        Uses a simple heuristic based on character count and line wrapping.

        Parameters
        ----------
        nodes : list of Node
            Text-bearing AST nodes.
        available_width_inches : float
            Width available for text rendering.

        Returns
        -------
        float
            Estimated height in inches.

        """
        font_size_pt = self.options.default_font_size
        # Approximate characters per line: width in points / (font_size * ~0.5 avg char width)
        chars_per_line = max(1, int(available_width_inches * 72 / (font_size_pt * 0.5)))

        total_lines = 0
        for node in nodes:
            total_lines += self._count_lines_for_node(node, chars_per_line)

        line_height_inches = (font_size_pt / 72.0) * DEFAULT_PPTX_LINE_HEIGHT_FACTOR
        return max(total_lines * line_height_inches, DEFAULT_PPTX_MIN_TEXT_HEIGHT)

    def _count_lines_for_node(self, node: Node, chars_per_line: int) -> int:
        """Count estimated rendered lines for a single AST node.

        Parameters
        ----------
        node : Node
            AST node to estimate.
        chars_per_line : int
            Characters that fit on one line.

        Returns
        -------
        int
            Estimated number of lines (minimum 1).

        """
        if isinstance(node, CodeBlock):
            code = node.content if isinstance(node.content, str) else ""
            return max(1, code.count("\n") + 1)

        if isinstance(node, List):
            count = 0
            for child in node.items:
                text = self._extract_text_from_nodes(child.content if hasattr(child, "content") else [])
                count += max(1, -(-len(text) // chars_per_line))  # ceiling division
            return max(1, count)

        # Paragraph / other text-bearing nodes
        text = self._extract_text_from_nodes(node.content if hasattr(node, "content") else [])
        if not text:
            return 1
        return max(1, -(-len(text) // chars_per_line))  # ceiling division

    def _render_table_at(
        self, slide: "Slide", table: Table, left_inches: float, top_inches: float, width_inches: float
    ) -> float:
        """Render a table at an explicit position and return its height.

        Parameters
        ----------
        slide : Slide
            Target slide.
        table : Table
            AST table node.
        left_inches, top_inches, width_inches : float
            Position and width in inches.

        Returns
        -------
        float
            Table height in inches.

        """
        from pptx.util import Inches

        all_rows: list[TableRow] = []
        if table.header:
            all_rows.append(table.header)
        all_rows.extend(table.rows)

        if not all_rows:
            return 0.0

        num_rows = len(all_rows)
        num_cols = self._compute_table_columns(all_rows)
        if num_cols == 0 or num_rows == 0:
            return 0.0

        height_inches = self.options.table_height_per_row * num_rows

        left = Inches(left_inches)
        top = Inches(top_inches)
        width = Inches(width_inches)
        height = Inches(height_inches)

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        pptx_table = table_shape.table

        # Track which grid cells are occupied by spanning cells
        occupied = [[False] * num_cols for _ in range(num_rows)]

        for row_idx, ast_row in enumerate(all_rows):
            col_idx = 0
            for ast_cell in ast_row.cells:
                while col_idx < num_cols and occupied[row_idx][col_idx]:
                    col_idx += 1
                if col_idx >= num_cols:
                    break

                pptx_cell = pptx_table.rows[row_idx].cells[col_idx]
                cell_text = self._extract_text_from_nodes(ast_cell.content)
                pptx_cell.text = cell_text

                colspan = ast_cell.colspan
                rowspan = ast_cell.rowspan
                for r in range(row_idx, min(row_idx + rowspan, num_rows)):
                    for c in range(col_idx, min(col_idx + colspan, num_cols)):
                        occupied[r][c] = True

                if colspan > 1 or rowspan > 1:
                    end_row = min(row_idx + rowspan - 1, num_rows - 1)
                    end_col = min(col_idx + colspan - 1, num_cols - 1)
                    if end_row > row_idx or end_col > col_idx:
                        pptx_cell.merge(pptx_table.rows[end_row].cells[end_col])

                col_idx += colspan

        return height_inches

    def _resolve_image_file(self, image: Image) -> str | None:
        """Resolve an Image node to a local file path (or temp file).

        Returns
        -------
        str or None
            Path to the image file, or None if unresolvable.

        """
        if not image.url:
            return None

        if image.url.startswith("data:"):
            return self._decode_base64_image(image.url)
        elif urlparse(image.url).scheme in ("http", "https"):
            return self._fetch_remote_image(image.url)
        else:
            return image.url

    def _render_image_at(
        self, slide: "Slide", image: Image, left_inches: float, top_inches: float, max_width_inches: float
    ) -> float:
        """Render an image at an explicit position and return its height.

        Parameters
        ----------
        slide : Slide
            Target slide.
        image : Image
            AST image node.
        left_inches, top_inches : float
            Position in inches.
        max_width_inches : float
            Maximum width in inches.

        Returns
        -------
        float
            Image height in inches (0.0 if the image could not be placed).

        """
        from pptx.util import Emu, Inches

        image_file = self._resolve_image_file(image)
        if not image_file:
            return 0.0

        try:
            pic = slide.shapes.add_picture(
                image_file,
                Inches(left_inches),
                Inches(top_inches),
                width=Inches(min(self.options.image_width, max_width_inches)),
            )
            # add_picture returns a Picture shape; convert EMU height to inches
            return pic.height / Emu(914400)  # 914400 EMU per inch
        except Exception as e:
            logger.warning(f"Failed to add image to slide: {e}")
            if self.options.fail_on_resource_errors:
                raise RenderingError(
                    f"Failed to add image to slide: {e!r}",
                    rendering_stage="image_processing",
                    original_error=e,
                ) from e
            return 0.0

    @staticmethod
    def _find_content_placeholder(slide: "Slide") -> Any | None:
        """Find the first non-title content placeholder on a slide.

        Returns the placeholder shape, or ``None`` if none exists.
        """
        content_ph = None
        for shape in slide.placeholders:
            if not hasattr(shape, "text_frame"):
                continue
            try:
                if shape.placeholder_format.idx == 0:
                    continue
            except Exception:
                pass
            try:
                if hasattr(shape.placeholder_format, "type") and shape.placeholder_format.type == 2:
                    return shape
            except Exception:
                pass
            if content_ph is None and shape.has_text_frame:
                content_ph = shape
        return content_ph

    @staticmethod
    def _remove_content_placeholders(slide: "Slide", *, keep_idx: int | None = None) -> None:
        """Remove non-title content placeholders from a slide.

        Template slides typically contain placeholder shapes with default text
        like "Click here to add text".  When using flow layout we create our
        own textboxes, so these placeholders must be removed to avoid overlap.

        Only placeholders with ``idx >= 1`` are removed (idx 0 is the title).

        Parameters
        ----------
        slide : Slide
            Slide to remove placeholders from.
        keep_idx : int or None
            Placeholder index to keep (e.g. the content placeholder being
            reused for the first text block).

        """
        sp_tree = slide.shapes._spTree
        for ph in list(slide.placeholders):
            try:
                idx = ph.placeholder_format.idx
            except Exception:
                idx = -1
            if idx >= 1 and idx != keep_idx:
                sp_element = ph._element
                sp_tree.remove(sp_element)

    def _render_slide_content_flow(self, slide: "Slide", nodes: list[Node]) -> None:
        """Render AST nodes using flow layout (no overlap).

        Positions text blocks, tables, and images sequentially from top to
        bottom, advancing a y-cursor after each element.

        The first text block reuses the template's content placeholder (if one
        exists) so that built-in font styles are inherited.  Remaining content
        placeholders are removed to prevent overlap.

        """
        from pptx.util import Inches

        content_left = self.options.content_margin_left
        content_width = self.options.slide_width - self.options.content_margin_left - self.options.content_margin_right
        content_top = self.options.content_margin_top
        content_bottom = self.options.slide_height - self.options.content_margin_bottom

        blocks = self._group_into_layout_blocks(nodes)

        # Find content placeholder to reuse for the first text block
        content_ph = self._find_content_placeholder(slide)
        first_text_rendered = False

        # Determine placeholder index to keep
        keep_idx: int | None = None
        if content_ph is not None:
            try:
                keep_idx = content_ph.placeholder_format.idx
            except Exception:
                pass

        # Remove other content placeholders (keep the one we'll reuse)
        self._remove_content_placeholders(slide, keep_idx=keep_idx)

        y_cursor = content_top

        for block_type, block_data in blocks:
            if y_cursor >= content_bottom:
                logger.warning("Content exceeds slide area; remaining blocks truncated.")
                break

            if block_type == "text":
                if not first_text_rendered and content_ph is not None:
                    # Reuse content placeholder for first text block (inherits
                    # template styles like font size and family).
                    content_ph.left = Inches(content_left)
                    content_ph.top = Inches(y_cursor)
                    content_ph.width = Inches(content_width)
                    height = self._estimate_text_block_height(block_data, content_width)
                    height = min(height, content_bottom - y_cursor)
                    content_ph.height = Inches(height)
                    text_frame = content_ph.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True
                    self._current_textbox = text_frame
                    self._using_placeholder = True
                    first_text_rendered = True
                else:
                    height = self._estimate_text_block_height(block_data, content_width)
                    height = min(height, content_bottom - y_cursor)
                    textbox = slide.shapes.add_textbox(
                        Inches(content_left),
                        Inches(y_cursor),
                        Inches(content_width),
                        Inches(height),
                    )
                    textbox.text_frame.word_wrap = True
                    self._current_textbox = textbox.text_frame
                    self._using_placeholder = False

                for node in block_data:
                    node.accept(self)
                y_cursor += height + self.options.element_gap

            elif block_type == "table":
                table_node = block_data[0]
                assert isinstance(table_node, Table)
                table_height = self._render_table_at(slide, table_node, content_left, y_cursor, content_width)
                y_cursor += table_height + self.options.element_gap

            elif block_type == "image":
                image_node = block_data[0]
                assert isinstance(image_node, Image)
                try:
                    image_height = self._render_image_at(slide, image_node, content_left, y_cursor, content_width)
                except Exception as e:
                    logger.warning(f"Failed to render image {image_node.url}: {e}")
                    if self.options.fail_on_resource_errors:
                        raise RenderingError(
                            f"Failed to render image {image_node.url}: {e!r}",
                            rendering_stage="image_processing",
                            original_error=e,
                        ) from e
                    image_height = 0.0
                y_cursor += image_height + self.options.element_gap

        # Clean up: if the content placeholder was kept but never used
        # (e.g. no text blocks), remove it now so it doesn't show through
        if not first_text_rendered and content_ph is not None and keep_idx is not None:
            self._remove_content_placeholders(slide, keep_idx=None)

        # Reset placeholder flag
        self._using_placeholder = False

    # ------------------------------------------------------------------
    # Thin wrappers (used by fixed-position path)
    # ------------------------------------------------------------------

    def _render_table(self, slide: "Slide", table: Table) -> None:
        """Render a table at the fixed configured position."""
        self._render_table_at(slide, table, self.options.table_left, self.options.table_top, self.options.table_width)

    def _render_image(self, slide: "Slide", image: Image) -> None:
        """Render an image at the fixed configured position."""
        if not image.url:
            return

        try:
            opts = self.options
            self._render_image_at(slide, image, opts.image_left, opts.image_top, opts.image_width)
        except Exception as e:
            logger.warning(f"Failed to render image {image.url}: {e}")
            if self.options.fail_on_resource_errors:
                raise RenderingError(
                    f"Failed to render image {image.url}: {e!r}",
                    rendering_stage="image_processing",
                    original_error=e,
                ) from e

    def _decode_base64_image(self, data_uri: str) -> str | None:
        """Decode base64 image to temporary file.

        Parameters
        ----------
        data_uri : str
            Data URI with base64 encoded image

        Returns
        -------
        str or None
            Path to temporary file, or None if decoding failed

        """
        # Use centralized image utility
        temp_path = decode_base64_image_to_file(data_uri, delete_on_exit=False)
        if temp_path:
            self._temp_files.append(temp_path)
        else:
            logger.warning(f"Failed to decode base64 image: {data_uri[:50]}...")
            if self.options.fail_on_resource_errors:
                raise RenderingError(
                    "Failed to decode base64 image", rendering_stage="image_processing", original_error=None
                )
        return temp_path

    def _fetch_remote_image(self, url: str) -> str | None:
        """Fetch remote image securely.

        Parameters
        ----------
        url : str
            Remote image URL

        Returns
        -------
        str or None
            Path to temporary file with image data, or None if fetch failed

        """
        # Check global network disable flag
        if is_network_disabled():
            logger.debug(f"Network disabled, skipping remote image: {url}")
            return None

        # Check if remote fetching is allowed
        if not self.options.network.allow_remote_fetch:
            logger.debug(f"Remote fetching disabled, skipping image: {url}")
            return None

        try:
            # Fetch image data securely
            image_data = fetch_image_securely(
                url=url,
                allowed_hosts=self.options.network.allowed_hosts,
                require_https=self.options.network.require_https,
                max_size_bytes=self.options.max_asset_size_bytes,
                timeout=self.options.network.network_timeout,
                require_head_success=self.options.network.require_head_success,
            )

            # Determine file extension from URL or content type
            parsed = urlparse(url)
            path_lower = parsed.path.lower()

            if path_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg")):
                ext = path_lower.split(".")[-1]
            else:
                # Default to png
                ext = "png"

            # Write to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as f:
                f.write(image_data)
                temp_path = f.name

            self._temp_files.append(temp_path)
            logger.debug(f"Successfully fetched remote image: {url}")
            return temp_path

        except Exception as e:
            logger.warning(f"Failed to fetch remote image {url}: {e}")
            if self.options.fail_on_resource_errors:
                raise RenderingError(
                    f"Failed to fetch remote image {url}: {e!r}", rendering_stage="image_processing", original_error=e
                ) from e
            return None

    def _extract_speaker_notes(self, nodes: list[Node]) -> tuple[list[Node], list[Node]]:
        """Extract speaker notes from slide content nodes.

        Scans for an H3 heading with "Speaker Notes" text and splits the nodes
        into slide content (before the heading) and notes content (after the heading).

        Parameters
        ----------
        nodes : list of Node
            All nodes for the slide

        Returns
        -------
        tuple of (list of Node, list of Node)
            (slide_content_nodes, notes_content_nodes)

        """
        slide_content: list[Node] = []
        notes_content: list[Node] = []
        found_notes_heading = False

        for i, node in enumerate(nodes):
            # Check if this is a "Speaker Notes" heading (H3)
            if isinstance(node, Heading) and node.level == 3:
                # Extract heading text
                heading_text = self._extract_text_from_nodes(node.content).strip()
                if heading_text == "Speaker Notes":
                    # Found it! Everything after this (excluding the heading itself) is notes
                    found_notes_heading = True
                    # Capture remaining nodes as notes content
                    notes_content = nodes[i + 1 :]
                    break

            # If we haven't found the notes heading yet, add to slide content
            if not found_notes_heading:
                slide_content.append(node)

        return slide_content, notes_content

    def _write_speaker_notes(self, slide: "Slide", notes_nodes: list[Node]) -> None:
        """Write speaker notes to a slide.

        Parameters
        ----------
        slide : Slide
            Slide to add notes to
        notes_nodes : list of Node
            AST nodes to render as speaker notes

        """
        try:
            # Access the notes slide
            notes_slide = slide.notes_slide

            # Get the notes text frame
            notes_text_frame = notes_slide.notes_text_frame

            # Clear any existing notes
            notes_text_frame.clear()

            # Temporarily set current textbox to notes frame
            saved_textbox = self._current_textbox
            self._current_textbox = notes_text_frame

            # Render notes content using visitor pattern
            for node in notes_nodes:
                node.accept(self)

            # Restore previous textbox
            self._current_textbox = saved_textbox

        except Exception as e:
            # Log warning but don't fail rendering
            logger.warning(f"Failed to write speaker notes to slide: {e}")

    def _extract_text_from_nodes(self, nodes: list[Node]) -> str:
        """Extract plain text from inline nodes.

        Parameters
        ----------
        nodes : list of Node
            Inline nodes to extract text from

        Returns
        -------
        str
            Plain text content

        """
        text_parts: list[str] = []

        def collect_text(node_list: list[Node]) -> None:
            """Recursively collect text."""
            for node in node_list:
                if isinstance(node, Text):
                    text_parts.append(node.content)
                elif hasattr(node, "content"):
                    if isinstance(node.content, list):
                        collect_text(node.content)
                    elif isinstance(node.content, str):
                        text_parts.append(node.content)

        collect_text(nodes)
        return "".join(text_parts)

    # Visitor methods for rendering to text frame

    def visit_paragraph(self, node: Paragraph) -> None:
        """Render a Paragraph node.

        Parameters
        ----------
        node : Paragraph
            Paragraph to render

        """
        if not self._current_textbox:
            return

        # Add paragraph to text frame
        p = self._current_textbox.add_paragraph()
        self._current_paragraph = p

        # Set font size only when not using a template placeholder
        # (placeholders inherit styles from the slide master)
        if not self._using_placeholder:
            p.font.size = self._Pt(self.options.default_font_size)

        # Render content
        for child in node.content:
            child.accept(self)

        self._current_paragraph = None

    def visit_heading(self, node: Heading) -> None:
        """Render a Heading node (non-title headings).

        Parameters
        ----------
        node : Heading
            Heading to render

        """
        if not self._current_textbox:
            return

        # Create paragraph for heading
        p = self._current_textbox.add_paragraph()
        self._current_paragraph = p

        # Make it bold and larger
        p.font.bold = True
        if not self._using_placeholder:
            p.font.size = self._Pt(self.options.default_font_size + 4)

        # Render content
        for child in node.content:
            child.accept(self)

        self._current_paragraph = None

    def visit_text(self, node: Text) -> None:
        """Render a Text node.

        Parameters
        ----------
        node : Text
            Text to render

        """
        if not self._current_paragraph:
            return

        run = self._current_paragraph.add_run()
        run.text = node.content

    def visit_strong(self, node: Strong) -> None:
        """Render a Strong node.

        Parameters
        ----------
        node : Strong
            Strong to render

        """
        if not self._current_paragraph:
            return

        # Render content with bold
        for child in node.content:
            if isinstance(child, Text):
                run = self._current_paragraph.add_run()
                run.text = child.content
                run.font.bold = True
            else:
                child.accept(self)

    def visit_emphasis(self, node: Emphasis) -> None:
        """Render an Emphasis node.

        Parameters
        ----------
        node : Emphasis
            Emphasis to render

        """
        if not self._current_paragraph:
            return

        # Render content with italic
        for child in node.content:
            if isinstance(child, Text):
                run = self._current_paragraph.add_run()
                run.text = child.content
                run.font.italic = True
            else:
                child.accept(self)

    def visit_code(self, node: Code) -> None:
        """Render a Code node.

        Parameters
        ----------
        node : Code
            Code to render

        """
        if not self._current_paragraph:
            return

        run = self._current_paragraph.add_run()
        run.text = node.content
        run.font.name = "Courier New"

    def visit_code_block(self, node: CodeBlock) -> None:
        """Render a CodeBlock node.

        Parameters
        ----------
        node : CodeBlock
            Code block to render

        """
        if not self._current_textbox:
            return

        p = self._current_textbox.add_paragraph()
        run = p.add_run()
        run.text = node.content
        run.font.name = "Courier New"
        if not self._using_placeholder:
            run.font.size = self._Pt(self.options.default_font_size - 2)

    def visit_list(self, node: List) -> None:
        """Render a List node.

        Parameters
        ----------
        node : List
            List to render

        """
        # Track ordered/unordered and initialize counter for ordered lists
        self._list_ordered_stack.append(node.ordered)
        if node.ordered:
            self._list_item_counters.append(node.start if hasattr(node, "start") else 1)
        else:
            self._list_item_counters.append(0)  # Not used for unordered

        # Render list items
        for item in node.items:
            item.accept(self)
            # Increment counter for ordered lists
            if node.ordered and self._list_item_counters:
                self._list_item_counters[-1] += 1

        # Clean up state
        self._list_ordered_stack.pop()
        self._list_item_counters.pop()

    def visit_list_item(self, node: ListItem) -> None:
        """Render a ListItem node.

        Parameters
        ----------
        node : ListItem
            List item to render

        Notes
        -----
        **Bullet Handling in Text Boxes vs Placeholders:**

        PowerPoint behaves differently for text boxes vs content placeholders:

        - **Content placeholders**: Bullets are enabled by default, setting ``p.level``
          is sufficient to display bullets at the appropriate nesting level.
        - **Text boxes**: Bullets are disabled by default. Setting ``p.level`` alone
          will NOT show bullets. We must explicitly enable bullets via OOXML.

        This implementation uses python-pptx's OOXML API to programmatically enable
        bullets for unordered lists in all contexts (both placeholders and text boxes).

        """
        if not self._current_textbox:
            return

        # Create paragraph for list item
        p = self._current_textbox.add_paragraph()

        # Determine if this is ordered or unordered
        is_ordered = self._list_ordered_stack[-1] if self._list_ordered_stack else False

        # Calculate nesting level (depth in the stack)
        nesting_level = len(self._list_ordered_stack) - 1
        # Limit to reasonable depth for PowerPoint (0-8)
        nesting_level = min(nesting_level, 8)

        if is_ordered:
            # For ordered lists, manually add number prefix since python-pptx
            # has limited support for numbered lists
            item_number = self._list_item_counters[-1] if self._list_item_counters else 1

            # Set indentation level but no bullet (we're adding our own number)
            p.level = nesting_level

            # Add numbered prefix as the first run with configurable spacing
            run = p.add_run()
            spaces = " " * self.options.list_number_spacing
            run.text = f"{item_number}.{spaces}"
            run.font.bold = True  # Make number bold for visibility

            # Apply configurable indentation for nested lists
            # Note: Actual spacing may vary across templates
            if nesting_level > 0:
                try:
                    # python-pptx's level property handles indentation automatically,
                    # but templates may override this. The list_indent_per_level option
                    # documents the intent, though actual rendering depends on template.
                    p.level = nesting_level  # Already set above, but explicit for clarity
                except Exception as e:
                    logger.debug(f"Failed to apply list indentation: {e}")
        else:
            # For unordered lists, use PowerPoint's built-in bullet system
            p.level = nesting_level

            # Explicitly enable bullets via OOXML for text boxes (if enabled)
            # This ensures bullets appear in both text boxes and content placeholders
            if self.options.force_textbox_bullets:
                try:
                    # Access paragraph properties element
                    pPr = p._element.get_or_add_pPr()
                    # Enable bullet numbering by adding/updating buFont, buChar, or buAutoNum
                    # For simple bullets, we add a buChar element (bullet character)
                    from pptx.oxml import parse_xml

                    # Check if bullet is already configured
                    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                    if pPr.find(".//a:buChar", namespaces=ns) is None:
                        # Add bullet character (standard bullet: U+2022)
                        xml_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                        bu_char_xml = f'<a:buChar xmlns:a="{xml_ns}" char="\u2022"/>'
                        bu_char = parse_xml(bu_char_xml)
                        pPr.append(bu_char)
                except Exception as e:
                    # If OOXML manipulation fails, log warning but continue
                    # Bullets may not appear in text boxes, but rendering won't fail
                    logger.debug(f"Failed to enable bullets via OOXML: {e}")

            # Apply configurable indentation for nested lists
            # Note: list_indent_per_level option documents intent, but actual
            # rendering depends on template as python-pptx's level property
            # handles indentation automatically based on template settings
            if nesting_level > 0:
                # Indentation is controlled by p.level which was set above
                # The list_indent_per_level option documents expected behavior
                pass

        self._current_paragraph = p

        # Render children
        for child in node.children:
            if isinstance(child, Paragraph):
                # Render paragraph content directly
                for inline in child.content:
                    inline.accept(self)
            elif isinstance(child, List):
                # Nested list - will be handled by visit_list
                child.accept(self)
            else:
                child.accept(self)

        self._current_paragraph = None

    # Stub methods for other node types

    def visit_document(self, node: Document) -> None:
        """Document handled by render() method."""
        pass

    def visit_link(self, node: Link) -> None:
        """Render link as plain text (hyperlinks not implemented yet)."""
        for child in node.content:
            child.accept(self)

    def visit_image(self, node: Image) -> None:
        """Images handled separately by _render_image()."""
        pass

    def visit_table(self, node: Table) -> None:
        """Tables handled separately by _render_table()."""
        pass

    def visit_table_row(self, node: TableRow) -> None:
        """Handle table row (delegated to visit_table)."""
        pass

    def visit_table_cell(self, node: TableCell) -> None:
        """Handle table cell (delegated to visit_table)."""
        pass

    def visit_line_break(self, node: LineBreak) -> None:
        """Render line break."""
        if self._current_paragraph:
            self._current_paragraph.add_run().text = "\n"

    def visit_underline(self, node: Underline) -> None:
        """Render underline."""
        for child in node.content:
            child.accept(self)

    def visit_block_quote(self, node: "BlockQuote") -> None:
        """Render block quote as indented text.

        Parameters
        ----------
        node : BlockQuote
            Block quote to render

        """
        if not self._current_textbox:
            return

        # Render children with indentation
        for child in node.children:
            child.accept(self)

    def visit_thematic_break(self, node: "ThematicBreak") -> None:
        """Render thematic break as separator line.

        Parameters
        ----------
        node : ThematicBreak
            Thematic break to render

        """
        if not self._current_textbox:
            return

        # Add a separator paragraph
        p = self._current_textbox.add_paragraph()
        run = p.add_run()
        run.text = "---"

    def visit_strikethrough(self, node: "Strikethrough") -> None:
        """Render strikethrough text.

        Parameters
        ----------
        node : Strikethrough
            Strikethrough to render

        """
        if not self._current_paragraph:
            return

        # Render content with strikethrough (python-pptx doesn't support strikethrough directly)
        for child in node.content:
            child.accept(self)

    def visit_subscript(self, node: "Subscript") -> None:
        """Render subscript text.

        Parameters
        ----------
        node : Subscript
            Subscript to render

        """
        if not self._current_paragraph:
            return

        # Render content (python-pptx has limited subscript support)
        for child in node.content:
            child.accept(self)

    def visit_superscript(self, node: "Superscript") -> None:
        """Render superscript text.

        Parameters
        ----------
        node : Superscript
            Superscript to render

        """
        if not self._current_paragraph:
            return

        # Render content (python-pptx has limited superscript support)
        for child in node.content:
            child.accept(self)

    def visit_html_block(self, node: "HTMLBlock") -> None:
        """Skip HTML blocks in PPTX rendering.

        Parameters
        ----------
        node : HTMLBlock
            HTML block to skip

        """
        pass

    def visit_html_inline(self, node: "HTMLInline") -> None:
        """Skip inline HTML in PPTX rendering.

        Parameters
        ----------
        node : HTMLInline
            Inline HTML to skip

        """
        pass

    def visit_footnote_reference(self, node: "FootnoteReference") -> None:
        """Render footnote reference as superscript number.

        Parameters
        ----------
        node : FootnoteReference
            Footnote reference to render

        """
        if not self._current_paragraph:
            return

        run = self._current_paragraph.add_run()
        run.text = f"[{node.identifier}]"

    def visit_footnote_definition(self, node: "FootnoteDefinition") -> None:
        """Skip footnote definitions in PPTX rendering.

        Parameters
        ----------
        node : FootnoteDefinition
            Footnote definition to skip

        """
        pass

    def visit_math_inline(self, node: "MathInline") -> None:
        """Render inline math as plain text.

        Parameters
        ----------
        node : MathInline
            Inline math to render

        """
        if not self._current_paragraph:
            return

        # Get preferred representation (prefer latex, fallback to others)
        content, notation = node.get_preferred_representation("latex")

        run = self._current_paragraph.add_run()
        # Wrap latex in $ delimiters for clarity
        if notation == "latex":
            run.text = f"${content}$"
        else:
            run.text = content

    def visit_math_block(self, node: "MathBlock") -> None:
        """Render math block as plain text.

        Parameters
        ----------
        node : MathBlock
            Math block to render

        """
        if not self._current_textbox:
            return

        # Get preferred representation (prefer latex, fallback to others)
        content, notation = node.get_preferred_representation("latex")

        p = self._current_textbox.add_paragraph()
        run = p.add_run()
        # Wrap latex in $$ delimiters for clarity
        if notation == "latex":
            run.text = f"$${content}$$"
        else:
            run.text = content

    def visit_definition_list(self, node: "DefinitionList") -> None:
        """Render definition list.

        Parameters
        ----------
        node : DefinitionList
            Definition list to render

        """
        for term, descriptions in node.items:
            term.accept(self)
            for desc in descriptions:
                desc.accept(self)

    def visit_definition_term(self, node: "DefinitionTerm") -> None:
        """Render definition term as bold text.

        Parameters
        ----------
        node : DefinitionTerm
            Definition term to render

        """
        if not self._current_textbox:
            return

        p = self._current_textbox.add_paragraph()
        self._current_paragraph = p
        p.font.bold = True

        for child in node.content:
            child.accept(self)

        self._current_paragraph = None

    def visit_definition_description(self, node: "DefinitionDescription") -> None:
        """Render definition description as indented text.

        Parameters
        ----------
        node : DefinitionDescription
            Definition description to render

        """
        if not self._current_textbox:
            return

        p = self._current_textbox.add_paragraph()
        p.level = 1
        self._current_paragraph = p

        for child in node.content:
            child.accept(self)

        self._current_paragraph = None

    def visit_comment(self, node: Comment) -> None:
        """Render a Comment node according to comment_mode option.

        Parameters
        ----------
        node : Comment
            Comment block to render

        Notes
        -----
        Supports multiple rendering modes via comment_mode option:
        - "speaker_notes": Render in slide speaker notes (default)
        - "visible": Render as visible italic text in slide content
        - "ignore": Skip comment entirely

        """
        comment_mode = self.options.comment_mode

        if comment_mode == "ignore":
            return

        # Format comment with metadata
        comment_parts = []

        # Add metadata header if available
        if node.metadata.get("author") or node.metadata.get("date") or node.metadata.get("label"):
            header_parts = []
            if node.metadata.get("label"):
                header_parts.append(f"Comment {node.metadata['label']}")
            else:
                header_parts.append("Comment")

            if node.metadata.get("author"):
                header_parts.append(f"by {node.metadata['author']}")

            if node.metadata.get("date"):
                header_parts.append(f"({node.metadata['date']})")

            comment_parts.append(" ".join(header_parts))
            comment_parts.append(": ")

        comment_parts.append(node.content)
        comment_text = "".join(comment_parts)

        if comment_mode == "speaker_notes":
            # Write to the slide's speaker notes pane
            if self._current_slide is not None:
                try:
                    notes_slide = self._current_slide.notes_slide
                    notes_tf = notes_slide.notes_text_frame
                    p = notes_tf.add_paragraph()
                    run = p.add_run()
                    run.text = comment_text
                except Exception as e:
                    logger.warning(f"Failed to write comment to speaker notes: {e}")
        elif comment_mode == "visible":
            # Render as visible italic text in slide content
            if self._current_textbox:
                p = self._current_textbox.add_paragraph()
                run = p.add_run()
                run.text = comment_text
                run.font.italic = True

    def visit_comment_inline(self, node: CommentInline) -> None:
        """Render a CommentInline node according to comment_mode option.

        Parameters
        ----------
        node : CommentInline
            Inline comment to render

        Notes
        -----
        Supports multiple rendering modes via comment_mode option:
        - "speaker_notes": Render as italic text (for speaker notes context)
        - "visible": Render as italic text (for slide content context)
        - "ignore": Skip comment entirely

        Both speaker_notes and visible modes render the same way for inline
        comments - as italic text within the current paragraph.

        """
        comment_mode = self.options.comment_mode

        if comment_mode == "ignore":
            return

        if not self._current_paragraph:
            return

        # Format comment with metadata
        comment_parts = []

        # Add metadata prefix if available
        if node.metadata.get("author"):
            prefix_parts = []
            if node.metadata.get("label"):
                prefix_parts.append(f"[Comment {node.metadata['label']}")
            else:
                prefix_parts.append("[Comment")

            prefix_parts.append(f"by {node.metadata['author']}")

            if node.metadata.get("date"):
                prefix_parts.append(f"({node.metadata['date']})")

            comment_parts.append(" ".join(prefix_parts))
            comment_parts.append(": ")
        elif node.metadata.get("label"):
            comment_parts.append(f"[Comment {node.metadata['label']}: ")

        comment_parts.append(node.content)

        if node.metadata.get("author") or node.metadata.get("label"):
            comment_parts.append("]")

        comment_text = "".join(comment_parts)

        if comment_mode == "speaker_notes":
            # Route to speaker notes instead of inline on the slide
            if self._current_slide is not None:
                try:
                    notes_slide = self._current_slide.notes_slide
                    notes_tf = notes_slide.notes_text_frame
                    p = notes_tf.add_paragraph()
                    run = p.add_run()
                    run.text = comment_text
                except Exception as e:
                    logger.warning(f"Failed to write inline comment to speaker notes: {e}")
        else:
            # "visible" mode: render as italic text run in slide content
            run = self._current_paragraph.add_run()
            run.text = comment_text
            run.font.italic = True
