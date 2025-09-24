"""PowerPoint presentation to Markdown conversion module.

This module provides functionality to extract content from Microsoft PowerPoint
presentations (PPTX format) and convert it to structured Markdown. It processes
slides individually, extracting text content, formatting, tables, and embedded
images while preserving the logical structure and hierarchy of the presentation.

The converter analyzes each slide's layout and content, handling various slide
elements including title text, body content, tables, charts, and embedded media.
Special attention is paid to maintaining the presentation's logical flow and
structure in the resulting Markdown format.

Key Features
------------
- Slide-by-slide content extraction and conversion
- Text formatting preservation (bold, italic, lists)
- Table extraction with proper Markdown table formatting
- Image embedding as base64 data URIs
- Chart and graphic frame handling
- Hierarchical content structure preservation
- Bullet point and numbering conversion
- Slide title and content separation

Supported Elements
------------------
- Text content with formatting (bold, italic, underline)
- Slide titles and subtitles
- Bulleted and numbered lists with indentation
- Tables with cell content and structure
- Images and graphics (embedded as base64)
- Charts and diagrams (metadata extraction)
- Text boxes and content placeholders

Processing Features
-------------------
- Automatic slide numbering and separation
- Paragraph-level formatting detection
- List indentation and nesting preservation
- Image format detection and encoding
- Error handling for corrupted or complex elements
- Logging for debugging and issue tracking

Dependencies
------------
- python-pptx: For PowerPoint file parsing and content extraction
- imghdr: For image format detection
- base64: For image encoding
- logging: For debug and error reporting

Examples
--------
Basic presentation conversion:

    >>> from pptx2markdown import pptx_to_markdown
    >>> with open('presentation.pptx', 'rb') as f:
    ...     markdown = pptx_to_markdown(f)
    >>> print(markdown)

Note
----
Requires python-pptx package. Complex animations, transitions, and some
advanced PowerPoint features will be omitted in the conversion process.
The focus is on extracting textual and structural content.
"""

#  Copyright (c) 2023-2025 Tom Villani, Ph.D.
#
#  Permission is hereby granted, free of charge, to any person obtaining a copy of this software and associated
#  documentation files (the “Software”), to deal in the Software without restriction, including without limitation
#  the rights to use, copy, modify, merge, publish, distribute, sublicense, and/or sell copies of the Software,
#  and to permit persons to whom the Software is furnished to do so, subject to the following conditions:
#
#  The above copyright notice and this permission notice shall be included in all copies or substantial
#  portions of the Software.
#
#  THE SOFTWARE IS PROVIDED “AS IS”, WITHOUT WARRANTY OF ANY KIND, EXPRESS OR IMPLIED, INCLUDING
#  BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT.
#  IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY,
#  WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN CONNECTION WITH THE
#  SOFTWARE OR THE USE OR OTHER DEALINGS IN THE SOFTWARE.

import base64
import imghdr
import logging
from typing import Any, Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Inches

from ._input_utils import validate_and_convert_input, escape_markdown_special, format_special_text
from .constants import DEFAULT_CONVERT_IMAGES_TO_BASE64, DEFAULT_SLIDE_NUMBERS
from .exceptions import MdparseInputError, MdparseConversionError
from .options import PptxOptions, MarkdownOptions

logger = logging.getLogger(__name__)


def _process_paragraph_format(paragraph: Any) -> list[tuple[str, str]]:
    """Extract formatting information from a paragraph."""
    formats = []
    if paragraph.level:  # Handle indentation levels
        formats.append(("  " * paragraph.level, ""))
    return formats


def _process_run_format(run: Any) -> list[tuple[str, str]]:
    """Get formatting markers for text run."""
    formats = []
    if run.font.bold:
        formats.append(("**", "**"))
    if run.font.italic:
        formats.append(("*", "*"))
    if run.font.underline:
        formats.append(("__", "__"))
    return formats


def _process_text_frame(frame: Any) -> str:
    """Convert a text frame to markdown, preserving formatting."""
    lines = []

    for paragraph in frame.paragraphs:
        if not paragraph.text.strip():
            continue

        para_formats = _process_paragraph_format(paragraph)
        text_parts = []

        for run in paragraph.runs:
            text = run.text.strip()
            if not text:
                continue

            # Apply run-level formatting
            for start, end in _process_run_format(run):
                text = f"{start}{text}{end}"

            text_parts.append(text)

        # Combine runs and apply paragraph-level formatting
        if text_parts:
            line = " ".join(text_parts)
            for start, end in para_formats:
                line = f"{start}{line}{end}"
            lines.append(line)

    return "\n".join(lines)


def _process_table(table: Any) -> str:
    """Convert a table shape to markdown format."""
    markdown_rows = []

    # Process headers
    if table.rows:
        header_cells = []
        for cell in table.rows[0].cells:
            cell_text = _process_text_frame(cell.text_frame)
            header_cells.append(cell_text.replace("\n", " "))

        markdown_rows.append("| " + " | ".join(header_cells) + " |")
        markdown_rows.append("| " + " | ".join(["---"] * len(header_cells)) + " |")

        first = False
        # Process data rows
        for row in table.rows:
            if not first:
                first = True
                continue
            row_cells = []
            for cell in row.cells:
                cell_text = _process_text_frame(cell.text_frame)
                row_cells.append(cell_text.replace("\n", " "))
            markdown_rows.append("| " + " | ".join(row_cells) + " |")

    return "\n".join(markdown_rows)


def _extract_image_data(shape: Any) -> str:
    """Extract image data and convert to base64."""
    # try:
    image = shape.image
    image_bytes = image.blob

    # Determine image type from blob
    image_type = imghdr.what(None, h=image_bytes)

    # Convert to base64
    b64_data = base64.b64encode(image_bytes).decode("utf-8")
    return f"data:image/{image_type};base64,{b64_data}"
    # except Exception as e:
    #     logger.error(f"Failed to extract image data: {str(e)}")
    #     return None


def _process_shape(shape: Any, convert_images_to_base64: bool = False) -> str | None:
    """Process a single shape and convert to markdown."""
    # try:
    # Handle text boxes and other shapes with text
    if shape.has_text_frame:
        return _process_text_frame(shape.text_frame)

    # Handle tables
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return _process_table(shape.table)

    # Handle pictures
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        if convert_images_to_base64:
            image_data = _extract_image_data(shape)
            if image_data:
                alt_text = shape.alt_text or "image"
                return f"![{alt_text}]({image_data})"
        else:
            alt_text = shape.alt_text or "image"
            return f"![{alt_text}]()"

    # Handle charts (convert to tables if possible)
    elif isinstance(shape, GraphicFrame) and hasattr(shape, "chart"):
        chart = shape.chart
        # Basic chart data extraction
        data_rows = []

        # Extract categories (x-axis)
        categories = [cat.label for cat in chart.plots[0].categories if hasattr(cat, "label")]
        if categories:
            data_rows.append("| Category | " + " | ".join(str(cat) for cat in categories) + " |")
            data_rows.append("| --- | " + " | ".join(["---"] * len(categories)) + " |")

        # Extract series data
        for series in chart.series:
            values = list(series.values)
            data_rows.append(f"| {series.name} | " + " | ".join(str(val) for val in values) + " |")

        return "\n".join(data_rows)

    return None

    # except Exception as e:
    #     logger.error(f"Error processing shape: {str(e)}")
    #     return ''


def pptx_to_markdown(
    input_data: Union[str, Any],
    options: PptxOptions | None = None,
    convert_images_to_base64: bool | None = None,  # Deprecated, use options.convert_images_to_base64
    slide_numbers: bool | None = None  # Deprecated, use options.slide_numbers
) -> str:
    """Convert a PowerPoint presentation to Markdown format.

    Parameters
    ----------
    pptx_file : str or file-like object
        The PowerPoint file to convert
    convert_images_to_base64 : bool, default=False
        Whether to convert images to base64 data URLs
    slide_numbers : bool, default=True
        Whether to include slide numbers as headers

    Returns
    -------
    str
        The presentation content in Markdown format
    """
    # Handle backward compatibility and merge options
    if options is None:
        options = PptxOptions()

    # Handle deprecated parameters
    if convert_images_to_base64 is not None:
        options.convert_images_to_base64 = convert_images_to_base64
    if slide_numbers is not None:
        options.slide_numbers = slide_numbers

    # Validate and convert input
    try:
        from pptx.presentation import Presentation as PresentationType

        doc_input, input_type = validate_and_convert_input(
            input_data,
            supported_types=["path-like", "file-like", "pptx.Presentation objects"]
        )

        # Open presentation based on input type
        if input_type == "object" and isinstance(doc_input, PresentationType):
            prs = doc_input
        else:
            prs = Presentation(doc_input)

    except Exception as e:
        if "python-pptx" in str(e).lower() or "pptx" in str(e).lower():
            raise MdparseConversionError(
                "python-pptx library is required for PPTX conversion. "
                "Install with: pip install python-pptx",
                conversion_stage="dependency_check",
                original_error=e
            ) from e
        else:
            raise MdparseConversionError(
                f"Failed to open PPTX presentation: {str(e)}",
                conversion_stage="document_opening",
                original_error=e
            ) from e

    # Get Markdown options (create default if not provided)
    md_options = options.markdown_options or MarkdownOptions()

    markdown_content = []

    for i, slide in enumerate(prs.slides, 1):
        slide_content = []

        # Process slide title if present
        if slide.shapes.title:
            title_text = _process_text_frame(slide.shapes.title.text_frame)
            if slide_numbers:
                slide_content.append(f"# Slide {i}: {title_text.strip()}\n")
            else:
                slide_content.append(f"# {title_text.strip()}\n")

        # Process all shapes in the slide
        for shape in slide.shapes:
            shape_content = _process_shape(shape, convert_images_to_base64)
            if shape_content:
                slide_content.append(shape_content + "\n")

        # Add slide content with spacing
        if slide_content:
            markdown_content.extend(slide_content)
            markdown_content.append("\n---\n")  # Add separator between slides

    return "\n".join(markdown_content).strip()


def create_test_presentation() -> Any:
    """Create a test PowerPoint presentation with various elements."""
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Test PowerPoint Presentation"
    subtitle.text = "For Markdown Conversion Testing"

    # Content slide with bullet points
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes

    title = shapes.title
    title.text = "Bullet Points Test"

    body = shapes.placeholders[1]
    tf = body.text_frame

    p = tf.add_paragraph()
    p.text = "First Level Bullet"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Second Level Bullet"
    p.level = 1

    # Slide with table
    table_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = table_slide.shapes

    title = shapes.title
    title.text = "Table Test"

    rows, cols = 3, 3
    left = top = width = height = Inches(2.0)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # Add header row
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = f"Header {i + 1}"

    # Add data rows
    for row in range(1, rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = f"Cell {row},{col}"

    return prs
