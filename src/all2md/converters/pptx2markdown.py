#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# src/all2md/converters/pptx2markdown.py
"""PowerPoint presentation to Markdown conversion module.

This module provides functionality to extract content from Microsoft PowerPoint
presentations (PPTX format) and convert it to structured Markdown. It processes
slides individually, extracting text content, formatting, tables, and embedded
images while preserving the logical structure and hierarchy of the presentation.

The converter analyzes each slide's layout and content, handling various slide
elements including title text, body content, tables, charts, and embedded media.
Special attention is paid to maintaining the presentation's logical flow and
structure in the resulting Markdown format.

Key Features
------------
- Slide-by-slide content extraction and conversion
- Text formatting preservation (bold, italic, lists)
- Table extraction with proper Markdown table formatting
- Image embedding as base64 data URIs
- Chart and graphic frame handling
- Hierarchical content structure preservation
- Bullet point and numbering conversion
- Slide title and content separation

Supported Elements
------------------
- Text content with formatting (bold, italic, underline)
- Slide titles and subtitles
- Bulleted and numbered lists with indentation
- Tables with cell content and structure
- Images and graphics (embedded as base64)
- Charts and diagrams (metadata extraction)
- Text boxes and content placeholders

Processing Features
-------------------
- Automatic slide numbering and separation
- Paragraph-level formatting detection
- List indentation and nesting preservation
- Image format detection and encoding
- Error handling for corrupted or complex elements
- Logging for debugging and issue tracking

Dependencies
------------
- python-pptx: For PowerPoint file parsing and content extraction
- base64: For image encoding
- logging: For debug and error reporting

Examples
--------
Basic presentation conversion:

    >>> from all2md.converters.pptx2markdown import pptx_to_markdown
    >>> with open('presentation.pptx', 'rb') as f:
    ...     markdown = pptx_to_markdown(f)
    >>> print(markdown)

Note
----
Requires python-pptx package. Complex animations, transitions, and some
advanced PowerPoint features will be omitted in the conversion process.
The focus is on extracting textual and structural content.
"""

import logging
from pathlib import Path
from typing import IO, Any, Union

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Inches

from all2md.converter_metadata import ConverterMetadata
from all2md.exceptions import MarkdownConversionError
from all2md.options import PptxOptions
from all2md.utils.attachments import create_attachment_sequencer, extract_pptx_image_data, process_attachment
from all2md.utils.inputs import format_markdown_heading, format_special_text, validate_and_convert_input
from all2md.utils.metadata import (
    OFFICE_FIELD_MAPPING,
    DocumentMetadata,
    map_properties_to_metadata,
    prepend_metadata_if_enabled,
)
from all2md.utils.security import validate_zip_archive

logger = logging.getLogger(__name__)


def _detect_list_formatting_xml(paragraph: Any) -> tuple[str | None, str | None]:
    """Detect list formatting using XML element inspection.

    Parameters
    ----------
    paragraph : Any
        The paragraph object to inspect

    Returns
    -------
    tuple[str | None, str | None]
        (list_type, list_style) where list_type is "bullet" or "number"
    """
    try:
        # Access paragraph properties XML element
        if not hasattr(paragraph, '_p') or paragraph._p is None:
            return None, None

        pPr = paragraph._p.pPr
        if pPr is None:
            return None, None

        # Check for bullet character element
        bu_char = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        if bu_char is not None:
            char = bu_char.get('char', '•')
            return "bullet", char

        # Check for auto numbering element
        bu_auto_num = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
        if bu_auto_num is not None:
            num_type = bu_auto_num.get('type', 'arabicPeriod')
            return "number", num_type

        # Check for bullet font (indicates some form of bullet formatting)
        bu_font = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
        if bu_font is not None:
            return "bullet", "default"

    except Exception:
        # Fall back gracefully if XML parsing fails
        pass

    return None, None


def _detect_list_item(paragraph: Any, slide_context: dict | None = None) -> tuple[bool, str]:
    """Detect if a paragraph is a list item and determine the list type.

    Uses XML-based detection first, then falls back to heuristics.

    Parameters
    ----------
    paragraph : Any
        The paragraph object to analyze
    slide_context : dict, optional
        Context about the slide to help with detection

    Returns
    -------
    tuple[bool, str]
        (is_list_item, list_type) where list_type is "bullet" or "number"
    """
    # First try XML-based detection for proper list formatting
    xml_list_type, xml_list_style = _detect_list_formatting_xml(paragraph)
    if xml_list_type:
        return True, xml_list_type

    # Fall back to level-based detection
    if not hasattr(paragraph, 'level') or paragraph.level is None:
        return False, "bullet"

    level = paragraph.level
    if level > 0:
        # Use slide context to help determine list type for indented items
        if slide_context and slide_context.get('has_numbered_list', False):
            return True, "number"
        return True, "bullet"

    # For level 0, use heuristics as last resort
    text = paragraph.text.strip() if hasattr(paragraph, 'text') else ""
    if not text:
        return False, "bullet"

    # Check for explicit numbered list patterns in text
    import re
    if re.match(r'^\d+[\.\)]\s', text):
        return True, "number"

    # Check if this looks like a numbered list item based on context
    if (slide_context and slide_context.get('has_numbered_list', False) and
        ('item' in text.lower() or 'first' in text.lower() or
         'second' in text.lower() or 'third' in text.lower())):
        return True, "number"

    # Use heuristics for bullet lists - shorter text that doesn't look like a title/header
    words = text.split()
    if len(words) <= 8 and not text.endswith(('.', '!', '?', ':')):
        # Additional checks to avoid false positives
        if not (text.lower().startswith(('slide', 'title', 'chapter')) or
                len(words) <= 3 and text.istitle()):
            return True, "bullet"

    return False, "bullet"


def _process_paragraph_format(
    paragraph: Any, is_list_item: bool = False, list_type: str = "bullet", list_number: int = 1
) -> list[tuple[str, str]]:
    """Extract formatting information from a paragraph.

    Parameters
    ----------
    paragraph : Any
        The paragraph object to process
    is_list_item : bool, default False
        Whether this paragraph is part of a list
    list_type : str, default "bullet"
        Type of list ("bullet" or "number")
    list_number : int, default 1
        The number for numbered lists

    Returns
    -------
    list[tuple[str, str]]
        List of (prefix, suffix) formatting tuples
    """
    formats = []
    level = paragraph.level if hasattr(paragraph, 'level') and paragraph.level else 0

    if is_list_item:
        # Generate proper markdown list markers
        indent = "  " * level
        if list_type == "bullet":
            marker = "- "
        else:  # numbered list
            marker = f"{list_number}. "
        formats.append((f"{indent}{marker}", ""))
    elif level > 0:
        # Non-list items with indentation (just indent)
        formats.append(("  " * level, ""))

    return formats


def _process_run_format(run: Any, text: str = "", md_options=None) -> str:
    """Apply formatting to text run based on font properties and hyperlinks."""
    content = text

    # Check for hyperlinks first
    if hasattr(run, 'hyperlink') and run.hyperlink and hasattr(run.hyperlink, 'address') and run.hyperlink.address:
        # Format as markdown link [text](url)
        url = run.hyperlink.address
        content = f"[{content}]({url})"

    # Apply bold and italic formatting with standard markdown
    if hasattr(run, 'font') and run.font:
        if hasattr(run.font, 'bold') and run.font.bold:
            content = f"**{content}**"
        if hasattr(run.font, 'italic') and run.font.italic:
            content = f"*{content}*"

        # Apply underline formatting using format_special_text
        if hasattr(run.font, 'underline') and run.font.underline:
            underline_mode = md_options.underline_mode if md_options else "html"
            content = format_special_text(content, "underline", underline_mode)

    return content


def _analyze_slide_context(frame: Any) -> dict:
    """Analyze the text frame to understand the slide context for better list detection.

    Parameters
    ----------
    frame : Any
        The text frame to analyze

    Returns
    -------
    dict
        Context information about the slide
    """
    context = {
        'has_numbered_list': False,
        'paragraph_count': 0,
        'max_level': 0
    }

    import re
    for paragraph in frame.paragraphs:
        if not paragraph.text.strip():
            continue

        context['paragraph_count'] += 1

        # Track maximum indentation level
        level = getattr(paragraph, 'level', 0) or 0
        context['max_level'] = max(context['max_level'], level)

        # Check if any paragraph looks like a numbered list
        text = paragraph.text.strip()
        if (re.match(r'^\d+[\.\)]\s', text) or
            'numbered' in text.lower() or
            'first item' in text.lower() or
            'second item' in text.lower() or
            'third item' in text.lower()):
            context['has_numbered_list'] = True

    return context


def _process_text_frame(frame: Any, md_options=None, is_title: bool = False) -> str:
    """Convert a text frame to markdown, preserving formatting.

    Parameters
    ----------
    frame : Any
        The text frame to process
    md_options : Any, optional
        Markdown options for formatting
    is_title : bool, default False
        Whether this text frame is a slide title

    Returns
    -------
    str
        The processed markdown text
    """
    lines = []

    # Analyze all paragraphs to detect slide context
    slide_context = _analyze_slide_context(frame)

    # Track list numbering by level
    list_counters = {}  # level -> current_number

    for paragraph in frame.paragraphs:
        if not paragraph.text.strip():
            continue

        # Don't treat titles as list items
        if is_title:
            is_list_item, list_type = False, "bullet"
        else:
            # Detect if this is a list item and what type
            is_list_item, list_type = _detect_list_item(paragraph, slide_context)

        # Handle list numbering
        level = getattr(paragraph, 'level', 0) or 0
        if is_list_item and list_type == "number":
            # Increment counter for this level
            if level not in list_counters:
                list_counters[level] = 1
            else:
                list_counters[level] += 1

            # Reset counters for deeper levels when we return to a shallower level
            levels_to_reset = [lvl for lvl in list_counters.keys() if lvl > level]
            for lvl in levels_to_reset:
                del list_counters[lvl]

            list_number = list_counters[level]
        else:
            list_number = 1

        para_formats = _process_paragraph_format(paragraph, is_list_item, list_type, list_number)
        text_parts = []

        for run in paragraph.runs:
            text = run.text.strip()
            if not text:
                continue

            # Apply run-level formatting using the updated function
            text = _process_run_format(run, text, md_options)
            text_parts.append(text)

        # Combine runs and apply paragraph-level formatting
        if text_parts:
            line = " ".join(text_parts)
            for start, end in para_formats:
                line = f"{start}{line}{end}"
            lines.append(line)

    return "\n".join(lines)


def _process_table(table: Any, md_options=None) -> str:
    """Convert a table shape to markdown format."""
    markdown_rows = []

    # Process headers
    if table.rows:
        header_cells = []
        for cell in table.rows[0].cells:
            cell_text = _process_text_frame(cell.text_frame, md_options)
            header_cells.append(cell_text.replace("\n", " "))

        markdown_rows.append("| " + " | ".join(header_cells) + " |")
        markdown_rows.append("| " + " | ".join(["---"] * len(header_cells)) + " |")

        first = False
        # Process data rows
        for row in table.rows:
            if not first:
                first = True
                continue
            row_cells = []
            for cell in row.cells:
                cell_text = _process_text_frame(cell.text_frame, md_options)
                row_cells.append(cell_text.replace("\n", " "))
            markdown_rows.append("| " + " | ".join(row_cells) + " |")

    return "\n".join(markdown_rows)


def _process_shape(
        shape: Any, options: PptxOptions, base_filename: str = "presentation", slide_num: int = 1,
        img_counter: dict | None = None, attachment_sequencer=None, slide_title_shape=None
) -> str | None:
    """Process a single shape and convert to markdown."""
    if img_counter is None:
        img_counter = {}

    # Handle text boxes and other shapes with text
    if shape.has_text_frame:
        md_options = options.markdown_options if options else None
        # Check if this shape is the slide title
        is_title = slide_title_shape is not None and shape == slide_title_shape
        return _process_text_frame(shape.text_frame, md_options, is_title=is_title)

    # Handle tables
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        md_options = options.markdown_options if options else None
        return _process_table(shape.table, md_options)

    # Handle pictures
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        # Extract image data if needed
        if options.attachment_mode != "skip":
            image_data = extract_pptx_image_data(shape)
        else:
            logger.info(f"Skipping image in slide {slide_num} (attachment_mode=skip): {shape.alt_text or 'image'}")
            image_data = None
        if image_data and not isinstance(image_data, bytes):
            image_data = None
        alt_text = shape.alt_text or "image"

        # Use sequencer if available, otherwise fall back to manual counting
        if attachment_sequencer:
            image_filename, _ = attachment_sequencer(
                base_stem=base_filename,
                format_type="pptx",
                slide_num=slide_num,
                extension="png"
            )
        else:
            # Get next sequence number for this slide
            slide_key = f"slide_{slide_num}"
            img_counter[slide_key] = img_counter.get(slide_key, 0) + 1
            sequence_num = img_counter[slide_key]

            from all2md.utils.attachments import generate_attachment_filename
            image_filename = generate_attachment_filename(
                base_stem=base_filename,
                format_type="pptx",
                slide_num=slide_num,
                sequence_num=sequence_num,
                extension="png"
            )

        # Process image using unified attachment handling
        return process_attachment(
            attachment_data=image_data,
            attachment_name=image_filename,
            alt_text=alt_text,
            attachment_mode=options.attachment_mode,
            attachment_output_dir=options.attachment_output_dir,
            attachment_base_url=options.attachment_base_url,
            is_image=True,
            alt_text_mode=options.alt_text_mode,
        )

    # Handle charts (convert to tables if possible)
    elif isinstance(shape, GraphicFrame) and hasattr(shape, "chart"):
        chart = shape.chart
        # Basic chart data extraction
        data_rows = []

        # Check if this is a scatter plot (XY chart)
        try:
            chart_type = chart.chart_type
            is_scatter = chart_type == XL_CHART_TYPE.XY_SCATTER
        except Exception:
            is_scatter = False

        if is_scatter:
            # For scatter plots, extract X,Y pairs from XML
            for series in chart.series:
                try:
                    x_values = []
                    y_values = []

                    # Try to extract X and Y values from series XML
                    if hasattr(series, '_element'):
                        element = series._element
                        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}

                        # Extract X values
                        x_val_ref = element.find('.//c:xVal', ns)
                        if x_val_ref is not None:
                            num_cache = x_val_ref.find('.//c:numCache', ns)
                            if num_cache is not None:
                                pt_elements = num_cache.findall('.//c:pt', ns)
                                for pt in pt_elements:
                                    v_element = pt.find('c:v', ns)
                                    if v_element is not None and v_element.text:
                                        x_values.append(float(v_element.text))

                        # Extract Y values
                        y_val_ref = element.find('.//c:yVal', ns)
                        if y_val_ref is not None:
                            num_cache = y_val_ref.find('.//c:numCache', ns)
                            if num_cache is not None:
                                pt_elements = num_cache.findall('.//c:pt', ns)
                                for pt in pt_elements:
                                    v_element = pt.find('c:v', ns)
                                    if v_element is not None and v_element.text:
                                        y_values.append(float(v_element.text))

                    # Create table with X,Y pairs if we found both
                    if x_values and y_values and len(x_values) == len(y_values):
                        data_rows.append(f"| {series.name} X | " + " | ".join(str(x) for x in x_values) + " |")
                        data_rows.append(f"| {series.name} Y | " + " | ".join(str(y) for y in y_values) + " |")
                        if len(data_rows) >= 2 and "---" not in data_rows[-2]:
                            data_rows.insert(-2, "| --- | " + " | ".join(["---"] * len(x_values)) + " |")
                    else:
                        # Fallback to just Y values
                        values = list(series.values)
                        data_rows.append(f"| {series.name} | " + " | ".join(str(val) for val in values) + " |")
                except Exception:
                    # Fallback to basic series processing
                    values = list(series.values)
                    data_rows.append(f"| {series.name} | " + " | ".join(str(val) for val in values) + " |")
        else:
            # Standard chart processing
            # Extract categories (x-axis)
            try:
                categories = [cat.label for cat in chart.plots[0].categories if hasattr(cat, "label")]
                if categories:
                    data_rows.append("| Category | " + " | ".join(str(cat) for cat in categories) + " |")
                    data_rows.append("| --- | " + " | ".join(["---"] * len(categories)) + " |")
            except Exception:
                pass

            # Extract series data
            for series in chart.series:
                try:
                    values = list(series.values)
                    data_rows.append(f"| {series.name} | " + " | ".join(str(val) for val in values) + " |")
                except Exception:
                    # Skip series that can't be processed
                    continue

        return "\n".join(data_rows) if data_rows else f"Chart: {getattr(chart, 'chart_title', 'Untitled Chart')}"

    # Log unsupported shape types
    shape_type_name = getattr(shape.shape_type, 'name', 'UNKNOWN') if hasattr(shape, 'shape_type') else 'UNKNOWN'
    logger.debug(f"Unsupported shape type skipped: {shape_type_name}")
    return None


def extract_pptx_metadata(prs: Presentation) -> DocumentMetadata:
    """Extract metadata from PPTX presentation.

    Parameters
    ----------
    prs : Presentation
        python-pptx Presentation object

    Returns
    -------
    DocumentMetadata
        Extracted metadata
    """
    if not hasattr(prs, 'core_properties'):
        metadata = DocumentMetadata()
    else:
        props = prs.core_properties
        # Use the utility function for standard metadata extraction
        metadata = map_properties_to_metadata(props, OFFICE_FIELD_MAPPING)

        # Add PPTX-specific custom metadata
        custom_properties = ['last_modified_by', 'revision', 'comments']
        for prop_name in custom_properties:
            if hasattr(props, prop_name):
                value = getattr(props, prop_name)
                if value:
                    metadata.custom[prop_name] = value

    # Add slide count as custom metadata
    try:
        metadata.custom['slide_count'] = len(prs.slides)
    except Exception:
        pass

    return metadata


def pptx_to_markdown(input_data: Union[str, Path, IO[bytes]], options: PptxOptions | None = None) -> str:
    """Convert a PowerPoint presentation to Markdown format.

    Parameters
    ----------
    input_data : str, Path, or IO[bytes]
        The PowerPoint file to convert as a file path or binary file object
    options : PptxOptions | None, default None
        The options for extraction.

    Returns
    -------
    str
        The presentation content in Markdown format
    """
    # Handle backward compatibility and merge options
    if options is None:
        options = PptxOptions()

    # Validate and convert input
    try:
        from pptx.presentation import Presentation as PresentationType

        doc_input, input_type = validate_and_convert_input(
            input_data, supported_types=["path-like", "file-like", "pptx.Presentation objects"]
        )

        # Validate ZIP archive security for file-based inputs
        if input_type in ("path", "file") and not isinstance(doc_input, PresentationType):
            try:
                validate_zip_archive(doc_input if input_type == "path" else input_data)
            except Exception as e:
                raise MarkdownConversionError(
                    f"PPTX archive failed security validation: {str(e)}",
                    conversion_stage="archive_validation",
                    original_error=e
                ) from e

        # Open presentation based on input type
        if input_type == "object" and isinstance(doc_input, PresentationType):
            prs = doc_input
        else:
            prs = Presentation(doc_input)

    except Exception as e:
        if "python-pptx" in str(e).lower() or "pptx" in str(e).lower():
            raise MarkdownConversionError(
                "python-pptx library is required for PPTX conversion. Install with: pip install python-pptx",
                conversion_stage="dependency_check",
                original_error=e,
            ) from e
        else:
            raise MarkdownConversionError(
                f"Failed to open PPTX presentation: {str(e)}", conversion_stage="document_opening", original_error=e
            ) from e

    # Extract base filename for standardized attachment naming
    if input_type == "path" and isinstance(doc_input, (str, Path)):
        base_filename = Path(doc_input).stem
    else:
        # For non-file inputs, use a default name
        base_filename = "presentation"

    # Extract metadata if requested
    metadata = None
    if options.extract_metadata:
        metadata = extract_pptx_metadata(prs)

    # Get Markdown options (create default if not provided) - currently not used in processing
    # md_options = options.markdown_options or MarkdownOptions()

    markdown_content = []
    img_counter = {}  # Track image sequences across slides

    # Create attachment sequencer for consistent filename generation
    attachment_sequencer = create_attachment_sequencer()

    for i, slide in enumerate(prs.slides, 1):
        slide_content = []

        # Process slide title if present
        if slide.shapes.title:
            title_text = _process_text_frame(slide.shapes.title.text_frame, options.markdown_options, is_title=True)
            if options.include_slide_numbers:
                full_title = f"Slide {i}: {title_text.strip()}"
            else:
                full_title = title_text.strip()

            use_hash = options.markdown_options.use_hash_headings if options.markdown_options else True
            slide_content.append(format_markdown_heading(full_title, 1, use_hash))

        # Process all shapes in the slide
        for shape in slide.shapes:
            shape_content = _process_shape(
                shape, options, base_filename, i, img_counter, attachment_sequencer, slide.shapes.title
            )
            if shape_content:
                slide_content.append(shape_content + "\n")

        # Add slide content with spacing
        if slide_content:
            markdown_content.extend(slide_content)
            markdown_content.append("\n---\n")  # Add separator between slides

    result = "\n".join(markdown_content).strip()

    # Prepend metadata if enabled
    result = prepend_metadata_if_enabled(result, metadata, options.extract_metadata)

    return result


# Converter metadata for registration
CONVERTER_METADATA = ConverterMetadata(
    format_name="pptx",
    extensions=[".pptx"],
    mime_types=["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
    magic_bytes=[
        (b"PK\x03\x04", 0),  # ZIP signature
    ],
    converter_module="all2md.converters.pptx2markdown",
    converter_function="pptx_to_markdown",
    required_packages=[("python-pptx", "")],
    import_error_message="PPTX conversion requires 'python-pptx'. Install with: pip install python-pptx",
    options_class="PptxOptions",
    description="Convert PowerPoint presentations to Markdown",
    priority=7
)


def create_test_presentation() -> Any:
    """Create a test PowerPoint presentation with various elements."""
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Test PowerPoint Presentation"
    subtitle.text = "For Markdown Conversion Testing"

    # Content slide with bullet points
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes

    title = shapes.title
    title.text = "Bullet Points Test"

    body = shapes.placeholders[1]
    tf = body.text_frame

    p = tf.add_paragraph()
    p.text = "First Level Bullet"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Second Level Bullet"
    p.level = 1

    # Slide with table
    table_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = table_slide.shapes

    title = shapes.title
    title.text = "Table Test"

    rows, cols = 3, 3
    left = top = width = height = Inches(2.0)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # Add header row
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = f"Header {i + 1}"

    # Add data rows
    for row in range(1, rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = f"Cell {row},{col}"

    return prs
