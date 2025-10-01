#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# src/all2md/converters/pptx2ast.py
"""PPTX to AST converter.

This module provides conversion from Microsoft PowerPoint presentations to AST representation.
It replaces direct markdown string generation with structured AST building,
enabling multiple rendering strategies and improved testability.

"""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.text.text import TextFrame

from all2md.ast import (
    Document,
    Emphasis,
    Heading,
    Image,
    List,
    ListItem,
    Node,
    Paragraph as AstParagraph,
    Strong,
    Table as AstTable,
    TableCell,
    TableRow,
    Text,
    ThematicBreak,
    Underline,
)
from all2md.options import PptxOptions

logger = logging.getLogger(__name__)


def parse_slide_ranges(slide_spec: str, total_slides: int) -> list[int]:
    """Parse slide range specification into list of 0-based slide indices.

    Supports various formats:
    - "1-3" → [0, 1, 2]
    - "5" → [4]
    - "10-" → [9, 10, ..., total_slides-1]
    - "1-3,5,10-" → combined ranges

    Parameters
    ----------
    slide_spec : str
        Slide range specification
    total_slides : int
        Total number of slides in presentation

    Returns
    -------
    list of int
        0-based slide indices

    Examples
    --------
    >>> parse_slide_ranges("1-3,5", 10)
    [0, 1, 2, 4]
    >>> parse_slide_ranges("8-", 10)
    [7, 8, 9]

    """
    slides = set()

    # Split by comma to handle multiple ranges
    parts = slide_spec.split(',')

    for part in parts:
        part = part.strip()
        if not part:
            continue

        # Handle range (e.g., "1-3" or "10-")
        if '-' in part:
            range_parts = part.split('-', 1)
            start_str = range_parts[0].strip()
            end_str = range_parts[1].strip()

            # Parse start (1-based to 0-based)
            if start_str:
                start = int(start_str) - 1
            else:
                start = 0

            # Parse end (1-based to 0-based, or use total_slides if empty)
            if end_str:
                end = int(end_str) - 1
            else:
                end = total_slides - 1

            # Add all slides in range
            for s in range(start, end + 1):
                if 0 <= s < total_slides:
                    slides.add(s)
        else:
            # Single slide (1-based to 0-based)
            slide = int(part) - 1
            if 0 <= slide < total_slides:
                slides.add(slide)

    # Return sorted list
    return sorted(slides)


class PptxToAstConverter:
    """Convert PPTX to AST representation.

    This converter parses PowerPoint presentations using python-pptx and builds an AST
    that can be rendered to various markdown flavors.

    Parameters
    ----------
    options : PptxOptions or None, default = None
        Conversion options
    base_filename : str
        Base filename for image attachments
    attachment_sequencer : callable or None
        Sequencer for generating attachment filenames

    """

    def __init__(
        self,
        options: PptxOptions | None = None,
        base_filename: str = "presentation",
        attachment_sequencer: Any = None,
    ):
        self.options = options or PptxOptions()
        self.base_filename = base_filename
        self.attachment_sequencer = attachment_sequencer
        self._current_slide_num = 0

    def convert_to_ast(self, prs: "Presentation") -> Document:
        """Convert PPTX presentation to AST Document.

        Parameters
        ----------
        prs : Presentation
            PPTX presentation to convert

        Returns
        -------
        Document
            AST document node

        """
        children: list[Node] = []

        # Determine which slides to process
        all_slides = list(prs.slides)
        total_slides = len(all_slides)

        if self.options.slides:
            # Parse slide range specification
            slide_indices = parse_slide_ranges(self.options.slides, total_slides)
        else:
            # Process all slides
            slide_indices = list(range(total_slides))

        # Process selected slides
        for idx in slide_indices:
            slide = all_slides[idx]
            self._current_slide_num = idx + 1  # 1-based for display
            slide_nodes = self._process_slide_to_ast(slide)
            if slide_nodes:
                children.extend(slide_nodes)

                # Add slide separator (thematic break) after every slide
                children.append(ThematicBreak())

        return Document(children=children)

    def _process_slide_to_ast(self, slide: Any) -> list[Node]:
        """Process a PPTX slide to AST nodes.

        Parameters
        ----------
        slide : Slide
            PPTX slide to process

        Returns
        -------
        list of Node
            List of AST nodes representing the slide

        """
        nodes: list[Node] = []

        # Process slide title if present and enabled
        if self.options.include_titles_as_h2 and slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            if self.options.include_slide_numbers:
                title_text = f"Slide {self._current_slide_num}: {title_text}"

            # Slide titles are level 2 headings
            nodes.append(Heading(level=2, content=[Text(content=title_text)]))

        # Process all shapes in the slide
        for shape in slide.shapes:
            # Skip the title shape if it was already processed as a heading
            if self.options.include_titles_as_h2 and shape == slide.shapes.title:
                continue

            shape_nodes = self._process_shape_to_ast(shape)
            if shape_nodes:
                if isinstance(shape_nodes, list):
                    nodes.extend(shape_nodes)
                else:
                    nodes.append(shape_nodes)

        return nodes

    def _process_shape_to_ast(self, shape: "BaseShape") -> Node | list[Node] | None:
        """Process a PPTX shape to AST nodes.

        Parameters
        ----------
        shape : BaseShape
            PPTX shape to process

        Returns
        -------
        Node, list of Node, or None
            Resulting AST node(s)

        """
        # Check if shape has text_frame
        if hasattr(shape, "text_frame") and shape.text_frame:
            return self._process_text_frame_to_ast(shape.text_frame)

        # Check if shape is a table
        if hasattr(shape, "table"):
            return self._process_table_to_ast(shape.table)

        # Check if shape is an image
        if hasattr(shape, "image"):
            return self._process_image_to_ast(shape)

        # For other shapes, skip or handle as needed
        return None

    def _process_text_frame_to_ast(self, frame: "TextFrame") -> list[Node] | None:
        """Process a text frame to AST nodes.

        Parameters
        ----------
        frame : TextFrame
            Text frame to process

        Returns
        -------
        list of Node or None
            List of AST nodes (paragraphs, lists, etc.)

        """
        # Import here to avoid circular dependencies
        from all2md.converters.pptx2markdown import _analyze_slide_context, _detect_list_item

        nodes: list[Node] = []
        slide_context = _analyze_slide_context(frame)

        # Track list accumulation
        current_list: List | None = None
        current_list_type: str | None = None

        for paragraph in frame.paragraphs:
            if not paragraph.text.strip():
                continue

            # Detect if this is a list item
            is_list_item, list_type = _detect_list_item(paragraph, slide_context)

            if is_list_item:
                # Check if we need to start a new list or continue current
                if current_list is None or current_list_type != list_type:
                    # Finalize previous list
                    if current_list:
                        nodes.append(current_list)

                    # Start new list
                    current_list_type = list_type
                    current_list = List(ordered=(list_type == "number"), items=[], tight=True)

                # Add item to current list
                item_content = self._process_paragraph_runs_to_inline(paragraph)
                if item_content:
                    list_item = ListItem(children=[AstParagraph(content=item_content)])
                    current_list.items.append(list_item)
            else:
                # Not a list item - finalize any current list
                if current_list:
                    nodes.append(current_list)
                    current_list = None
                    current_list_type = None

                # Process as regular paragraph
                para_content = self._process_paragraph_runs_to_inline(paragraph)
                if para_content:
                    nodes.append(AstParagraph(content=para_content))

        # Finalize any remaining list
        if current_list:
            nodes.append(current_list)

        return nodes if nodes else None

    def _process_paragraph_runs_to_inline(self, paragraph: Any) -> list[Node]:
        """Process paragraph runs to inline AST nodes.

        Parameters
        ----------
        paragraph : Paragraph
            Paragraph containing runs

        Returns
        -------
        list of Node
            List of inline AST nodes

        """
        result: list[Node] = []

        # Group runs by formatting to avoid excessive nesting
        grouped_runs: list[tuple[str, tuple[bool, bool, bool] | None]] = []
        current_text: list[str] = []
        current_format: tuple[bool, bool, bool] | None = None

        for run in paragraph.runs:
            text = run.text
            if not text.strip():  # Skip empty runs
                continue

            # Get formatting key (bold, italic, underline)
            format_key = (
                run.font.bold or False,
                run.font.italic or False,
                run.font.underline or False,
            )

            # Start new group if format changes
            if format_key != current_format:
                if current_text:
                    # Join with single space and strip
                    grouped_runs.append((" ".join(current_text).strip(), current_format))
                    current_text = []
                current_format = format_key

            current_text.append(text.strip())

        # Add final group
        if current_text:
            grouped_runs.append((" ".join(current_text).strip(), current_format))

        # Convert groups to AST nodes
        for text, format_key in grouped_runs:
            if not text:
                continue

            # Build inline nodes with formatting
            inline_node: Node = Text(content=text)

            if format_key:
                # Apply formatting layers from innermost to outermost
                if format_key[2]:  # underline
                    inline_node = Underline(content=[inline_node])
                if format_key[0]:  # bold
                    inline_node = Strong(content=[inline_node])
                if format_key[1]:  # italic
                    inline_node = Emphasis(content=[inline_node])

            result.append(inline_node)

        return result

    def _process_table_to_ast(self, table: Any) -> AstTable | None:
        """Process a PPTX table to AST Table node.

        Parameters
        ----------
        table : Table
            PPTX table to convert

        Returns
        -------
        AstTable or None
            Table node if table has content

        """
        if len(table.rows) == 0:
            return None

        # First row is header
        header_cells: list[TableCell] = []
        for cell in table.rows[0].cells:
            if cell.text_frame:
                cell_content = self._process_paragraph_runs_to_inline_simple(cell.text_frame)
                header_cells.append(TableCell(content=cell_content))
            else:
                header_cells.append(TableCell(content=[]))

        header_row = TableRow(cells=header_cells, is_header=True)

        # Data rows (skip first row which is header)
        data_rows: list[TableRow] = []
        for i, row in enumerate(table.rows):
            if i == 0:
                continue  # Skip header row
            row_cells: list[TableCell] = []
            for cell in row.cells:
                if cell.text_frame:
                    cell_content = self._process_paragraph_runs_to_inline_simple(cell.text_frame)
                    row_cells.append(TableCell(content=cell_content))
                else:
                    row_cells.append(TableCell(content=[]))
            data_rows.append(TableRow(cells=row_cells))

        return AstTable(header=header_row, rows=data_rows)

    def _process_paragraph_runs_to_inline_simple(self, frame: "TextFrame") -> list[Node]:
        """Process all paragraphs in a text frame to inline nodes (for tables).

        Parameters
        ----------
        frame : TextFrame
            Text frame to process

        Returns
        -------
        list of Node
            List of inline nodes

        """
        result: list[Node] = []

        for paragraph in frame.paragraphs:
            para_inlines = self._process_paragraph_runs_to_inline(paragraph)
            result.extend(para_inlines)

            # Add space between paragraphs (except last)
            if paragraph != frame.paragraphs[-1]:
                result.append(Text(content=" "))

        return result

    def _process_image_to_ast(self, shape: Any) -> Image | None:
        """Process an image shape to AST Image node.

        Parameters
        ----------
        shape : Shape
            Image shape to process

        Returns
        -------
        Image or None
            Image node if image can be processed

        """
        # Import here to avoid circular dependencies
        from all2md.converters.pptx2markdown import extract_pptx_image_data
        from all2md.utils.attachments import process_attachment

        try:
            # Extract image data
            image_data, extension = extract_pptx_image_data(shape)

            if not image_data:
                return None

            # Use sequencer if available
            if self.attachment_sequencer:
                image_filename, _ = self.attachment_sequencer(
                    base_stem=self.base_filename,
                    format_type="general",
                    extension=extension or "png"
                )
            else:
                from all2md.utils.attachments import generate_attachment_filename
                image_filename = generate_attachment_filename(
                    base_stem=self.base_filename,
                    format_type="general",
                    sequence_num=self._current_slide_num,
                    extension=extension or "png"
                )

            # Process attachment
            processed_image = process_attachment(
                attachment_data=image_data,
                attachment_name=image_filename,
                alt_text="image",
                attachment_mode=self.options.attachment_mode,
                attachment_output_dir=self.options.attachment_output_dir,
                attachment_base_url=self.options.attachment_base_url,
                is_image=True,
                alt_text_mode=self.options.alt_text_mode,
            )

            # Parse the markdown image string
            import re
            match = re.match(r'^!\[([^\]]*)\](?:\(([^)]+)\))?$', processed_image)
            if match:
                alt_text = match.group(1)
                url = match.group(2) or ""
                return Image(url=url, alt_text=alt_text, title=None)

        except Exception as e:
            logger.debug(f"Failed to process image: {e}")
            return None

        return None
