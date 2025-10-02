#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# src/all2md/parsers/pptx2markdown.py

"""PowerPoint presentation to Markdown conversion module.

This module provides functionality to extract content from Microsoft PowerPoint
presentations (PPTX format) and convert it to structured Markdown. It processes
slides individually, extracting text content, formatting, tables, and embedded
images while preserving the logical structure and hierarchy of the presentation.

The converter analyzes each slide's layout and content, handling various slide
elements including title text, body content, tables, charts, and embedded media.
Special attention is paid to maintaining the presentation's logical flow and
structure in the resulting Markdown format.

Key Features
------------
- Slide-by-slide content extraction and conversion
- Text formatting preservation (bold, italic, lists)
- Table extraction with proper Markdown table formatting
- Image embedding as base64 data URIs
- Chart and graphic frame handling
- Hierarchical content structure preservation
- Bullet point and numbering conversion
- Slide title and content separation

Supported Elements
------------------
- Text content with formatting (bold, italic, underline)
- Slide titles and subtitles
- Bulleted and numbered lists with indentation
- Tables with cell content and structure
- Images and graphics (embedded as base64)
- Charts and diagrams (metadata extraction)
- Text boxes and content placeholders

Processing Features
-------------------
- Automatic slide numbering and separation
- Paragraph-level formatting detection
- List indentation and nesting preservation
- Image format detection and encoding
- Error handling for corrupted or complex elements
- Logging for debugging and issue tracking

Dependencies
------------
- python-pptx: For PowerPoint file parsing and content extraction
- base64: For image encoding
- logging: For debug and error reporting

Examples
--------
Basic presentation conversion:

    >>> from all2md.parsers.pptx2markdown import pptx_to_markdown
    >>> with open('presentation.pptx', 'rb') as f:
    ...     markdown = pptx_to_markdown(f)
    >>> print(markdown)

Note
----
Requires python-pptx package. Complex animations, transitions, and some
advanced PowerPoint features will be omitted in the conversion process.
The focus is on extracting textual and structural content.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import IO, TYPE_CHECKING, Any, Union

from all2md.converter_metadata import ConverterMetadata
from all2md.exceptions import MarkdownConversionError
from all2md.options import PptxOptions
from all2md.utils.attachments import create_attachment_sequencer
from all2md.utils.inputs import validate_and_convert_input
from all2md.utils.metadata import (
    OFFICE_FIELD_MAPPING,
    DocumentMetadata,
    map_properties_to_metadata,
    prepend_metadata_if_enabled,
)
from all2md.utils.security import validate_zip_archive

# Type checking imports for static analysis without runtime overhead
if TYPE_CHECKING:
    from pptx import Presentation

logger = logging.getLogger(__name__)


# To be removed shortly
def extract_pptx_metadata(prs: Presentation) -> DocumentMetadata:
    """Extract metadata from PPTX presentation.

    Parameters
    ----------
    prs : Presentation
        python-pptx Presentation object

    Returns
    -------
    DocumentMetadata
        Extracted metadata
    """
    if not hasattr(prs, 'core_properties'):
        metadata = DocumentMetadata()
    else:
        props = prs.core_properties
        # Use the utility function for standard metadata extraction
        metadata = map_properties_to_metadata(props, OFFICE_FIELD_MAPPING)

        # Add PPTX-specific custom metadata
        custom_properties = ['last_modified_by', 'revision', 'comments']
        for prop_name in custom_properties:
            if hasattr(props, prop_name):
                value = getattr(props, prop_name)
                if value:
                    metadata.custom[prop_name] = value

    # Add slide count as custom metadata
    try:
        metadata.custom['slide_count'] = len(prs.slides)
    except Exception:
        pass

    return metadata

# To be removed shortly.
def pptx_to_markdown(input_data: Union[str, Path, IO[bytes]], options: PptxOptions | None = None) -> str:
    """Convert a PowerPoint presentation to Markdown format.

    Parameters
    ----------
    input_data : str, Path, or IO[bytes]
        The PowerPoint file to convert as a file path or binary file object
    options : PptxOptions | None, default None
        The options for extraction.

    Returns
    -------
    str
        The presentation content in Markdown format
    """
    # Lazy import of heavy python-pptx dependencies
    from pptx import Presentation

    # Handle backward compatibility and merge options
    if options is None:
        options = PptxOptions()

    # Validate and convert input
    try:
        from pptx.presentation import Presentation as PresentationType

        doc_input, input_type = validate_and_convert_input(
            input_data, supported_types=["path-like", "file-like", "pptx.Presentation objects"]
        )

        # Validate ZIP archive security for file-based inputs
        if input_type in ("path", "file") and not isinstance(doc_input, PresentationType):
            try:
                validate_zip_archive(doc_input if input_type == "path" else input_data)
            except Exception as e:
                raise MarkdownConversionError(
                    f"PPTX archive failed security validation: {str(e)}",
                    conversion_stage="archive_validation",
                    original_error=e
                ) from e

        # Open presentation based on input type
        if input_type == "object" and isinstance(doc_input, PresentationType):
            prs = doc_input
        else:
            prs = Presentation(doc_input)

    except Exception as e:
        if "python-pptx" in str(e).lower() or "pptx" in str(e).lower():
            raise MarkdownConversionError(
                "python-pptx library is required for PPTX conversion. Install with: pip install python-pptx",
                conversion_stage="dependency_check",
                original_error=e,
            ) from e
        else:
            raise MarkdownConversionError(
                f"Failed to open PPTX presentation: {str(e)}", conversion_stage="document_opening", original_error=e
            ) from e

    # Extract base filename for standardized attachment naming
    if input_type == "path" and isinstance(doc_input, (str, Path)):
        base_filename = Path(doc_input).stem
    else:
        # For non-file inputs, use a default name
        base_filename = "presentation"

    # Extract metadata if requested
    metadata = None
    if options.extract_metadata:
        metadata = extract_pptx_metadata(prs)

    # Use new AST-based conversion path
    from all2md.parsers.pptx import PptxToAstConverter
    from all2md.ast import MarkdownRenderer
    from all2md.options import MarkdownOptions

    # Create attachment sequencer for consistent filename generation
    attachment_sequencer = create_attachment_sequencer()

    # Convert PPTX to AST
    ast_converter = PptxToAstConverter(
        options=options,
        base_filename=base_filename,
        attachment_sequencer=attachment_sequencer,
    )
    ast_document = ast_converter.convert_to_ast(prs)

    # Get MarkdownOptions (use provided or create default)
    md_opts = options.markdown_options if options.markdown_options else MarkdownOptions()

    # Render AST to markdown using MarkdownOptions directly
    renderer = MarkdownRenderer(md_opts)
    result = renderer.render(ast_document)

    # Prepend metadata if enabled
    result = prepend_metadata_if_enabled(result, metadata, options.extract_metadata)

    return result


# Converter metadata for registration
CONVERTER_METADATA = ConverterMetadata(
    format_name="pptx",
    extensions=[".pptx"],
    mime_types=["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
    magic_bytes=[
        (b"PK\x03\x04", 0),  # ZIP signature
    ],
    converter_module="all2md.parsers.pptx2markdown",
    converter_function="pptx_to_markdown",
    required_packages=[("python-pptx", "pptx", "")],
    import_error_message="PPTX conversion requires 'python-pptx'. Install with: pip install python-pptx",
    options_class="PptxOptions",
    description="Convert PowerPoint presentations to Markdown",
    priority=7
)


def create_test_presentation() -> Any:
    """Create a test PowerPoint presentation with various elements."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Test PowerPoint Presentation"
    subtitle.text = "For Markdown Conversion Testing"

    # Content slide with bullet points
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes

    title = shapes.title
    title.text = "Bullet Points Test"

    body = shapes.placeholders[1]
    tf = body.text_frame

    p = tf.add_paragraph()
    p.text = "First Level Bullet"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Second Level Bullet"
    p.level = 1

    # Slide with table
    table_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = table_slide.shapes

    title = shapes.title
    title.text = "Table Test"

    rows, cols = 3, 3
    left = top = width = height = Inches(2.0)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # Add header row
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = f"Header {i + 1}"

    # Add data rows
    for row in range(1, rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = f"Cell {row},{col}"

    return prs
