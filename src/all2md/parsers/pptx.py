#  Copyright (c) 2025 Tom Villani, Ph.D.
#
# src/all2md/parsers/pptx.py
"""PPTX to AST converter.

This module provides conversion from Microsoft PowerPoint presentations to AST representation.
It replaces direct markdown string generation with structured AST building,
enabling multiple rendering strategies and improved testability.

"""

from __future__ import annotations

import logging
import re
import tempfile
import zipfile
from pathlib import Path
from typing import IO, TYPE_CHECKING, Any, Optional, Union, cast

import defusedxml.ElementTree as ET

from all2md.exceptions import MalformedFileError, ZipFileSecurityError
from all2md.options.pptx import PptxOptions
from all2md.progress import ProgressCallback
from all2md.utils.inputs import parse_page_ranges, validate_and_convert_input

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.text.text import TextFrame

from all2md.ast import (
    CodeBlock,
    Comment,
    Document,
    Emphasis,
    Heading,
    Image,
    Link,
    List,
    ListItem,
    Node,
    Strikethrough,
    Strong,
    Subscript,
    Superscript,
    TableCell,
    TableRow,
    Text,
    ThematicBreak,
    Underline,
)
from all2md.ast import (
    Paragraph as AstParagraph,
)
from all2md.ast import (
    Table as AstTable,
)
from all2md.constants import DEPS_PPTX
from all2md.converter_metadata import ConverterMetadata
from all2md.parsers.base import BaseParser
from all2md.utils.attachments import (
    create_attachment_sequencer,
    extract_pptx_image_data,
    process_attachment,
)
from all2md.utils.chart_helpers import build_chart_table
from all2md.utils.decorators import requires_dependencies
from all2md.utils.metadata import OFFICE_FIELD_MAPPING, DocumentMetadata, map_properties_to_metadata
from all2md.utils.parser_helpers import attachment_result_to_image_node, group_and_format_runs

logger = logging.getLogger(__name__)


class PptxToAstConverter(BaseParser):
    """Convert PPTX to AST representation.

    This converter parses PowerPoint presentations using python-pptx and builds an AST
    that can be rendered to various markdown flavors.

    Parameters
    ----------
    options : PptxOptions or None, default = None
        Conversion options

    """

    def __init__(self, options: PptxOptions | None = None, progress_callback: Optional[ProgressCallback] = None):
        """Initialize the PPTX parser with options and progress callback."""
        BaseParser._validate_options_type(options, PptxOptions, "pptx")
        options = options or PptxOptions()
        super().__init__(options, progress_callback)
        self.options: PptxOptions = options

        self._current_slide_num = 0
        self._base_filename = "presentation"
        self._attachment_sequencer = create_attachment_sequencer()
        self._attachment_footnotes: dict[str, str] = {}  # label -> content for footnote definitions
        self._pptx_comments: dict[int, list[Comment]] = {}  # slide_num -> list of Comment nodes
        self._input_data: Union[str, Path, IO[bytes], bytes, None] = None  # Store input for comment extraction

    @requires_dependencies("pptx", DEPS_PPTX)
    def parse(self, input_data: Union[str, Path, IO[bytes], bytes]) -> Document:
        """Parse PPTX input into an AST Document.

        Parameters
        ----------
        input_data : str, Path, IO[bytes], or bytes
            The PPTX document to parse. Can be:
            - File path (str or Path)
            - File-like object in binary mode
            - Raw PPTX bytes

        Returns
        -------
        Document
            AST Document node representing the parsed PPTX structure

        Raises
        ------
        ParsingError
            If parsing fails due to invalid format or corruption
        DependencyError
            If python-pptx is not installed

        """
        from pptx import Presentation
        from pptx.presentation import Presentation as PresentationType

        # Validate and convert input
        try:
            doc_input, input_type = validate_and_convert_input(
                input_data, supported_types=["path-like", "file-like", "pptx.Presentation objects"]
            )

            # Validate ZIP archive security for all input types
            if not isinstance(doc_input, PresentationType):
                self._validate_zip_security(input_data, suffix=".pptx")

            # Open presentation based on input type
            if input_type == "object" and isinstance(doc_input, PresentationType):
                prs = doc_input
            else:
                prs = Presentation(doc_input)

        except ZipFileSecurityError:
            raise

        except Exception as e:
            raise MalformedFileError(f"Failed to open PPTX presentation: {e!r}", original_error=e) from e

        # Extract base filename for standardized attachment naming
        if input_type == "path" and isinstance(doc_input, (str, Path)):
            self._base_filename = Path(doc_input).stem
        else:
            self._base_filename = "presentation"

        # Store input data for comment extraction
        self._input_data = input_data

        # Convert PPTX to AST
        return self.convert_to_ast(prs)

    def convert_to_ast(self, prs: "Presentation") -> Document:
        """Convert PPTX presentation to AST Document.

        Parameters
        ----------
        prs : Presentation
            PPTX presentation to convert

        Returns
        -------
        Document
            AST document node

        """
        # Reset parser state to prevent leakage across parse calls
        self._attachment_footnotes = {}
        self._current_slide_num = 0
        self._attachment_sequencer = create_attachment_sequencer()
        self._pptx_comments = {}

        # Extract PPTX comments from XML if input data is available
        if self._input_data is not None:
            try:
                self._pptx_comments = self._extract_pptx_comments(self._input_data)
                if self._pptx_comments:
                    logger.debug(f"Extracted comments from {len(self._pptx_comments)} slides")
            except Exception as e:
                logger.debug(f"Failed to extract PPTX comments: {e}")

        children: list[Node] = []

        # Determine which slides to process
        all_slides = list(prs.slides)
        total_slides = len(all_slides)

        # Emit started event
        self._emit_progress(
            "started",
            f"Converting PPTX with {total_slides} slide{'s' if total_slides != 1 else ''}",
            current=0,
            total=total_slides,
        )

        if self.options.slides:
            # Parse slide range specification
            slide_indices = parse_page_ranges(self.options.slides, total_slides)
        else:
            # Process all slides
            slide_indices = list(range(total_slides))

        # Process selected slides
        for idx in slide_indices:
            slide = all_slides[idx]
            self._current_slide_num = idx + 1  # 1-based for display
            slide_nodes = self._process_slide_to_ast(slide)
            if slide_nodes:
                children.extend(slide_nodes)

                # Add slide separator (thematic break) after every slide
                children.append(ThematicBreak())

        # Extract and attach metadata
        metadata = self.extract_metadata(prs)

        # Append footnote definitions if any were collected
        if self.options.attachments_footnotes_section:
            self._append_attachment_footnotes(
                children, self._attachment_footnotes, self.options.attachments_footnotes_section
            )

        # Emit finished event
        self._emit_progress(
            "finished",
            f"PPTX conversion completed ({len(slide_indices)} slide{'s' if len(slide_indices) != 1 else ''})",
            current=len(slide_indices),
            total=len(slide_indices),
        )

        return Document(children=children, metadata=metadata.to_dict())

    def _process_slide_to_ast(self, slide: Any) -> list[Node]:
        """Process a PPTX slide to AST nodes.

        Parameters
        ----------
        slide : Slide
            PPTX slide to process

        Returns
        -------
        list of Node
            List of AST nodes representing the slide

        """
        nodes: list[Node] = []

        # Process slide title if present and enabled
        if self.options.include_titles_as_h2 and slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            if self.options.include_slide_numbers:
                title_text = f"Slide {self._current_slide_num}: {title_text}"

            # Slide titles are level 2 headings
            nodes.append(Heading(level=2, content=[Text(content=title_text)]))

        # Process all shapes in the slide
        for shape in slide.shapes:
            # Skip the title shape if it was already processed as a heading
            if self.options.include_titles_as_h2 and shape == slide.shapes.title:
                continue

            shape_nodes = self._process_shape_to_ast(shape)
            if shape_nodes:
                if isinstance(shape_nodes, list):
                    nodes.extend(shape_nodes)
                else:
                    nodes.append(shape_nodes)

        # Process speaker notes if requested
        if self.options.include_notes:
            notes_nodes = self._extract_slide_notes(slide)
            if notes_nodes:
                nodes.extend(notes_nodes)

        # Add PPTX comments for this slide if available
        if self._current_slide_num in self._pptx_comments:
            slide_comments = self._pptx_comments[self._current_slide_num]
            nodes.extend(slide_comments)

        return nodes

    def _process_shape_to_ast(self, shape: "BaseShape") -> Node | list[Node] | None:
        """Process a PPTX shape to AST nodes.

        Parameters
        ----------
        shape : BaseShape
            PPTX shape to process

        Returns
        -------
        Node, list of Node, or None
            Resulting AST node(s)

        """
        # Check if shape has text_frame
        if hasattr(shape, "text_frame") and shape.text_frame:
            return self._process_text_frame_to_ast(shape.text_frame)

        # Check if shape is a table
        if hasattr(shape, "has_table") and shape.has_table:
            return self._process_table_to_ast(shape.table)

        # Check if shape is an image
        if hasattr(shape, "image"):
            return self._process_image_to_ast(shape)

        from pptx.shapes.graphfrm import GraphicFrame

        if isinstance(shape, GraphicFrame) and hasattr(shape, "has_chart") and shape.has_chart:
            return self._process_chart_to_ast(shape.chart)

        # For other shapes, skip or handle as needed
        return None

    def _process_chart_to_ast(self, chart: Any) -> Node | list[Node] | None:
        """Process a PPTX chart to AST Table node.

        Parameters
        ----------
        chart : Chart
            PPTX chart to convert

        Returns
        -------
        AstTable or None
            Table node representing chart data

        """
        from pptx.enum.chart import XL_CHART_TYPE

        try:
            chart_type = chart.chart_type
        except Exception:
            chart_type = None

        is_scatter = chart_type == XL_CHART_TYPE.XY_SCATTER if chart_type is not None else False
        mode = self.options.charts_mode.lower()

        table_node: AstTable | None = None
        mermaid_code: str | None = None

        if is_scatter:
            scatter_data = self._extract_scatter_chart_data(chart)
            if mode in {"data", "both"}:
                table_node = self._scatter_data_to_table(scatter_data)
            if mode in {"mermaid", "both"}:
                mermaid_code = self._scatter_chart_to_mermaid(chart, scatter_data)
            if mode == "data" and table_node is None:
                table_node = self._scatter_data_to_table(scatter_data)
        else:
            categories, series_rows = self._extract_standard_chart_data(chart)
            if mode in {"data", "both"}:
                table_node = self._standard_data_to_table(categories, series_rows)
            if mode in {"mermaid", "both"}:
                mermaid_code = self._standard_chart_to_mermaid(chart, chart_type, categories, series_rows)
            if mode == "data" and table_node is None:
                table_node = self._standard_data_to_table(categories, series_rows)

        mermaid_node: CodeBlock | None = None
        if mermaid_code:
            mermaid_node = CodeBlock(content=mermaid_code, language="mermaid")

        output_nodes: list[Node] = []

        if table_node and mode in {"data", "both"}:
            output_nodes.append(table_node)
        if mermaid_node and mode in {"mermaid", "both"}:
            output_nodes.append(mermaid_node)

        if not output_nodes:
            # Fallback to any available representation
            if mermaid_node:
                output_nodes.append(mermaid_node)
            elif table_node:
                output_nodes.append(table_node)
            else:
                return None

        if len(output_nodes) == 1:
            return output_nodes[0]
        return output_nodes

    def _extract_scatter_chart_data(self, chart: Any) -> list[tuple[str, list[float], list[float]]]:
        data: list[tuple[str, list[float], list[float]]] = []

        for series in chart.series:
            try:
                x_values: list[float] = []
                y_values: list[float] = []

                if hasattr(series, "_element"):
                    element = series._element
                    ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}

                    x_val_ref = element.find(".//c:xVal", ns)
                    if x_val_ref is not None:
                        num_cache = x_val_ref.find(".//c:numCache", ns)
                        if num_cache is not None:
                            for pt in num_cache.findall(".//c:pt", ns):
                                v_element = pt.find("c:v", ns)
                                if v_element is not None and v_element.text:
                                    try:
                                        x_values.append(float(v_element.text))
                                    except ValueError:
                                        continue

                    y_val_ref = element.find(".//c:yVal", ns)
                    if y_val_ref is not None:
                        num_cache = y_val_ref.find(".//c:numCache", ns)
                        if num_cache is not None:
                            for pt in num_cache.findall(".//c:pt", ns):
                                v_element = pt.find("c:v", ns)
                                if v_element is not None and v_element.text:
                                    try:
                                        y_values.append(float(v_element.text))
                                    except ValueError:
                                        continue

                if x_values and y_values and len(x_values) == len(y_values):
                    data.append((series.name or "Series", x_values, y_values))
                elif hasattr(series, "values") and series.values:
                    y_values = [float(v) for v in series.values if v is not None]
                    if y_values:
                        x_values = list(range(len(y_values)))
                        data.append((series.name or "Series", x_values, y_values))
            except Exception:
                continue

        return data

    def _process_scatter_chart_to_table(self, chart: Any) -> AstTable | None:
        """Process scatter plot to AST Table with X/Y rows."""
        series_data = self._extract_scatter_chart_data(chart)
        if not series_data:
            return None

        return self._scatter_data_to_table(series_data)

    def _extract_standard_chart_data(self, chart: Any) -> tuple[list[str], list[tuple[str, list[Any]]]]:
        categories: list[str] = []
        try:
            if hasattr(chart, "plots") and chart.plots:
                categories = [
                    cat.label if hasattr(cat, "label") else str(cat)
                    for cat in chart.plots[0].categories
                    if hasattr(cat, "label") or cat
                ]
        except Exception:
            categories = []

        series_rows: list[tuple[str, list[Any]]] = []
        for series in chart.series:
            try:
                if hasattr(series, "values") and series.values:
                    values = list(series.values)
                    series_name = series.name or "Series"
                    series_rows.append((series_name, values))
            except Exception:
                continue

        return categories, series_rows

    def _process_standard_chart_to_table(self, chart: Any) -> AstTable | None:
        """Process standard chart (bar, column, line, pie) to AST Table."""
        categories, series_rows = self._extract_standard_chart_data(chart)

        if not series_rows:
            return None

        return self._standard_data_to_table(categories, series_rows)

    def _scatter_data_to_table(
        self,
        series_data: list[tuple[str, list[float], list[float]]],
    ) -> AstTable | None:
        if not series_data:
            return None

        max_points = max((len(x_vals) for _, x_vals, _ in series_data), default=0)
        if max_points == 0:
            return None

        header_cells = [TableCell(content=[Text(content="Series")])]
        for i in range(max_points):
            header_cells.append(TableCell(content=[Text(content=f"Point {i + 1}")]))
        header_row = TableRow(cells=header_cells, is_header=True)

        all_rows: list[TableRow] = []
        for series_name, x_values, y_values in series_data:
            x_cells = [TableCell(content=[Text(content=f"{series_name} X")])]
            for val in x_values:
                x_cells.append(TableCell(content=[Text(content=self._format_scalar(val))]))
            all_rows.append(TableRow(cells=x_cells))

            y_cells = [TableCell(content=[Text(content=f"{series_name} Y")])]
            for val in y_values:
                y_cells.append(TableCell(content=[Text(content=self._format_scalar(val))]))
            all_rows.append(TableRow(cells=y_cells))

        return AstTable(header=header_row, rows=all_rows)

    def _standard_data_to_table(
        self,
        categories: list[str],
        series_rows: list[tuple[str, list[Any]]],
    ) -> AstTable | None:
        if not series_rows:
            return None

        # Use the build_chart_table helper to create consistent table structure
        # Note: build_chart_table expects categories as row labels and series as columns
        # For PPTX charts, we transpose the data to match this expected structure
        return build_chart_table(categories=categories, series_data=series_rows, category_header="Category")

    def _scatter_chart_to_mermaid(
        self,
        chart: Any,
        series_data: list[tuple[str, list[float], list[float]]],
    ) -> str | None:
        if not series_data:
            return None

        x_all: list[float] = []
        y_all: list[float] = []
        series_lines: list[str] = []

        for series_name, x_values, y_values in series_data:
            pairs: list[str] = []
            for x_val, y_val in zip(x_values, y_values, strict=False):
                x_num = self._coerce_float(x_val)
                y_num = self._coerce_float(y_val)
                if x_num is None or y_num is None:
                    continue
                x_all.append(x_num)
                y_all.append(y_num)
                pairs.append(f"({self._format_axis_value(x_num)}, {self._format_axis_value(y_num)})")

            if pairs:
                series_lines.append(f'  scatter "{self._escape_mermaid_text(series_name)}" [{", ".join(pairs)}]')

        if not series_lines:
            return None

        x_axis_label = self._get_axis_title(getattr(chart, "category_axis", None)) or "X"
        y_axis_label = self._get_axis_title(getattr(chart, "value_axis", None)) or "Y"

        x_min = min(x_all) if x_all else None
        x_max = max(x_all) if x_all else None
        y_min = min(y_all) if y_all else None
        y_max = max(y_all) if y_all else None

        if x_min is None or x_max is None or y_min is None or y_max is None:
            return None

        lines: list[str] = ["xychart-beta"]
        title = self._get_chart_title(chart)
        if title:
            lines.append(f'  title "{self._escape_mermaid_text(title)}"')

        lines.append(
            f'  x-axis "{self._escape_mermaid_text(x_axis_label)}" '
            f"{self._format_axis_value(x_min)} --> {self._format_axis_value(x_max)}"
        )
        lines.append(
            f'  y-axis "{self._escape_mermaid_text(y_axis_label)}" '
            f"{self._format_axis_value(y_min)} --> {self._format_axis_value(y_max)}"
        )
        lines.extend(series_lines)

        return "\n".join(lines)

    def _standard_chart_to_mermaid(
        self,
        chart: Any,
        chart_type: Any,
        categories: list[str],
        series_rows: list[tuple[str, list[Any]]],
    ) -> str | None:
        if not series_rows:
            return None

        mermaid_series_type = self._map_chart_type_to_mermaid(chart_type)
        if mermaid_series_type is None:
            return None

        max_cols = max((len(values) for _, values in series_rows), default=0)
        if max_cols == 0:
            return None

        if categories and len(categories) == max_cols:
            x_labels = [self._escape_mermaid_text(str(label)) for label in categories]
        else:
            x_labels = [f"Col {i + 1}" for i in range(max_cols)]

        numeric_values: list[float] = []
        for _, values in series_rows:
            for value in values:
                numeric = self._coerce_float(value)
                if numeric is not None:
                    numeric_values.append(numeric)

        y_axis_label = self._get_axis_title(getattr(chart, "value_axis", None)) or "Value"
        y_min = min(numeric_values) if numeric_values else None
        y_max = max(numeric_values) if numeric_values else None

        lines: list[str] = ["xychart-beta"]
        title = self._get_chart_title(chart)
        if title:
            lines.append(f'  title "{self._escape_mermaid_text(title)}"')

        x_axis_list = ", ".join(f'"{label}"' for label in x_labels)
        lines.append(f"  x-axis [{x_axis_list}]")

        if y_min is not None and y_max is not None and y_min != y_max:
            lines.append(
                f'  y-axis "{self._escape_mermaid_text(y_axis_label)}" '
                f"{self._format_axis_value(y_min)} --> {self._format_axis_value(y_max)}"
            )

        for series_name, values in series_rows:
            series_label = self._escape_mermaid_text(series_name)
            formatted_values = [self._format_mermaid_number(val) for val in values]
            # Pad values to match axis length
            if len(formatted_values) < max_cols:
                formatted_values.extend(["null"] * (max_cols - len(formatted_values)))
            lines.append(f'  {mermaid_series_type} "{series_label}" [{", ".join(formatted_values)}]')

        return "\n".join(lines)

    @staticmethod
    def _map_chart_type_to_mermaid(chart_type: Any) -> str | None:
        if chart_type is None:
            return None

        from pptx.enum.chart import XL_CHART_TYPE

        bar_types = {
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.COLUMN_STACKED_100,
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.BAR_STACKED,
            XL_CHART_TYPE.BAR_STACKED_100,
        }

        line_types = {
            XL_CHART_TYPE.LINE,
            XL_CHART_TYPE.LINE_MARKERS,
            XL_CHART_TYPE.LINE_STACKED,
            XL_CHART_TYPE.LINE_STACKED_100,
        }

        area_types = {
            XL_CHART_TYPE.AREA,
            XL_CHART_TYPE.AREA_STACKED,
            XL_CHART_TYPE.AREA_STACKED_100,
        }

        if chart_type in bar_types:
            return "bar"
        if chart_type in line_types:
            return "line"
        if chart_type in area_types:
            return "area"

        return None

    @staticmethod
    def _get_chart_title(chart: Any) -> str | None:
        try:
            if hasattr(chart, "has_title") and chart.has_title and chart.chart_title:
                text_frame = chart.chart_title.text_frame
                if text_frame and text_frame.text:
                    return text_frame.text.strip()
        except Exception:
            return None
        return None

    @staticmethod
    def _get_axis_title(axis: Any) -> str | None:
        if axis is None:
            return None
        try:
            if getattr(axis, "has_title", False) and axis.axis_title:
                text_frame = axis.axis_title.text_frame
                if text_frame and text_frame.text:
                    return text_frame.text.strip()
        except Exception:
            return None
        return None

    @staticmethod
    def _escape_mermaid_text(text: str) -> str:
        return str(text).replace('"', '\\"')

    @staticmethod
    def _coerce_float(value: Any) -> float | None:
        if value is None:
            return None
        if isinstance(value, (int, float)):
            return float(value)
        try:
            return float(value)
        except (TypeError, ValueError):
            return None

    @staticmethod
    def _format_axis_value(value: float) -> str:
        if float(value).is_integer():
            return str(int(round(value)))
        return f"{float(value):.6g}"

    def _format_mermaid_number(self, value: Any) -> str:
        numeric = self._coerce_float(value)
        if numeric is None:
            return "null"
        return self._format_axis_value(numeric)

    @staticmethod
    def _format_scalar(value: Any) -> str:
        numeric = PptxToAstConverter._coerce_float(value)
        if numeric is None:
            return str(value)
        if float(numeric).is_integer():
            return str(int(round(numeric)))
        return f"{float(numeric):.6g}"

    def _process_text_frame_to_ast(self, frame: "TextFrame") -> list[Node] | None:
        """Process a text frame to AST nodes.

        Parameters
        ----------
        frame : TextFrame
            Text frame to process

        Returns
        -------
        list of Node or None
            List of AST nodes (paragraphs, lists, etc.)

        """
        nodes: list[Node] = []
        slide_context = _analyze_slide_context(frame)

        # Track list accumulation
        current_list: List | None = None
        current_list_type: str | None = None

        for paragraph in frame.paragraphs:
            if not paragraph.text.strip():
                continue

            # Detect if this is a list item
            is_list_item, list_type = _detect_list_item(
                paragraph, slide_context, strict_mode=self.options.strict_list_detection
            )

            if is_list_item:
                # Check if we need to start a new list or continue current
                if current_list is None or current_list_type != list_type:
                    # Finalize previous list
                    if current_list:
                        nodes.append(current_list)

                    # Start new list
                    current_list_type = list_type
                    current_list = List(ordered=(list_type == "number"), items=[], tight=True)

                # Add item to current list
                item_content = self._process_paragraph_runs_to_inline(paragraph)
                if item_content:
                    list_item = ListItem(children=[AstParagraph(content=item_content)])
                    current_list.items.append(list_item)
            else:
                # Not a list item - finalize any current list
                if current_list:
                    nodes.append(current_list)
                    current_list = None
                    current_list_type = None

                # Process as regular paragraph
                para_content = self._process_paragraph_runs_to_inline(paragraph)
                if para_content:
                    nodes.append(AstParagraph(content=para_content))

        # Finalize any remaining list
        if current_list:
            nodes.append(current_list)

        return nodes if nodes else None

    def _process_paragraph_runs_to_inline(self, paragraph: Any) -> list[Node]:
        """Process paragraph runs to inline AST nodes using group_and_format_runs helper.

        Parameters
        ----------
        paragraph : Paragraph
            Paragraph containing runs

        Returns
        -------
        list of Node
            List of inline AST nodes

        """

        def text_extractor(run: Any) -> str:
            # Add space after text to preserve word boundaries
            # The helper strips and joins without separator, so we add space explicitly
            text = run.text if run.text else ""
            stripped = text.strip()
            if not stripped:
                return ""
            # Add space suffix if original text had trailing whitespace
            if text != stripped and text.endswith((" ", "\t", "\n")):
                return stripped + " "
            return stripped

        def format_extractor(run: Any) -> tuple[bool, bool, bool, bool, bool, bool]:
            # Return format flags in order that matches desired application order
            # group_and_format_runs processes from high index to low index
            # Original PPTX applies formats (innermost to outermost):
            # underline, bold, italic, strikethrough, subscript, superscript
            # So we return in reverse order: (sup, sub, strike, italic, bold, underline)
            if run.font:
                return (
                    bool(getattr(run.font, "superscript", False)),  # Index 0 - applied last (outermost)
                    bool(getattr(run.font, "subscript", False)),  # Index 1
                    bool(getattr(run.font, "strike", False)),  # Index 2 - strikethrough
                    bool(run.font.italic),  # Index 3
                    bool(run.font.bold),  # Index 4
                    bool(run.font.underline),  # Index 5 - applied first (innermost)
                )
            return (False, False, False, False, False, False)

        # Custom format builders matching format tuple order
        # Loop processes from high index to low: 5, 4, 3, 2, 1, 0
        # Application order: underline, bold, italic, strike, subscript, superscript
        format_builders = (
            lambda nodes: Superscript(content=nodes),  # Index 0 - applied last (outermost)
            lambda nodes: Subscript(content=nodes),  # Index 1
            lambda nodes: Strikethrough(content=nodes),  # Index 2
            lambda nodes: Emphasis(content=nodes),  # Index 3
            lambda nodes: Strong(content=nodes),  # Index 4
            lambda nodes: Underline(content=nodes),  # Index 5 - applied first (innermost)
        )

        # Process runs with formatting
        inline_nodes = group_and_format_runs(
            runs=paragraph.runs,
            text_extractor=text_extractor,
            format_extractor=format_extractor,
            format_builders=format_builders,
        )

        # Post-process to extract hyperlinks
        # PPTX hyperlinks are on runs, not separate nodes
        result: list[Node] = []
        for run in paragraph.runs:
            # Check if this run has a hyperlink
            if hasattr(run, "hyperlink") and run.hyperlink and hasattr(run.hyperlink, "address"):
                hyperlink_url = run.hyperlink.address
                if hyperlink_url:
                    # Find the corresponding text node(s) for this run
                    # For simplicity, wrap the run's text in a Link node
                    run_text = run.text.strip() if run.text else ""
                    if run_text:
                        # Create link with the formatted content
                        link_content = [Text(content=run_text)]
                        result.append(Link(url=hyperlink_url, content=cast(list[Node], link_content)))
                        continue

            # No hyperlink - use the nodes from group_and_format_runs
            # Note: This is a simplified approach - ideally we'd match runs to nodes
            # For now, if there are no hyperlinks in any runs, just return inline_nodes

        # If no hyperlinks were found, return the formatted nodes as-is
        if not result:
            return inline_nodes

        # Otherwise, we have a mix - this is complex, so for now just return inline_nodes
        # and log that hyperlinks might not be fully integrated with formatting
        # A full solution would require matching runs to nodes in group_and_format_runs result
        return inline_nodes

    def _process_table_to_ast(self, table: Any) -> AstTable | None:
        """Process a PPTX table to AST Table node.

        Parameters
        ----------
        table : Table
            PPTX table to convert

        Returns
        -------
        AstTable or None
            Table node if table has content

        """
        if len(table.rows) == 0:
            return None

        # First row is header
        header_cells: list[TableCell] = []
        for cell in table.rows[0].cells:
            if cell.text_frame:
                cell_content = self._process_paragraph_runs_to_inline_simple(cell.text_frame)
                header_cells.append(TableCell(content=cell_content))
            else:
                header_cells.append(TableCell(content=[]))

        header_row = TableRow(cells=header_cells, is_header=True)

        # Data rows (skip first row which is header)
        data_rows: list[TableRow] = []
        for i, row in enumerate(table.rows):
            if i == 0:
                continue  # Skip header row
            row_cells: list[TableCell] = []
            for cell in row.cells:
                if cell.text_frame:
                    cell_content = self._process_paragraph_runs_to_inline_simple(cell.text_frame)
                    row_cells.append(TableCell(content=cell_content))
                else:
                    row_cells.append(TableCell(content=[]))
            data_rows.append(TableRow(cells=row_cells))

        return AstTable(header=header_row, rows=data_rows)

    def _process_paragraph_runs_to_inline_simple(self, frame: "TextFrame") -> list[Node]:
        """Process all paragraphs in a text frame to inline nodes (for tables).

        Parameters
        ----------
        frame : TextFrame
            Text frame to process

        Returns
        -------
        list of Node
            List of inline nodes

        """
        result: list[Node] = []

        for paragraph in frame.paragraphs:
            para_inlines = self._process_paragraph_runs_to_inline(paragraph)
            result.extend(para_inlines)

            # Add space between paragraphs (except last)
            if paragraph != frame.paragraphs[-1]:
                result.append(Text(content=" "))

        return result

    def _process_image_to_ast(self, shape: Any) -> Image | None:
        """Process an image shape to AST Image node.

        Parameters
        ----------
        shape : Shape
            Image shape to process

        Returns
        -------
        Image or None
            Image node if image can be processed

        """
        try:
            # Extract image data
            image_data = extract_pptx_image_data(shape)

            if not image_data or not isinstance(image_data, bytes):
                return None

            # Default extension for PPTX images
            extension = "png"

            # Use sequencer for sequential attachment names
            image_filename, _ = self._attachment_sequencer(
                base_stem=self._base_filename, format_type="general", extension=extension
            )

            # Process attachment - returns dict with URL, markdown, and footnote info
            result = process_attachment(
                attachment_data=image_data,
                attachment_name=image_filename,
                alt_text="image",
                attachment_mode=self.options.attachment_mode,
                attachment_output_dir=self.options.attachment_output_dir,
                attachment_base_url=self.options.attachment_base_url,
                is_image=True,
                alt_text_mode=self.options.alt_text_mode,
            )

            # Collect footnote info if present
            if result.get("footnote_label") and result.get("footnote_content"):
                self._attachment_footnotes[result["footnote_label"]] = result["footnote_content"]

            # Use helper to convert result to Image node (eliminates regex parsing)
            image_node = attachment_result_to_image_node(result, fallback_alt_text="image")
            if image_node and isinstance(image_node, Image):
                return image_node
            return None

        except Exception as e:
            logger.debug(f"Failed to process image: {e}")
            return None

    def _extract_slide_notes(self, slide: Any) -> list[Node]:
        """Extract speaker notes from a slide.

        Parameters
        ----------
        slide : Slide
            PPTX slide to extract notes from

        Returns
        -------
        list of Node
            List of AST nodes representing the speaker notes

        Notes
        -----
        The parsing mode is controlled by self.options.comment_mode:
        - "content": Returns regular nodes with H3 heading (default, backward compatible)
        - "comment": Returns Comment node with metadata
        - "ignore": Returns empty list

        """
        notes_nodes: list[Node] = []

        # Check comment_mode option
        comment_mode = self.options.comment_mode

        # If ignore mode, return empty list
        if comment_mode == "ignore":
            return notes_nodes

        try:
            # Check if slide has notes
            if not hasattr(slide, "notes_slide"):
                return notes_nodes

            notes_slide = slide.notes_slide
            if not hasattr(notes_slide, "notes_text_frame"):
                return notes_nodes

            notes_text_frame = notes_slide.notes_text_frame
            if not notes_text_frame or not notes_text_frame.text.strip():
                return notes_nodes

            # Extract notes text
            notes_text = notes_text_frame.text.strip()

            if comment_mode == "comment":
                # Create a Comment node with metadata
                notes_comment = Comment(
                    content=notes_text,
                    metadata={
                        "comment_type": "pptx_speaker_notes",
                        "slide_number": self._current_slide_num,
                    },
                )
                notes_nodes.append(notes_comment)

            elif comment_mode == "content":
                # Original behavior: Add H3 heading and process as regular content
                notes_nodes.append(Heading(level=3, content=[Text(content="Speaker Notes")]))

                # Process the notes text frame using existing text frame processing
                frame_nodes = self._process_text_frame_to_ast(notes_text_frame)
                if frame_nodes:
                    notes_nodes.extend(frame_nodes)

        except Exception as e:
            logger.debug(f"Failed to extract speaker notes from slide: {e}")

        return notes_nodes

    def _open_pptx_as_zip(
        self, input_data: Union[str, Path, IO[bytes], bytes]
    ) -> tuple[zipfile.ZipFile | None, str | None]:
        """Open PPTX input data as a ZipFile.

        Parameters
        ----------
        input_data : str, Path, IO[bytes], or bytes
            PPTX input data

        Returns
        -------
        tuple[zipfile.ZipFile | None, str | None]
            Tuple of (zip_file, temp_file_path). temp_file_path is set if a temp file was created.

        """
        zip_file: zipfile.ZipFile | None = None
        temp_file_path: str | None = None

        if isinstance(input_data, (str, Path)):
            # File path - open directly
            zip_file = zipfile.ZipFile(str(input_data), "r")
        elif isinstance(input_data, bytes):
            # Bytes - write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                temp_file.write(input_data)
                temp_file_path = temp_file.name
            zip_file = zipfile.ZipFile(temp_file_path, "r")
        elif hasattr(input_data, "read"):
            # File-like object - write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                temp_file.write(input_data.read())
                temp_file_path = temp_file.name
                # Reset the file pointer if possible
                if hasattr(input_data, "seek"):
                    try:
                        input_data.seek(0)
                    except Exception:
                        pass
            zip_file = zipfile.ZipFile(temp_file_path, "r")

        return zip_file, temp_file_path

    def _parse_comment_authors(self, zip_file: zipfile.ZipFile, namespaces: dict[str, str]) -> dict[str, str]:
        """Parse comment authors from PPTX ZIP file.

        Parameters
        ----------
        zip_file : zipfile.ZipFile
            Open PPTX ZIP file
        namespaces : dict[str, str]
            XML namespaces

        Returns
        -------
        dict[str, str]
            Mapping of author IDs to author names

        """
        authors_map: dict[str, str] = {}
        try:
            if "ppt/commentAuthors.xml" in zip_file.namelist():
                authors_xml = zip_file.read("ppt/commentAuthors.xml")
                authors_root = ET.fromstring(authors_xml)

                # Parse author elements
                for author_elem in authors_root.findall(".//p:cmAuthor", namespaces):
                    author_id = author_elem.get("id")
                    author_name = author_elem.get("name")
                    if author_id and author_name:
                        authors_map[author_id] = author_name
        except Exception as e:
            logger.debug(f"Failed to parse comment authors: {e}")

        return authors_map

    def _parse_slide_comments(
        self,
        zip_file: zipfile.ZipFile,
        filename: str,
        authors_map: dict[str, str],
        namespaces: dict[str, str],
    ) -> tuple[int | None, list[Comment]]:
        """Parse comments from a single slide's comment file.

        Parameters
        ----------
        zip_file : zipfile.ZipFile
            Open PPTX ZIP file
        filename : str
            Comment file path within ZIP
        authors_map : dict[str, str]
            Mapping of author IDs to names
        namespaces : dict[str, str]
            XML namespaces

        Returns
        -------
        tuple[int | None, list[Comment]]
            Tuple of (slide_number, comments). Returns (None, []) on error.

        """
        try:
            # Extract slide number from filename
            slide_num_str = filename.replace("ppt/comments/comment", "").replace(".xml", "")
            slide_number = int(slide_num_str)

            # Parse comment XML
            comment_xml = zip_file.read(filename)
            comment_root = ET.fromstring(comment_xml)

            # Extract comments from this slide
            slide_comments: list[Comment] = []
            for cm_elem in comment_root.findall(".//p:cm", namespaces):
                comment_node = self._create_comment_node(cm_elem, slide_number, authors_map, namespaces)
                if comment_node:
                    slide_comments.append(comment_node)

            return slide_number, slide_comments

        except Exception as e:
            logger.debug(f"Failed to parse comments from {filename}: {e}")
            return None, []

    def _create_comment_node(
        self,
        cm_elem: Any,
        slide_number: int,
        authors_map: dict[str, str],
        namespaces: dict[str, str],
    ) -> Comment | None:
        """Create a Comment node from XML element.

        Parameters
        ----------
        cm_elem : Any
            Comment XML element
        slide_number : int
            Slide number (1-based)
        authors_map : dict[str, str]
            Mapping of author IDs to names
        namespaces : dict[str, str]
            XML namespaces

        Returns
        -------
        Comment | None
            Comment node, or None if creation fails

        """
        # Extract comment attributes
        comment_id = cm_elem.get("idx") or cm_elem.get("id")
        author_id = cm_elem.get("authorId")
        date_time = cm_elem.get("dt")

        # Extract comment text
        text_elem = cm_elem.find(".//p:text", namespaces)
        comment_text = text_elem.text if text_elem is not None and text_elem.text else ""

        # Get author name from map
        author_name = authors_map.get(author_id, f"Author {author_id}") if author_id else "Unknown"

        # Create Comment metadata
        comment_metadata: dict[str, Any] = {
            "comment_type": "pptx_comment",
            "slide_number": slide_number,
        }

        if author_name:
            comment_metadata["author"] = author_name
        if date_time:
            comment_metadata["date"] = date_time
        if comment_id:
            comment_metadata["identifier"] = comment_id

        # Extract position if available
        pos_elem = cm_elem.find(".//p:pos", namespaces)
        if pos_elem is not None:
            x_pos = pos_elem.get("x")
            y_pos = pos_elem.get("y")
            if x_pos and y_pos:
                comment_metadata["position"] = {"x": x_pos, "y": y_pos}

        return Comment(
            content=comment_text.strip() if comment_text else "",
            metadata=comment_metadata,
        )

    def _extract_pptx_comments(self, input_data: Union[str, Path, IO[bytes], bytes]) -> dict[int, list[Comment]]:
        """Extract PowerPoint comments from PPTX ZIP archive.

        Parameters
        ----------
        input_data : str, Path, IO[bytes], or bytes
            PPTX input data to extract comments from

        Returns
        -------
        dict of int to list of Comment
            Mapping of slide numbers (1-based) to lists of Comment nodes

        Notes
        -----
        PPTX comments are stored in the ZIP archive structure:
        - ppt/commentAuthors.xml - Author information
        - ppt/comments/comment1.xml - Comments for slide 1
        - ppt/comments/comment2.xml - Comments for slide 2
        - etc.

        This method uses XML parsing to extract comment data without requiring
        python-pptx support for comments.

        """
        comments_by_slide: dict[int, list[Comment]] = {}

        try:
            # Open PPTX as ZIP file
            zip_file, temp_file_path = self._open_pptx_as_zip(input_data)
            if not zip_file:
                logger.debug("Unsupported input type for comment extraction")
                return comments_by_slide

            # Define XML namespaces
            namespaces = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }

            # Parse comment authors
            authors_map = self._parse_comment_authors(zip_file, namespaces)

            # Parse comments for each slide
            for filename in zip_file.namelist():
                if filename.startswith("ppt/comments/comment") and filename.endswith(".xml"):
                    slide_number, slide_comments = self._parse_slide_comments(
                        zip_file, filename, authors_map, namespaces
                    )
                    if slide_number is not None and slide_comments:
                        comments_by_slide[slide_number] = slide_comments

            zip_file.close()

            # Clean up temp file if created
            if temp_file_path:
                try:
                    Path(temp_file_path).unlink()
                except Exception:
                    pass

        except Exception as e:
            logger.debug(f"Failed to extract PPTX comments: {e}")

        return comments_by_slide

    def extract_metadata(self, document: Any) -> DocumentMetadata:
        """Extract metadata from PPTX presentation.

        Parameters
        ----------
        document : Presentation
            python-pptx Presentation object

        Returns
        -------
        DocumentMetadata
            Extracted metadata

        """
        if not hasattr(document, "core_properties"):
            metadata = DocumentMetadata()
        else:
            props = document.core_properties
            # Use the utility function for standard metadata extraction
            metadata = map_properties_to_metadata(props, OFFICE_FIELD_MAPPING)

            # Add PPTX-specific custom metadata
            custom_properties = ["last_modified_by", "revision", "comments"]
            for prop_name in custom_properties:
                if hasattr(props, prop_name):
                    value = getattr(props, prop_name)
                    if value:
                        metadata.custom[prop_name] = value

        # Add slide count as custom metadata
        try:
            metadata.custom["slide_count"] = len(document.slides)
        except Exception:
            pass

        return metadata


# Converter metadata for registration
CONVERTER_METADATA = ConverterMetadata(
    format_name="pptx",
    extensions=[".pptx"],
    mime_types=["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
    magic_bytes=[
        (b"PK\x03\x04", 0),  # ZIP signature
    ],
    parser_class=PptxToAstConverter,
    renderer_class="all2md.renderers.pptx.PptxRenderer",
    parser_required_packages=[("python-pptx", "pptx", "")],
    renderer_required_packages=[("python-pptx", "pptx", ">=0.6.21")],
    import_error_message="PPTX conversion requires 'python-pptx'. Install with: pip install python-pptx",
    parser_options_class=PptxOptions,
    renderer_options_class="all2md.options.pptx.PptxRendererOptions",
    description="Convert PowerPoint presentations to/from Markdown",
    priority=7,
)


def _detect_list_formatting_xml(paragraph: Any) -> tuple[str | None, str | None]:
    """Detect list formatting using XML element inspection.

    Parameters
    ----------
    paragraph : Any
        The paragraph object to inspect

    Returns
    -------
    tuple[str | None, str | None]
        (list_type, list_style) where list_type is "bullet" or "number"

    """
    try:
        # Access paragraph properties XML element
        if not hasattr(paragraph, "_p") or paragraph._p is None:
            return None, None

        pPr = paragraph._p.pPr
        if pPr is None:
            return None, None

        # Check for bullet character element
        bu_char = pPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
        if bu_char is not None:
            char = bu_char.get("char", "•")
            return "bullet", char

        # Check for auto numbering element
        bu_auto_num = pPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
        if bu_auto_num is not None:
            num_type = bu_auto_num.get("type", "arabicPeriod")
            return "number", num_type

        # Check for bullet font (indicates some form of bullet formatting)
        bu_font = pPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont")
        if bu_font is not None:
            return "bullet", "default"

    except Exception:
        # Fall back gracefully if XML parsing fails
        pass

    return None, None


def _detect_list_item(paragraph: Any, slide_context: dict | None = None, strict_mode: bool = False) -> tuple[bool, str]:
    """Detect if a paragraph is a list item and determine the list type.

    Uses XML-based detection first, then falls back to heuristics unless strict_mode is enabled.

    Parameters
    ----------
    paragraph : Any
        The paragraph object to analyze
    slide_context : dict, optional
        Context about the slide to help with detection
    strict_mode : bool, default False
        If True, only use XML-based detection (no heuristics).
        If False, use XML detection with heuristic fallbacks.

    Returns
    -------
    tuple[bool, str]
        (is_list_item, list_type) where list_type is "bullet" or "number"

    """
    # First try XML-based detection for proper list formatting
    xml_list_type, xml_list_style = _detect_list_formatting_xml(paragraph)
    if xml_list_type:
        return True, xml_list_type

    # In strict mode, only trust XML detection
    if strict_mode:
        return False, "bullet"

    # Fall back to level-based detection
    if not hasattr(paragraph, "level") or paragraph.level is None:
        return False, "bullet"

    level = paragraph.level
    if level > 0:
        # Use slide context to help determine list type for indented items
        if slide_context and slide_context.get("has_numbered_list", False):
            return True, "number"
        return True, "bullet"

    # For level 0, use heuristics as last resort
    text = paragraph.text.strip() if hasattr(paragraph, "text") else ""
    if not text:
        return False, "bullet"

    # Check for explicit numbered list patterns in text
    if re.match(r"^\d+[\.\)]\s", text):
        return True, "number"

    # Check if this looks like a numbered list item based on context
    if (
        slide_context
        and slide_context.get("has_numbered_list", False)
        and ("item" in text.lower() or "first" in text.lower() or "second" in text.lower() or "third" in text.lower())
    ):
        return True, "number"

    # Use heuristics for bullet lists - shorter text that doesn't look like a title/header
    words = text.split()
    if len(words) <= 8 and not text.endswith((".", "!", "?", ":")):
        # Additional checks to avoid false positives
        if not (text.lower().startswith(("slide", "title", "chapter")) or len(words) <= 3 and text.istitle()):
            return True, "bullet"

    return False, "bullet"


def _analyze_slide_context(frame: Any) -> dict:
    """Analyze the text frame to understand the slide context for better list detection.

    Parameters
    ----------
    frame : Any
        The text frame to analyze

    Returns
    -------
    dict
        Context information about the slide

    """
    context = {"has_numbered_list": False, "paragraph_count": 0, "max_level": 0}

    for paragraph in frame.paragraphs:
        if not paragraph.text.strip():
            continue

        context["paragraph_count"] += 1

        # Track maximum indentation level
        level = getattr(paragraph, "level", 0) or 0
        context["max_level"] = max(context["max_level"], level)

        # Check if any paragraph looks like a numbered list
        text = paragraph.text.strip()
        if (
            re.match(r"^\d+[\.\)]\s", text)
            or "numbered" in text.lower()
            or "first item" in text.lower()
            or "second item" in text.lower()
            or "third item" in text.lower()
        ):
            context["has_numbered_list"] = True

    return context
